from fastapi import APIRouter, Request, BackgroundTasks, File, UploadFile, Form
from typing import List, Optional
from logger import get_logger

log = get_logger("notes")
router = APIRouter()

@router.get("/notes")
async def get_notes():
    from notes_manager import get_notes
    return get_notes()

@router.get("/api/graph")
async def get_graph():
    from notes_manager import get_db
    conn = get_db()
    c = conn.cursor()
    c.execute("SELECT id, title, project, parent_id FROM notes ORDER BY id ASC")
    rows = c.fetchall()
    
    nodes = []
    links = []
    
    workspaces = set()
    for r in rows:
        workspaces.add(r["project"] or "General")
        
    for ws in workspaces:
        nodes.append({
            "id": f"ws_{ws}",
            "name": ws,
            "group": 0,
            "val": 5, # Larger node for workspace
            "type": "workspace"
        })
        
    for r in rows:
        n_id = r["id"]
        title = r["title"] or "Untitled"
        project = r["project"] or "General"
        parent_id = r["parent_id"]
        
        nodes.append({
            "id": n_id,
            "name": title,
            "group": 1,
            "val": 2,
            "type": "note"
        })
        
        if parent_id is not None:
            # Link to parent note
            links.append({"source": n_id, "target": parent_id})
        else:
            # Link to workspace
            links.append({"source": n_id, "target": f"ws_{project}"})
            
    conn.close()
    return {"nodes": nodes, "links": links}

@router.post("/notes/update")
async def post_notes_update(request: Request, background_tasks: BackgroundTasks):
    from notes_manager import update_note, add_note
    from server import broadcast
    body = await request.json()
    index = body.get("index")
    title = body.get("title", "Untitled")
    text = body.get("text", "")
    project = body.get("project", "General")
    parent_id = body.get("parent_id", None)
    
    if index is not None and index > 0:
        success = update_note(index, title, text, project=project, parent_id=parent_id)
        if not success:
            note = add_note(text, title=title, project=project, parent_id=parent_id)
            index = note["id"]
    else:
        note = add_note(text, title=title, project=project, parent_id=parent_id)
        index = note["id"]
        
    def _summarize():
        from notes_manager import get_db
        from brain import think
        import json
        if len(text) < 50: return
        res = think(f"Extract exactly 3 short key bullet points from this text:\n\n{text}\n\nFormat as a JSON array of strings ONLY. No markdown, no explanation.")
        try:
            content = res.get("choices", [{}])[0].get("message", {}).get("content", "[]")
            import re
            arr_str = re.search(r'\[.*\]', content, re.DOTALL)
            if arr_str:
                kp = json.loads(arr_str.group())
                conn = get_db()
                c = conn.cursor()
                if index is not None and index > 0:
                    c.execute("UPDATE notes SET key_points=? WHERE id=?", (json.dumps(kp), index))
                    conn.commit()
                conn.close()
        except Exception as e:
            log.error(f"Auto-summarize failed: {e}")

    background_tasks.add_task(_summarize)
    broadcast("note_added", {})
    return {"success": True, "id": index}

@router.post("/api/notes/generate-batch")
async def generate_batch_notes(
    files: List[UploadFile] = File(...),
    prompt: Optional[str] = Form(None),
    project: str = Form("General"),
    mode: str = Form("separate")
):
    import base64
    from brain import think
    from notes_manager import add_note
    from server import broadcast
    import json
    
    parsed_files = []
    
    # Fail-Fast Parsing Loop
    for f in files:
        content = await f.read()
        filename = f.filename.lower()
        
        try:
            if filename.endswith(".pdf"):
                import io
                from pypdf import PdfReader
                reader = PdfReader(io.BytesIO(content))
                text = ""
                for page in reader.pages:
                    text += page.extract_text() + "\n"
                if not text.strip():
                    raise ValueError("No text could be extracted from PDF.")
                parsed_files.append({"type": "text", "content": text[:12000], "name": f.filename})
                
            elif filename.endswith((".pptx", ".ppt")):
                import io
                from pptx import Presentation
                prs = Presentation(io.BytesIO(content))
                text = ""
                for i, slide in enumerate(prs.slides):
                    text += f"--- Slide {i+1} ---\n"
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                if not text.strip():
                    raise ValueError("No text could be extracted from Presentation.")
                parsed_files.append({"type": "text", "content": text[:12000], "name": f.filename})
                
            elif filename.endswith((".png", ".jpg", ".jpeg", ".webp")):
                b64 = base64.b64encode(content).decode('utf-8')
                parsed_files.append({"type": "image", "content": b64, "name": f.filename})
                
            elif filename.endswith((".txt", ".md", ".csv", ".json")):
                text = content.decode('utf-8', errors='ignore')
                parsed_files.append({"type": "text", "content": text[:12000], "name": f.filename})
                
            else:
                raise ValueError(f"Unsupported file type: {filename}")
        except Exception as e:
            return {"success": False, "error": str(e), "failed_file": f.filename}
            
    system_instruction = 'You are an AI Note Generator. Respond ONLY with raw valid JSON in this exact format: {"title": "A short summary title", "body": "The detailed markdown notes..."}. Do NOT include markdown code blocks around the JSON.'
    
    if mode == "separate":
        for pf in parsed_files:
            if pf["type"] == "text":
                ctx = f"File Name: {pf['name']}\nContent:\n{pf['content']}"
                res = think(prompt or "Generate comprehensive notes for this document.", context=ctx + "\n\n" + system_instruction)
            else:
                ctx = f"File Name: {pf['name']}\nAnalyze this image and generate notes."
                res = think(prompt or ctx, context=system_instruction, image_base64=pf['content'])
                
            try:
                content_str = res.get("choices", [{}])[0].get("message", {}).get("content", "{}")
                content_str = content_str.strip()
                if content_str.startswith("```json"):
                    content_str = content_str[7:]
                if content_str.endswith("```"):
                    content_str = content_str[:-3]
                    
                note_data = json.loads(content_str.strip())
                add_note(note_data.get("body", "Failed to extract body."), title=note_data.get("title", pf["name"]), project=project)
            except Exception as e:
                log.error(f"Failed to parse JSON for {pf['name']}: {e}")
                add_note(f"Failed to generate structured notes. Raw output:\n{res}", title=pf["name"], project=project)
    else:
        combined_text = ""
        images = []
        for pf in parsed_files:
            if pf["type"] == "text":
                combined_text += f"\n\n--- File: {pf['name']} ---\n{pf['content']}"
            else:
                images.append(pf['content'])
                
        img_b64 = images[0] if images else None
        
        ctx = f"You are synthesizing multiple sources into one cohesive note.\n{combined_text}\n\n{system_instruction}"
        res = think(prompt or "Generate a unified, comprehensive note combining these sources.", context=ctx, image_base64=img_b64)
        
        try:
            content_str = res.get("choices", [{}])[0].get("message", {}).get("content", "{}")
            content_str = content_str.strip()
            if content_str.startswith("```json"):
                content_str = content_str[7:]
            if content_str.endswith("```"):
                content_str = content_str[:-3]
                
            note_data = json.loads(content_str.strip())
            add_note(note_data.get("body", "Failed to extract body."), title=note_data.get("title", "Combined Notes"), project=project)
        except Exception as e:
            log.error(f"Failed to parse JSON for combined notes: {e}")
            add_note(f"Failed to generate structured notes. Raw output:\n{res}", title="Combined Notes", project=project)

    broadcast("note_added", {})
    return {"success": True}

@router.delete("/notes/{index}")
async def delete_notes(index: int):
    from notes_manager import delete_note
    from server import broadcast
    success = delete_note(index)
    broadcast("note_added", {}) # tell frontend to reload notes
    return {"success": success}

@router.post("/notes/pin")
async def post_notes_pin(request: Request):
    from notes_manager import toggle_pin_note
    from server import broadcast
    body = await request.json()
    index = body.get("id")
    pinned = body.get("pinned", True)
    if index is not None:
        success = toggle_pin_note(index, pinned)
        broadcast("note_added", {})
        return {"success": success}
    return {"success": False, "error": "Missing id"}

@router.post("/notes/export")
async def export_notes():
    from notes_manager import get_notes
    import time
    from pathlib import Path
    
    notes = get_notes()
    timestamp = time.strftime("%Y%m%d_%H%M%S")
    export_path = Path.home() / "Documents" / "Primnox" / f"export_{timestamp}.md"
    export_path.parent.mkdir(parents=True, exist_ok=True)
    
    with open(export_path, "w", encoding="utf-8") as f:
        f.write("# Primnox Notes Export\n\n")
        for i, note in enumerate(notes):
            f.write(f"## Node {i}\n{note}\n\n")
            
    return {"success": True, "path": str(export_path)}

@router.get("/tasks")
async def get_tasks():
    from notes_manager import get_tasks
    return get_tasks()

@router.post("/tasks/{id}/complete")
async def complete_task(id: int):
    from notes_manager import complete_task
    return {"success": complete_task(id)}

@router.get("/conversations")
async def get_conversations():
    from notes_manager import get_conversations
    return get_conversations()
