# backend/server.py
from fastapi import FastAPI, WebSocket, WebSocketDisconnect, Request, BackgroundTasks, HTTPException, File, UploadFile, Form
from typing import List, Optional
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse
from core import PrimnoxCore
from logger import get_logger, get_log_buffer, APP_VERSION
import uvicorn
import asyncio
import json
import threading
import re
import logging
import os
from observer import start_clipboard_monitor, clear_clipboard_data, register_observer_callback
from pathlib import Path
import shutil
import sys


def _get_pii_model_status() -> str:
    try:
        from privacy_mirror import model_status
        return model_status()
    except Exception:
        return "unavailable"

# ── One-time DB migration: move databases from install dir → AppData ──────────
def _migrate_dbs_to_appdata():
    """
    Pre-0.0.10 databases lived next to the exe and were wiped on every update.
    On first run after an update, copy them to AppData if they aren't there yet.
    """
    appdata = os.environ.get("APPDATA")
    if not appdata:
        return
    dest_dir = Path(appdata) / "primnox_extension"
    dest_dir.mkdir(parents=True, exist_ok=True)

    # Candidate source locations: next to __file__ and next to the frozen exe
    candidates: list[Path] = [Path(__file__).parent]
    if getattr(sys, "frozen", False):
        candidates.append(Path(sys.executable).parent)

    for db_name in ("memory.db", "chat.db"):
        dest = dest_dir / db_name
        if dest.exists():
            continue  # already in AppData — nothing to do
        for src_dir in candidates:
            src = src_dir / db_name
            if src.exists() and src != dest:
                try:
                    shutil.copy2(src, dest)
                    log_migration = get_logger("migration")
                    log_migration.info(f"Migrated {db_name}: {src} → {dest}")
                except Exception as exc:
                    get_logger("migration").warning(f"Could not migrate {db_name}: {exc}")
                break

app = FastAPI()

# ── Feedback delivery ─────────────────────────────────────────────────────────
# Paste your Discord webhook URL here. Create one in:
# Discord channel → Edit Channel → Integrations → Webhooks → New Webhook → Copy URL
FEEDBACK_DISCORD_WEBHOOK = os.getenv("FEEDBACK_WEBHOOK", "")

# Zero-Trust Security Middleware
logger = logging.getLogger("primnox_firewall")

@app.middleware("http")
async def security_middleware(request: Request, call_next):
    # Absolute Ingress IP check
    client_host = request.client.host if request.client else None
    if client_host not in ["127.0.0.1", "localhost", "::1"]:
        logger.critical(f"SECURITY BREACH ATTEMPT: Ingress from {client_host}. Dropping connection.")
        raise HTTPException(status_code=403, detail="Offline-First Doctrine Violation. Connection Terminated.")

    # Host header check
    host_header = request.headers.get("host", "").lower()
    host_clean = re.sub(r':\d+$', '', host_header)
    if host_clean not in ["localhost", "127.0.0.1", "[::1]", "::1"]:
        logger.critical(f"SECURITY BREACH ATTEMPT: Host header '{host_header}' invalid. Dropping.")
        raise HTTPException(status_code=403, detail="Offline-First Doctrine Violation. Host header invalid.")

    response = await call_next(request)
    
    if "Server" in response.headers:
        del response.headers["Server"]
    return response

app.add_middleware(
    CORSMiddleware,
    # Origins the desktop shell can serve the UI from.
    #
    # Dev servers use the Vite ports. Packaged builds do NOT use http://localhost
    # at all — each shell has its own custom scheme, and a missing entry here
    # fails only in the packaged app while dev keeps working, which is the worst
    # possible place for this list to be wrong:
    #   Electron          app://.
    #   Tauri (Linux/mac) tauri://localhost
    #   Tauri (Windows)   http://tauri.localhost   (WebView2 maps the scheme to https-like URLs)
    allow_origins=[
        "http://localhost:3000",
        "http://127.0.0.1:3000",
        "http://localhost:5173",
        "http://127.0.0.1:5173",
        "app://.",
        "tauri://localhost",
        "http://tauri.localhost",
        "https://tauri.localhost",
    ],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


@app.middleware("http")
async def _restrict_to_loopback_host(request: Request, call_next):
    """Defense against DNS rebinding. This API is local-only (bound to 127.0.0.1)
    and unauthenticated; CORS doesn't stop a malicious page from *sending* simple
    cross-origin requests. Rejecting any request whose Host header isn't loopback
    means a rebound attacker domain (Host: evil.com) can't drive the API even if
    it resolves to 127.0.0.1. The real app always uses localhost/127.0.0.1."""
    host = (request.headers.get("host") or "").rsplit(":", 1)[0].strip("[]").lower()
    if host and host not in ("127.0.0.1", "localhost", "::1"):
        return JSONResponse(status_code=421, content={"detail": "Misdirected request"})
    return await call_next(request)

log = get_logger("server")
_migrate_dbs_to_appdata()   # must run before any DB module initialises
core = PrimnoxCore()
clients = set()
loop = None

_MASKED_KEYS = ["groq_api_key", "openai_api_key", "anthropic_api_key", "gemini_api_key"]

def _sanitize_settings(settings: dict) -> dict:
    """Return a copy of settings with all API keys replaced by a placeholder."""
    safe = dict(settings)
    for key in _MASKED_KEYS:
        if safe.get(key):
            safe[key] = "sk-****"
    # Never send the backup wordlist (2048 strings) over the wire
    safe.pop("backup_wordlist", None)
    return safe

def broadcast(event_type, data):
    """
    Sovereign V2: Thread-safe broadcast for frontend compatibility.
    Uses run_coroutine_threadsafe to bridge sync callback -> async WS.
    """
    global loop
    if not loop: return
    
    log.debug(f"Broadcasting event: {event_type}")
    message = json.dumps({"type": event_type, "data": data})
    
    async def send_to_all():
        for ws in list(clients):
            try:
                await ws.send_text(message)
            except Exception:
                clients.discard(ws)
                
    asyncio.run_coroutine_threadsafe(send_to_all(), loop)

core.register_broadcast_callback(broadcast)

# ── Wire reminders to the WS broadcast so a fired reminder reaches the frontend ─
# Without this the reminder loop marks reminders 'fired' but the callback is None,
# so the desktop Notification (usePrimnox handles 'reminder_triggered') never fires.
try:
    import reminder_manager
    reminder_manager.set_callback(broadcast)
except Exception as _e:
    log.warning(f"Reminder callback wiring skipped: {_e}")

# ── Auto-start backup scheduler if previously configured ──────────────────────
try:
    from backup_manager import backup_manager as _bm
    _bm._auto_start()
except Exception as _e:
    log.warning(f"Backup auto-start skipped: {_e}")

# ── Ensure local events table exists ─────────────────────────────────────────
try:
    from event_manager import init_events_table as _init_events
    _init_events()
except Exception as _e:
    log.warning(f"Events table init skipped: {_e}")

@app.websocket("/ws")
async def websocket_endpoint(ws: WebSocket):
    # Cross-Site WebSocket Hijacking guard: FastAPI's @app.middleware("http")
    # does NOT run for websocket connections, so any website could otherwise
    # open ws://127.0.0.1:4009/ws and receive core.settings (incl. API keys)
    # plus a live feed of mic/screen/chat events. Reject any connection whose
    # Origin header isn't one of our own frontends (or absent, e.g. Electron).
    origin = ws.headers.get("origin")
    allowed_origins = {
        "http://localhost:3000", "http://127.0.0.1:3000",
        "http://localhost:5173", "http://127.0.0.1:5173",
        "app://.",
    }
    if origin is not None and origin not in allowed_origins:
        log.critical(f"SECURITY: Rejected cross-origin WebSocket from origin '{origin}'")
        await ws.close(code=1008)
        return

    await ws.accept()
    log.info("WebSocket client connected")
    clients.add(ws)
    try:
        # Send initial states to client immediately
        await ws.send_text(json.dumps({"type": "mic_state", "data": {"muted": core.mic_muted}}))
        await ws.send_text(json.dumps({"type": "incognito_changed", "data": {"active": core.incognito}}))
        await ws.send_text(json.dumps({"type": "settings_updated", "data": _sanitize_settings(core.settings)}))
        while True:
            await ws.receive_text()
            await ws.send_text(json.dumps({"type": "pong", "data": {}}))
    except WebSocketDisconnect:
        log.info("WebSocket client disconnected")
        clients.discard(ws)

@app.post("/message")
async def post_message(request: Request, background_tasks: BackgroundTasks):
    content_type = request.headers.get("content-type", "")

    # ── JSON path (no files attached) ────────────────────────────────────
    if "application/json" in content_type:
        try:
            body = await request.json()
        except Exception:
            raise HTTPException(status_code=400, detail="Invalid JSON")
        text = body.get("text", "")
        session_id = body.get("sessionId", "current")
        log.info(f"Received message: {text[:50]}...")
        background_tasks.add_task(core.handle_text_input, text, session_id=session_id)
        return {"status": "ok"}

    # ── Multipart path (files attached) ──────────────────────────────────
    if "multipart/form-data" in content_type:
        import tempfile, os, io
        form = await request.form()
        text = form.get("text", "")
        session_id = form.get("sessionId", "current")
        uploaded_files = form.getlist("files")

        extracted_parts = []
        images_b64 = []
        for uf in uploaded_files:
            if not hasattr(uf, "read"):
                continue
            filename = getattr(uf, "filename", "unknown")
            content = await uf.read()
            lower = filename.lower()

            try:
                if lower.endswith(".pdf"):
                    from pypdf import PdfReader
                    reader = PdfReader(io.BytesIO(content))
                    pdf_text = "\n".join(p.extract_text() or "" for p in reader.pages)
                    extracted_parts.append(f"[File: {filename}]\n{pdf_text[:2500]}")

                elif lower.endswith((".pptx", ".ppt")):
                    from pptx import Presentation
                    prs = Presentation(io.BytesIO(content))
                    slides_text = []
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                slides_text.append(shape.text)
                    extracted_parts.append(f"[File: {filename}]\n{chr(10).join(slides_text)[:2500]}")

                elif lower.endswith((".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp")):
                    import base64
                    b64 = base64.b64encode(content).decode("utf-8")
                    images_b64.append(b64)
                    extracted_parts.append(f"[Image attached: {filename}]")

                elif lower.endswith((".txt", ".md", ".csv", ".json", ".py", ".js", ".ts", ".tsx", ".html", ".css", ".log", ".xml", ".yaml", ".yml", ".toml", ".ini", ".cfg", ".sh", ".bat", ".sql", ".rs", ".go", ".java", ".c", ".cpp", ".h", ".rb")):
                    file_text = content.decode("utf-8", errors="replace")
                    extracted_parts.append(f"[File: {filename}]\n```\n{file_text[:2500]}\n```")

                else:
                    # Unknown file type — just acknowledge it. We don't write a temp
                    # file: the path was never surfaced to the model or cleaned up, so
                    # writing one only leaked orphaned files into the temp directory.
                    extracted_parts.append(f"[File attached: {filename} ({len(content)} bytes)]")

            except Exception as e:
                log.warning(f"Failed to extract content from {filename}: {e}")
                extracted_parts.append(f"[File: {filename} — extraction failed: {e}]")

        # Build the final message with file contents injected for LLM
        full_text = text
        if extracted_parts:
            full_text = (text + "\n\n" + "\n\n".join(extracted_parts)).strip()

        # Build a clean display message with just file chips for the UI
        file_chips = " ".join(f"[📎 {getattr(uf, 'filename', 'file')}]" for uf in uploaded_files if hasattr(uf, "read"))
        display_text = (text + "\n" + file_chips).strip() if text else file_chips

        log.info(f"Received message with {len(uploaded_files)} file(s): {text[:50]}...")
        background_tasks.add_task(core.handle_text_input, full_text, session_id=session_id, display_text=display_text, images_b64=images_b64)
        return {"status": "ok", "files_processed": len(uploaded_files)}

    raise HTTPException(status_code=400, detail="Unsupported content type")

@app.get("/api/chats")
async def get_chats():
    from chat_manager import get_all_sessions
    return get_all_sessions()

@app.post("/api/chats")
async def post_chat():
    from chat_manager import create_session
    session = create_session()
    return session

@app.get("/api/chats/{session_id}")
async def get_chat_messages(session_id: str):
    from chat_manager import get_session_messages
    return get_session_messages(session_id)

@app.put("/api/chats/{session_id}")
async def put_chat(session_id: str, request: Request):
    from chat_manager import update_session
    body = await request.json()
    title = body.get("title")
    is_pinned = body.get("isPinned")
    folder_id = body.get("folderId")
    update_session(session_id, title=title, is_pinned=is_pinned, folder_id=folder_id)
    return {"status": "ok"}

@app.delete("/api/chats/{session_id}")
async def delete_chat(session_id: str):
    from chat_manager import delete_session
    delete_session(session_id)
    return {"status": "ok"}

@app.post("/api/folders")
async def create_folder_api(request: Request):
    """Create a new chat folder. Body: {title}"""
    from chat_manager import create_folder
    body = await request.json()
    title = body.get("title", "New Folder").strip() or "New Folder"
    folder = create_folder(title)
    return folder

@app.delete("/api/folders/{folder_id}")
async def delete_folder_api(folder_id: str):
    """Delete a chat folder (chats are moved out, not deleted)."""
    from chat_manager import delete_folder
    ok = delete_folder(folder_id)
    return {"status": "ok" if ok else "not_found"}

@app.post("/api/chats/{session_id}/auto_assign")
async def auto_assign_chat(session_id: str):
    from chat_manager import get_session_messages, update_session, get_db
    from brain import think

    conn = get_db()
    c = conn.cursor()
    c.execute('SELECT id, title FROM folders')
    folders = [{"id": r["id"], "title": r["title"]} for r in c.fetchall()]
    conn.close()

    if not folders:
        return {"status": "no_folders"}

    msgs = get_session_messages(session_id)[-20:]
    chat_text = "\n".join([f"{m['speaker']}: {m['text']}" for m in msgs])

    valid_ids   = [f["id"]    for f in folders]
    folder_list = "\n".join([f"{f['title']}  →  {f['id']}" for f in folders])

    prompt = (
        f"Pick the MOST relevant folder for this chat. "
        f"Reply with ONLY the exact ID from the list — no explanation, no quotes.\n\n"
        f"Folders:\n{folder_list}\n\nChat:\n{chat_text}"
    )

    try:
        result    = think(prompt)
        choices   = result.get("choices") or []
        raw       = (choices[0].get("message", {}).get("content", "") if choices else "").strip()
        # Strip any accidental surrounding quotes the LLM sometimes adds
        chosen_id = raw.strip('"\'')
    except Exception as e:
        log.warning(f"auto_assign LLM call failed: {e}")
        return {"status": "error", "reason": str(e)}

    if chosen_id in valid_ids:
        update_session(session_id, folder_id=chosen_id)
        log.info(f"Auto-assigned chat {session_id} → folder {chosen_id}")
        return {"status": "ok", "folder_id": chosen_id}

    # LLM returned garbage — log it so we can improve the prompt later
    log.warning(f"auto_assign returned invalid id '{chosen_id}' (valid: {valid_ids})")
    return {"status": "failed", "reason": "model returned invalid folder id", "raw": raw}


@app.get("/health")
async def health():
    return {"status": "ok", "version": APP_VERSION}


@app.get("/api/status")
async def get_status():
    """Rich status report: all subsystem states, DB sizes, model, feed state."""
    import datetime
    from memory import list_memories
    from notes_manager import get_notes
    from reminder_manager import list_reminders
    from skills.skill_router import list_skills
    from settings_manager import get_appdata_dir

    appdata = get_appdata_dir()

    # DB file sizes
    db_sizes = {}
    for db_name in ("memory.db", "chat.db"):
        p = appdata / db_name
        db_sizes[db_name] = p.stat().st_size // 1024 if p.exists() else 0  # KB

    # Last backup
    last_backup_name = None
    backups_dir = appdata / "backups"
    if backups_dir.exists():
        recent = sorted(backups_dir.glob("backup_*.zip"), reverse=True)
        if recent:
            last_backup_name = recent[0].name

    # Feed activity
    feed_len = len(core.feed.history)
    active_window = core.feed.active_window_title or "Unknown"

    # Counts
    try: mem_count = len(list_memories())
    except Exception: mem_count = 0
    try: notes_count = len(get_notes())
    except Exception: notes_count = 0
    try: reminders_count = len(list_reminders())
    except Exception: reminders_count = 0
    try: skills_count = len(list_skills())
    except Exception: skills_count = 0

    return {
        "status": "ok",
        "version": APP_VERSION,
        "timestamp": datetime.datetime.now().isoformat(),
        "active_window": active_window,
        "feed_events": feed_len,
        "mic_muted": core.mic_muted,
        "incognito": core.incognito,
        "active_model": core.settings.get("active_model", "Groq_Llama_3"),
        "has_api_key": bool(core.settings.get("groq_api_key") or core.settings.get("openai_api_key") or core.settings.get("anthropic_api_key")),
        "memories_count": mem_count,
        "notes_count": notes_count,
        "reminders_count": reminders_count,
        "skills_count": skills_count,
        "db_sizes_kb": db_sizes,
        "last_backup": last_backup_name,
        "pii_model_status": _get_pii_model_status(),
    }

@app.get("/api/dashboard")
async def get_dashboard():
    import datetime
    from notes_manager import get_notes
    from memory import list_memories
    from reminder_manager import list_reminders
    from skills.skill_router import list_skills

    # Feed data — ambient + window events
    history = list(core.feed.history[-25:])
    ambient_events = [h for h in core.feed.history if "Ambient:" in h]

    # Current focus
    active_window = core.feed.active_window_title or "Unknown"
    active_process = core.feed.active_process_name or "Unknown"

    # Meetings from disk
    meetings_dir = Path.home() / "Documents" / "Primnox" / "Meetings"
    meetings = []
    if meetings_dir.exists():
        for d in sorted(meetings_dir.iterdir(), reverse=True)[:5]:
            if not d.is_dir():
                continue
            summary_file = d / "summary.txt"
            has_summary = summary_file.exists()
            preview = None
            if has_summary:
                try:
                    preview = summary_file.read_text(encoding="utf-8")[:250]
                except Exception:
                    pass
            # Try to parse a date from the folder name (format: YYYYMMDD_HHMMSS_AppName)
            folder_date = None
            try:
                folder_date = datetime.datetime.strptime(d.name[:15], "%Y%m%d_%H%M%S").date()
            except Exception:
                pass
            meetings.append({
                "name": d.name,
                "has_summary": has_summary,
                "summary_preview": preview,
                "date": folder_date.isoformat() if folder_date else None,
                "is_today": folder_date == datetime.date.today() if folder_date else False,
            })

    # Counts
    try:
        notes_count = len(get_notes())
    except Exception:
        notes_count = 0
    try:
        memories_count = len(list_memories())
    except Exception:
        memories_count = 0

    # Pending reminders
    try:
        reminders = list_reminders()
        reminders_count = len(reminders)
    except Exception:
        reminders = []
        reminders_count = 0

    # Last backup info
    last_backup = None
    try:
        from settings_manager import get_appdata_dir
        backups_dir = get_appdata_dir() / "backups"
        if backups_dir.exists():
            recent = sorted(backups_dir.glob("backup_*.zip"), reverse=True)
            if recent:
                bp = recent[0]
                last_backup = {
                    "filename": bp.name,
                    "size_kb": bp.stat().st_size // 1024,
                    "timestamp": bp.stat().st_mtime,
                }
    except Exception:
        pass

    # Registered skills count
    try:
        skills_count = len(list_skills())
    except Exception:
        skills_count = 0

    # Flag whether the primary AI key is configured (don't expose the key itself)
    has_api_key = bool(core.settings.get("groq_api_key") or
                       core.settings.get("openai_api_key") or
                       core.settings.get("anthropic_api_key") or
                       core.settings.get("gemini_api_key") or
                       core.settings.get("active_model", "").endswith("_Local"))

    user_name = core.settings.get("operator_alias") or core.settings.get("nickname") or ""

    # Today's calendar events — merge local DB events with live provider events
    # (Google / Outlook / iCal / Notion). Both are normalised to the shape the
    # dashboard renders: {id, title, start_dt, all_day, color, location}.
    today_events: list = []
    now = datetime.datetime.now()
    day_start = now.replace(hour=0, minute=0, second=0, microsecond=0)
    day_end   = now.replace(hour=23, minute=59, second=59, microsecond=0)

    # 1) Local events (event_manager) — already stored in the frontend shape
    try:
        from event_manager import list_events
        today_events.extend(
            list_events(start_iso=day_start.isoformat(), end_iso=day_end.isoformat()) or []
        )
    except Exception as _e:
        log.warning(f"today_events: local events failed: {_e}")

    # 2) Live provider events — only if the user has configured any provider,
    #    so the common (no-provider) case makes zero network calls.
    try:
        if core.settings.get("calendar_providers"):
            from skills.calendar_skill import CalendarIslandSkill
            today = now.date()
            for i, ev in enumerate(CalendarIslandSkill()._fetch_events() or []):
                start = getattr(ev, "start", None)
                if start is None:
                    continue
                # astimezone() localises both naive and aware datetimes for a fair date compare
                if start.astimezone().date() != today:
                    continue
                today_events.append({
                    "id":       f"cal_{i}_{start.isoformat()}",
                    "title":    ev.title,
                    "start_dt": start.isoformat(),
                    "all_day":  False,
                    "color":    getattr(ev, "color", "#6366f1"),
                    "location": getattr(ev, "location", ""),
                })
    except Exception as _e:
        log.warning(f"today_events: provider events failed: {_e}")

    # Sort by start time and cap at 8
    today_events = sorted(today_events, key=lambda e: str(e.get("start_dt") or ""))[:8]

    return {
        "active_window": active_window,
        "active_process": active_process,
        "feed_history": history,
        "ambient_count": len(ambient_events),
        "meetings": meetings,
        "notes_count": notes_count,
        "memories_count": memories_count,
        "has_api_key": has_api_key,
        "reminders_count": reminders_count,
        "reminders": reminders,
        "last_backup": last_backup,
        "skills_count": skills_count,
        "user_name": user_name,
        "today_events": today_events,
        "pii_model_status": _get_pii_model_status(),
    }

@app.get("/api/ollama/status")
async def ollama_status():
    """Returns Ollama running status + list of installed models."""
    from brain import get_ollama_status, _safe_local_url
    from settings_manager import load_settings
    s = load_settings()
    base_url = _safe_local_url(s.get("ollama_base_url", "http://localhost:11434"), 11434)
    return get_ollama_status(base_url)

@app.get("/api/profile")
async def get_profile():
    """Return the current user emotion + learning profile."""
    from settings_manager import load_settings
    s = load_settings()
    return {
        "mood": s.get("current_mood"),
        "onboarding_profile": s.get("onboarding_profile", {}),
        "onboarding_completed": s.get("onboarding_completed", False),
    }


@app.post("/api/daily_brief")
async def post_daily_brief(background_tasks: BackgroundTasks):
    """Trigger daily debrief via DailyBriefSkill — result is broadcast via WS."""
    def _run_brief():
        from notes_manager import get_notes
        from skills.skill_router import route_skill
        notes_count = 0
        try:
            notes_count = len(get_notes())
        except Exception:
            pass
        result = route_skill(
            user_message="daily brief",
            metadata={
                "notes_count": notes_count,
                "feed_history": list(core.feed.history[-100:]),
            }
        )
        brief_text = result.get("output_text") or "Daily brief unavailable."
        broadcast("daily_debrief", {"debrief": brief_text})

    background_tasks.add_task(_run_brief)
    return {"status": "generating"}

@app.post("/api/error_explain")
async def explain_error(request: Request):
    """Feed an error string to the dynamic island error handler and return a structured payload."""
    from brain import think
    from system_prompts import ERROR_HANDLER_PROMPT
    body = await request.json()
    error_message = body.get("error_message", "")
    context = body.get("context", "")
    prompt = f"Error: {error_message}" + (f"\nContext: {context}" if context else "")
    try:
        response = think(prompt, system_override=ERROR_HANDLER_PROMPT)
        content = response.get("choices", [{}])[0].get("message", {}).get("content", "")
        # Strip markdown fences if the model wrapped its response anyway
        if "```json" in content:
            content = content.split("```json")[1].split("```")[0].strip()
        elif "```" in content:
            content = content.split("```")[1].split("```")[0].strip()
        try:
            payload = json.loads(content)
        except json.JSONDecodeError:
            # LLM response contained raw backslashes or was not valid JSON;
            # fall back to returning the plain text in a structured envelope.
            payload = {
                "summary": content[:300] if content else "AI returned an unparseable response.",
                "fix": "Review the raw AI output above.",
                "hover_text": "click to copy the fix"
            }
        return payload
    except Exception as e:
        log = get_logger("error_explain")
        log.error(f"error_explain failed: {e}")
        return {
            "summary": "something broke and i can't even explain it properly — check the logs",
            "fix": "check the logs bro",
            "hover_text": "click to copy the fix"
        }


@app.post("/api/media/control")
async def media_control(request: Request):
    """Send a playback command (play_pause / next / prev / stop) to whatever
    app is currently registered with Windows SMTC.
    Uses the same persistent SMTC event loop as feed_manager so the Windows
    Runtime COM apartment stays stable."""
    try:
        from winsdk.windows.media.control import GlobalSystemMediaTransportControlsSessionManager as _M
        import asyncio as _aio
        from feed_manager import _get_smtc_loop
    except ImportError:
        return {"success": False, "error": "winsdk not available"}

    body   = await request.json()
    action = body.get("action", "")

    async def _do():
        mgr     = await _M.request_async()
        session = mgr.get_current_session()
        if not session:
            return False
        if action == "play_pause": await session.try_toggle_play_pause_async()
        elif action == "next":     await session.try_skip_next_async()
        elif action == "prev":     await session.try_skip_previous_async()
        elif action == "stop":     await session.try_stop_async()
        else: return False
        return True

    try:
        future = _aio.run_coroutine_threadsafe(_do(), _get_smtc_loop())
        ok = await asyncio.to_thread(lambda: future.result(timeout=4))
        return {"success": ok}
    except Exception as e:
        return {"success": False, "error": str(e)}


@app.get("/api/calendar/events")
async def get_calendar_events(days: int = 7):
    """Legacy iCal endpoint kept for backwards compatibility."""
    try:
        from skills.calendar_skill import CalendarIslandSkill
        skill  = CalendarIslandSkill()
        events = skill._fetch_events()
        from datetime import datetime, timezone, timedelta
        cutoff = datetime.now(timezone.utc) + timedelta(days=days)
        events = [e for e in events if e.start <= cutoff]
        return {"events": [e.to_dict() for e in events], "count": len(events)}
    except Exception as e:
        return {"events": [], "count": 0, "error": str(e)}


# ── Local calendar CRUD ────────────────────────────────────────────────────────

@app.get("/api/events")
async def api_list_events(start: str = "", end: str = ""):
    from event_manager import list_events
    return {"events": list_events(start or None, end or None)}


@app.post("/api/events")
async def api_create_event(request: Request):
    data = await request.json()
    if not data.get("title") or not data.get("start_dt") or not data.get("end_dt"):
        raise HTTPException(status_code=400, detail="title, start_dt, end_dt are required")
    from event_manager import create_event
    return create_event(data)


@app.put("/api/events/{event_id}")
async def api_update_event(event_id: str, request: Request):
    data = await request.json()
    from event_manager import update_event
    updated = update_event(event_id, data)
    if not updated:
        raise HTTPException(status_code=404, detail="Event not found")
    return updated


@app.delete("/api/events/{event_id}")
async def api_delete_event(event_id: str):
    from event_manager import delete_event
    if not delete_event(event_id):
        raise HTTPException(status_code=404, detail="Event not found")
    return {"success": True}


@app.post("/api/events/parse-nl")
async def api_parse_nl_event(request: Request):
    """Parse a natural-language event description into a structured event dict."""
    body = await request.json()
    text = (body.get("text") or "").strip()
    if not text:
        raise HTTPException(status_code=400, detail="text is required")

    from brain import think
    from datetime import datetime as _dt
    today_str = _dt.now().strftime("%A, %B %d, %Y at %H:%M")
    prompt = (
        f'Today is {today_str}.\n\n'
        'Parse the following natural-language event description into a JSON object '
        'with these exact fields (no other text, no markdown fences):\n'
        '  title        (string)\n'
        '  start_dt     (ISO 8601 local datetime, e.g. "2026-06-12T14:00:00")\n'
        '  end_dt       (ISO 8601 local datetime; assume 1 h if not specified)\n'
        '  all_day      (boolean)\n'
        '  location     (string, empty if not mentioned)\n'
        '  description  (string, empty if not mentioned)\n'
        '  color        ("#6366f1" as default)\n\n'
        f'Input: "{text}"\n\n'
        'Respond with ONLY the JSON object.'
    )
    try:
        response = await asyncio.to_thread(think, prompt)
        content = (response.get("choices", [{}])[0]
                           .get("message", {})
                           .get("content", "")).strip()
        # Strip markdown code fences if model wraps output
        if content.startswith("```"):
            lines = content.splitlines()
            content = "\n".join(lines[1:-1] if lines[-1].strip() == "```" else lines[1:])
        import json as _json
        parsed = _json.loads(content)
        return {"event": parsed}
    except Exception as e:
        log.error(f"NL event parse failed: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/smart_paste")
async def smart_paste(request: Request):
    """Transform clipboard content via LLM to fit the current target application."""
    from brain import think
    from system_prompts import SMART_PASTE_PROMPT
    body = await request.json()
    content = body.get("content", "")
    if not content.strip():
        return {"transformed": content}
    try:
        from screen_reader import read_screen
        uia = read_screen()
        target_app = uia.get("window_title", "") if isinstance(uia, dict) else ""
    except Exception:
        target_app = ""
    prompt = f"Target app: {target_app or 'unknown'}\nContent to transform:\n{content}"
    try:
        response = think(prompt, system_override=SMART_PASTE_PROMPT)
        text = response.get("choices", [{}])[0].get("message", {}).get("content", content)
        return {"transformed": text.strip() or content}
    except Exception as e:
        log = get_logger("smart_paste")
        log.error(f"smart_paste failed: {e}")
        return {"transformed": content}


@app.get("/logs")
async def get_logs(limit: int = 200, level: str = "all"):
    return get_log_buffer(limit=limit, level=level)


# ── Backup ─────────────────────────────────────────────────────────────────────

@app.get("/api/backup/status")
async def backup_status():
    from backup_manager import backup_manager
    return backup_manager.status()


@app.post("/api/backup/validate-mnemonic")
async def backup_validate_mnemonic(body: dict):
    """Validate a 12-word phrase against the cached custom wordlist."""
    from backup_manager import validate_mnemonic
    from settings_manager import load_settings
    mnemonic = (body.get("mnemonic") or "").strip()
    if not mnemonic:
        raise HTTPException(400, "mnemonic is required")
    wordlist = load_settings().get("backup_wordlist", [])
    if not wordlist:
        raise HTTPException(400, "Wordlist not loaded — complete onboarding first")
    valid, err = validate_mnemonic(mnemonic, wordlist)
    return {"valid": valid, "error": err}


@app.post("/api/backup/setup")
async def backup_setup(body: dict):
    """Configure backup provider and activate the scheduler."""
    from backup_manager import backup_manager, validate_mnemonic
    from settings_manager import load_settings, save_settings

    mnemonic       = (body.get("mnemonic") or "").strip()
    provider       = body.get("provider", "s3")
    provider_cfg   = body.get("provider_config", {})
    interval_hours = int(float(body.get("interval_hours", 24)))

    if not mnemonic:
        raise HTTPException(400, "mnemonic is required")

    # Validation is mandatory: the wordlist checksum is the only thing that catches
    # a typo'd phrase. Skipping it would silently store a key that can never be
    # reproduced from the user's real mnemonic, making the backup unrecoverable.
    wordlist = load_settings().get("backup_wordlist", [])
    if not wordlist:
        raise HTTPException(400, "Wordlist not loaded — generate a mnemonic first")
    valid, err = validate_mnemonic(mnemonic, wordlist)
    if not valid:
        raise HTTPException(400, f"Invalid mnemonic: {err}")

    try:
        backup_manager.setup(mnemonic, provider, provider_cfg, interval_hours)
    except Exception as e:
        raise HTTPException(500, str(e))

    return {"ok": True, "message": f"Backup configured — syncing every {interval_hours}h"}


@app.post("/api/backup/now")
async def backup_now(body: dict, background_tasks: BackgroundTasks):
    """Trigger an immediate backup. Mnemonic optional if already unlocked."""
    from backup_manager import backup_manager
    mnemonic = (body.get("mnemonic") or "").strip() or None

    def _run():
        try:
            filename = backup_manager.backup_now(mnemonic)
            log.info(f"Manual backup done: {filename}")
        except Exception as e:
            log.error(f"Manual backup failed: {e}")

    background_tasks.add_task(_run)
    return {"ok": True, "message": "Backup started in background"}


@app.get("/api/backup/list")
async def backup_list():
    from backup_manager import backup_manager
    backups = await asyncio.to_thread(backup_manager.list_backups)
    return {"backups": backups}


@app.post("/api/backup/restore")
async def backup_restore(body: dict):
    """Download + decrypt + restore a named backup. Runs synchronously so errors propagate."""
    from backup_manager import backup_manager
    filename = body.get("filename", "").strip()
    mnemonic = (body.get("mnemonic") or "").strip()
    if not filename:
        raise HTTPException(400, "filename is required")
    if not mnemonic:
        raise HTTPException(400, "mnemonic is required")

    try:
        await asyncio.to_thread(backup_manager.restore, filename, mnemonic)
    except Exception as e:
        log.error(f"Restore failed: {e}")
        raise HTTPException(500, str(e))

    return {"ok": True, "message": "Restore complete — restart Primnox for changes to take effect"}


@app.post("/api/backup/import")
async def backup_import(file: UploadFile = File(...), mnemonic: str = Form(...)):
    """Import + decrypt + restore a .prx backup uploaded from disk. No cloud
    provider needed — only the file and the seed phrase. Lets a fresh install
    recover without first reconfiguring the original cloud provider."""
    mnemonic = (mnemonic or "").strip()
    if not mnemonic:
        raise HTTPException(400, "mnemonic is required")
    data = await file.read()
    if not data:
        raise HTTPException(400, "empty file")
    if data[:4] != b"PRNX":
        raise HTTPException(400, "Not a Primnox backup (.prx) file")

    from backup_manager import backup_manager
    try:
        await asyncio.to_thread(backup_manager.restore_from_bytes, data, mnemonic)
    except Exception as e:
        log.error(f"Import restore failed: {e}")
        # InvalidTag (wrong seed) is the common case — give a friendlier message
        msg = "Wrong seed phrase for this backup" if "InvalidTag" in type(e).__name__ else str(e)
        raise HTTPException(500, msg)

    return {"ok": True, "message": "Import complete — restart Primnox for changes to take effect"}


@app.delete("/api/backup/{filename}")
async def backup_delete(filename: str):
    import re as _re
    if not _re.fullmatch(r'[A-Za-z0-9_\-]+\.prx', filename):
        raise HTTPException(400, "Invalid backup filename")
    from backup_manager import backup_manager
    try:
        await asyncio.to_thread(backup_manager.delete_backup, filename)
        return {"ok": True}
    except Exception as e:
        raise HTTPException(500, str(e))


@app.post("/api/backup/test-connection")
async def backup_test_connection():
    from backup_manager import backup_manager
    ok = await asyncio.to_thread(backup_manager.test_connection)
    return {"ok": ok, "message": "Connected" if ok else "Connection failed — check credentials"}


@app.post("/api/backup/disable")
async def backup_disable():
    from backup_manager import backup_manager
    backup_manager.disable()
    return {"ok": True}


@app.post("/api/backup/wordlist")
async def backup_store_wordlist(body: dict):
    """
    Store the custom Primnox wordlist (fetched from seed.primnox.com during onboarding).
    Expects: {"wordlist": ["word1", "word2", ...]}  — exactly 2048 words.
    """
    from settings_manager import load_settings, save_settings
    wordlist = body.get("wordlist", [])
    if len(wordlist) != 2048:
        raise HTTPException(400, f"Wordlist must have 2048 words, got {len(wordlist)}")
    s = load_settings()
    s["backup_wordlist"] = wordlist
    save_settings(s)
    return {"ok": True}


@app.post("/api/backup/generate-mnemonic")
async def backup_generate_mnemonic():
    """
    Generate a fresh 12-word BIP-39-style mnemonic.

    Uses the stored custom Primnox wordlist if available (see POST /api/backup/wordlist).
    Falls back to fetching from seed.primnox.com, then to the BIP-39 English wordlist
    from the canonical GitHub source. The fetched list is cached to settings so subsequent
    calls work offline.
    """
    from backup_manager import generate_mnemonic
    from settings_manager import load_settings, save_settings

    s = load_settings()
    wordlist: list[str] = s.get("backup_wordlist", [])

    if len(wordlist) != 2048:
        # Try to fetch and cache from remote sources
        fetched: list[str] = []
        urls = [
            "https://seed.primnox.com/wordlist.txt",
            "https://raw.githubusercontent.com/trezor/python-mnemonic/master/src/mnemonic/wordlist/english.txt",
        ]
        for url in urls:
            try:
                import httpx
                async with httpx.AsyncClient(follow_redirects=True) as client:
                    r = await client.get(url, timeout=8)
                if r.status_code == 200:
                    words = [w.strip() for w in r.text.splitlines()
                             if w.strip() and not w.startswith("#")]
                    if len(words) >= 2048:
                        fetched = words[:2048]
                        break
            except Exception:
                continue

        if len(fetched) == 2048:
            wordlist = fetched
            s["backup_wordlist"] = wordlist
            save_settings(s)
        else:
            raise HTTPException(
                503,
                "Wordlist unavailable — visit seed.primnox.com or check your internet connection",
            )

    try:
        phrase = generate_mnemonic(wordlist)
        return {"mnemonic": phrase}
    except Exception as e:
        raise HTTPException(500, str(e))


@app.get("/notes")
async def get_notes():
    from notes_manager import get_notes
    return get_notes()

@app.get("/api/graph")
async def get_graph():
    import re as _re
    import json as _json
    from collections import Counter
    from notes_manager import get_db as notes_get_db

    # ── Fetch notes ───────────────────────────────────────────────────────────
    conn = notes_get_db()
    c = conn.cursor()
    c.execute("SELECT id, title, project, parent_id, key_points, text FROM notes ORDER BY id ASC")
    note_rows = c.fetchall()
    conn.close()

    # ── Fetch memories ────────────────────────────────────────────────────────
    try:
        from memory import list_memories
        mem_rows = list_memories(include_stale=False)
    except Exception:
        mem_rows = []

    STOP = {
        'the','a','an','and','or','but','in','on','at','to','for','of','with','by','from',
        'is','are','was','were','be','been','have','has','had','do','does','did','will',
        'would','could','should','may','might','shall','can','this','that','these','those',
        'it','its','they','them','their','we','our','you','your','i','my','me','he','she',
        'his','her','not','no','so','up','out','if','then','than','just','also','about',
        'into','over','after','before','between','through','during','without','like','when',
        'what','how','why','which','who','where','untitled','note','notes','new','user',
        'said','says','want','wants','asked','told','uses','used','need','needs','know',
    }

    def keywords(text: str) -> set:
        if not text:
            return set()
        words = _re.findall(r'\b[a-z]{4,}\b', text.lower())
        return {w for w in words if w not in STOP}

    nodes: list = []
    links: list = []
    seen_links: set = set()

    def add_link(src, tgt, link_type="related"):
        key = f"{min(str(src), str(tgt))}___{max(str(src), str(tgt))}"
        if key not in seen_links:
            seen_links.add(key)
            links.append({"source": src, "target": tgt, "type": link_type})

    # ── Workspace nodes ───────────────────────────────────────────────────────
    workspaces = {r["project"] or "General" for r in note_rows}
    for ws in workspaces:
        nodes.append({"id": f"ws_{ws}", "name": ws, "group": 0, "val": 5, "type": "workspace"})

    # ── Note nodes + keyword extraction ──────────────────────────────────────
    item_kw: dict[str, set] = {}   # unified id → keywords for both notes and memories

    for r in note_rows:
        n_id  = r["id"]
        title = r["title"] or "Untitled"
        proj  = r["project"] or "General"
        pid   = r["parent_id"]

        nodes.append({"id": n_id, "name": title, "group": 1, "val": 2, "type": "note"})

        if pid is not None:
            add_link(n_id, pid, "hierarchy")
        else:
            add_link(n_id, f"ws_{proj}", "hierarchy")

        kw = keywords(title)
        try:
            kp_list = _json.loads(r["key_points"] or "[]")
            for kp in (kp_list if isinstance(kp_list, list) else []):
                kw |= keywords(str(kp))
        except Exception:
            pass
        kw |= keywords((r["text"] or "")[:400])
        item_kw[str(n_id)] = kw

    # ── Memory nodes + keyword extraction ────────────────────────────────────
    for m in mem_rows:
        m_id  = f"mem_{m['key']}"
        label = m["text"][:50] + ("…" if len(m["text"]) > 50 else "")
        cat   = m.get("category", "session")
        nodes.append({"id": m_id, "name": label, "group": 3, "val": 1.5, "type": "memory", "category": cat})
        item_kw[m_id] = keywords(m["text"])

    # ── Concept nodes: terms shared across 3+ items (notes + memories) ────────
    term_count: Counter = Counter()
    for kw_set in item_kw.values():
        term_count.update(kw_set)

    concept_terms = {term for term, cnt in term_count.items() if cnt >= 3}
    for term in concept_terms:
        nodes.append({"id": f"tag_{term}", "name": f"#{term}", "group": 2, "val": 1, "type": "tag"})
        for item_id, kw_set in item_kw.items():
            if term in kw_set:
                # Use the numeric id for notes, string id for memories/tags
                src = int(item_id) if item_id.lstrip('-').isdigit() else item_id
                add_link(src, f"tag_{term}", "concept")

    # ── Direct cross-links for strong pairwise overlap ────────────────────────
    # Build an inverted index over non-concept terms so we only compare items
    # that actually share a term, instead of scanning all O(n²) pairs. Concept
    # terms (3+ items) are excluded here, so each remaining term maps to ≤2 items
    # and the work is linear in the number of term occurrences.
    from collections import defaultdict
    postings: dict = defaultdict(list)
    for item_id, kw_set in item_kw.items():
        for term in (kw_set - concept_terms):
            postings[term].append(item_id)

    pair_shared: dict = defaultdict(int)
    for ids in postings.values():
        for a in range(len(ids)):
            for b in range(a + 1, len(ids)):
                key = (ids[a], ids[b]) if ids[a] < ids[b] else (ids[b], ids[a])
                pair_shared[key] += 1

    for (id_a, id_b), shared_count in pair_shared.items():
        if shared_count >= 2:
            src_a = int(id_a) if id_a.lstrip('-').isdigit() else id_a
            src_b = int(id_b) if id_b.lstrip('-').isdigit() else id_b
            add_link(src_a, src_b, "related")

    return {"nodes": nodes, "links": links}

@app.post("/notes/update")
async def post_notes_update(request: Request, background_tasks: BackgroundTasks):
    from notes_manager import update_note, add_note
    body = await request.json()
    index = body.get("index")
    title = body.get("title", "Untitled")
    text = body.get("text", "")
    project = body.get("project", "General")
    parent_id = body.get("parent_id", None)
    
    if index is not None and index > 0:
        success = update_note(index, title, text, project=project, parent_id=parent_id)
        if not success:
            note = add_note(text, title=title, project=project, parent_id=parent_id)
            index = note["id"]
    else:
        note = add_note(text, title=title, project=project, parent_id=parent_id)
        index = note["id"]
        
    def _summarize():
        from notes_manager import get_db
        from brain import think
        import json
        if len(text) < 50: return
        res = think(f"Extract exactly 3 short key bullet points from this text:\n\n{text}\n\nFormat as a JSON array of strings ONLY. No markdown, no explanation.")
        try:
            content = res.get("choices", [{}])[0].get("message", {}).get("content", "[]")
            import re
            arr_str = re.search(r'\[.*\]', content, re.DOTALL)
            if arr_str:
                kp = json.loads(arr_str.group())
                conn = get_db()
                c = conn.cursor()
                if index is not None and index > 0:
                    c.execute("UPDATE notes SET key_points=? WHERE id=?", (json.dumps(kp), index))
                    conn.commit()
                conn.close()
        except Exception as e:
            log.error(f"Auto-summarize failed: {e}")

    background_tasks.add_task(_summarize)
    broadcast("note_added", {})
    return {"success": True, "id": index}

@app.post("/api/notes/generate-batch")
async def generate_batch_notes(
    files: List[UploadFile] = File(...),
    prompt: Optional[str] = Form(None),
    project: str = Form("General"),
    mode: str = Form("separate")
):
    import base64
    from brain import think
    from notes_manager import add_note
    import json
    
    parsed_files = []
    
    # Fail-Fast Parsing Loop
    for f in files:
        content = await f.read()
        filename = f.filename.lower()
        
        try:
            if filename.endswith(".pdf"):
                import io
                from pypdf import PdfReader
                reader = PdfReader(io.BytesIO(content))
                text = ""
                for page in reader.pages:
                    text += page.extract_text() + "\n"
                if not text.strip():
                    raise ValueError("No text could be extracted from PDF.")
                parsed_files.append({"type": "text", "content": text[:12000], "name": f.filename})
                
            elif filename.endswith((".pptx", ".ppt")):
                import io
                from pptx import Presentation
                prs = Presentation(io.BytesIO(content))
                text = ""
                for i, slide in enumerate(prs.slides):
                    text += f"--- Slide {i+1} ---\n"
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                if not text.strip():
                    raise ValueError("No text could be extracted from Presentation.")
                parsed_files.append({"type": "text", "content": text[:12000], "name": f.filename})
                
            elif filename.endswith((".png", ".jpg", ".jpeg", ".webp")):
                b64 = base64.b64encode(content).decode('utf-8')
                parsed_files.append({"type": "image", "content": b64, "name": f.filename})
                
            elif filename.endswith((".txt", ".md", ".csv", ".json")):
                text = content.decode('utf-8', errors='ignore')
                parsed_files.append({"type": "text", "content": text[:12000], "name": f.filename})
                
            else:
                raise ValueError(f"Unsupported file type: {filename}")
        except Exception as e:
            return {"success": False, "error": str(e), "failed_file": f.filename}
            
    system_instruction = 'You are an AI Note Generator. Respond ONLY with raw valid JSON in this exact format: {"title": "A short summary title", "body": "The detailed markdown notes..."}. Do NOT include markdown code blocks around the JSON.'
    
    if mode == "separate":
        for pf in parsed_files:
            if pf["type"] == "text":
                ctx = f"File Name: {pf['name']}\nContent:\n{pf['content']}"
                res = think(prompt or "Generate comprehensive notes for this document.", context=ctx + "\n\n" + system_instruction)
            else:
                ctx = f"File Name: {pf['name']}\nAnalyze this image and generate notes."
                res = think(prompt or ctx, context=system_instruction, image_base64=pf['content'])
                
            try:
                content_str = res.get("choices", [{}])[0].get("message", {}).get("content", "{}")
                content_str = content_str.strip()
                if content_str.startswith("```json"):
                    content_str = content_str[7:]
                if content_str.endswith("```"):
                    content_str = content_str[:-3]
                    
                note_data = json.loads(content_str.strip())
                add_note(note_data.get("body", "Failed to extract body."), title=note_data.get("title", pf["name"]), project=project)
            except Exception as e:
                log.error(f"Failed to parse JSON for {pf['name']}: {e}")
                add_note(f"Failed to generate structured notes. Raw output:\n{res}", title=pf["name"], project=project)
    else:
        combined_text = ""
        images = []
        for pf in parsed_files:
            if pf["type"] == "text":
                combined_text += f"\n\n--- File: {pf['name']} ---\n{pf['content']}"
            else:
                images.append(pf['content'])
                
        img_b64 = images[0] if images else None
        
        ctx = f"You are synthesizing multiple sources into one cohesive note.\n{combined_text}\n\n{system_instruction}"
        res = think(prompt or "Generate a unified, comprehensive note combining these sources.", context=ctx, image_base64=img_b64)
        
        try:
            content_str = res.get("choices", [{}])[0].get("message", {}).get("content", "{}")
            content_str = content_str.strip()
            if content_str.startswith("```json"):
                content_str = content_str[7:]
            if content_str.endswith("```"):
                content_str = content_str[:-3]
                
            note_data = json.loads(content_str.strip())
            add_note(note_data.get("body", "Failed to extract body."), title=note_data.get("title", "Combined Notes"), project=project)
        except Exception as e:
            log.error(f"Failed to parse JSON for combined notes: {e}")
            add_note(f"Failed to generate structured notes. Raw output:\n{res}", title="Combined Notes", project=project)

    broadcast("note_added", {})
    return {"success": True}

@app.post("/api/generate")
async def post_generate(request: Request):
    try:
        body = await request.json()
    except Exception:
        raise HTTPException(status_code=400, detail="Invalid JSON")
    
    prompt = body.get("prompt", "")
    from brain import think
    result = think(prompt)
    return result

@app.delete("/notes/{index}")
async def delete_notes(index: int):
    from notes_manager import delete_note
    success = delete_note(index)
    broadcast("note_added", {}) # tell frontend to reload notes
    return {"success": success}

@app.post("/notes/pin")
async def post_notes_pin(request: Request):
    from notes_manager import toggle_pin_note
    body = await request.json()
    index = body.get("id")
    pinned = body.get("pinned", True)
    if index is not None:
        success = toggle_pin_note(index, pinned)
        broadcast("note_added", {})
        return {"success": success}
    return {"success": False, "error": "Missing id"}

@app.post("/notes/export")
async def export_notes():
    from notes_manager import get_notes
    import time
    from pathlib import Path
    
    notes = get_notes()
    timestamp = time.strftime("%Y%m%d_%H%M%S")
    export_path = Path.home() / "Documents" / "Primnox" / f"export_{timestamp}.md"
    export_path.parent.mkdir(parents=True, exist_ok=True)
    
    with open(export_path, "w", encoding="utf-8") as f:
        f.write("# Primnox Notes Export\n\n")
        for i, note in enumerate(notes):
            f.write(f"## Node {i}\n{note}\n\n")
            
    return {"success": True, "filename": export_path.name}

@app.get("/tasks")
async def get_tasks():
    from notes_manager import get_tasks
    return get_tasks()

@app.post("/tasks")
async def create_task(request: Request):
    """Create a new task. Body: {text, priority?, due_date?}"""
    from notes_manager import add_task
    body = await request.json()
    text = body.get("text", "").strip()
    if not text:
        raise HTTPException(status_code=400, detail="text is required")
    priority = body.get("priority", "normal")
    due_date = body.get("due_date")
    task = add_task(text, priority=priority, due_date=due_date)
    broadcast("task_added", {"text": text})
    return {"success": True, "task": task}

@app.post("/tasks/{id}/complete")
async def complete_task(id: int):
    from notes_manager import complete_task
    return {"success": complete_task(id)}

@app.delete("/tasks/{id}")
async def delete_task_endpoint(id: int):
    from notes_manager import delete_task
    ok = delete_task(id)
    if ok:
        broadcast("task_added", {})  # refresh frontend task list
    return {"success": ok}

@app.get("/memory")
async def get_memory():
    from memory import get_memory
    return get_memory()

@app.delete("/memory/{key}")
async def delete_memory(key: str):
    from memory import delete_memory
    delete_memory(key)
    return {"success": True}


@app.get("/api/vault/status")
async def vault_status():
    """Status of the local 12-word-seed encryption of memory.db."""
    import local_vault
    from memory import DB_PATH
    return {
        "enabled": local_vault.is_enabled(DB_PATH),
        "locked": local_vault.is_locked(DB_PATH),
    }


import time as _time_mod
_vault_phrase_tokens: dict[str, float] = {}  # token → expiry timestamp

@app.post("/api/vault/phrase-token")
async def vault_phrase_issue_token():
    """Issue a 60-second single-use token required to fetch the recovery phrase.
    Vault must be enabled and unlocked to obtain a token."""
    import local_vault
    from memory import DB_PATH
    if not local_vault.is_enabled(DB_PATH):
        raise HTTPException(status_code=403, detail="Vault not enabled")
    if local_vault.is_locked(DB_PATH):
        raise HTTPException(status_code=403, detail="Vault is locked — unlock first")
    import secrets as _sec
    token = _sec.token_hex(32)
    _vault_phrase_tokens[token] = _time_mod.time() + 60
    return {"token": token}

@app.get("/api/vault/phrase")
async def api_vault_phrase(request: Request):
    """Return the stored recovery phrase — requires a valid one-time token from
    /api/vault/phrase-token, passed in the X-Vault-Token header (not the URL, so
    the capability token never lands in access logs)."""
    token = request.headers.get("x-vault-token", "")
    now = _time_mod.time()
    # Prune expired tokens
    expired = [t for t, exp in _vault_phrase_tokens.items() if now > exp]
    for t in expired:
        _vault_phrase_tokens.pop(t, None)
    if not token or token not in _vault_phrase_tokens or now > _vault_phrase_tokens[token]:
        raise HTTPException(status_code=403, detail="Valid one-time token required. Request one via POST /api/vault/phrase-token.")
    _vault_phrase_tokens.pop(token)  # single-use
    import local_vault
    phrase = local_vault.get_stored_phrase()
    if phrase is None:
        raise HTTPException(status_code=404, detail="No recovery phrase stored. You must re-enable the vault to generate a new one.")
    return {"phrase": phrase}


@app.post("/api/vault/setup")
async def vault_setup(body: dict = None):
    """
    Enable local encryption for memory.db.

    Body: {} to auto-generate a new 12-word mnemonic (returned ONCE — the
    user must save it, it cannot be recovered), or {"mnemonic": "..."} to
    use a user-supplied phrase (e.g. reuse their backup seed).
    """
    import local_vault
    from memory import DB_PATH
    body = body or {}
    try:
        mnemonic = await asyncio.to_thread(local_vault.setup_vault, DB_PATH, body.get("mnemonic"))
    except ValueError as e:
        raise HTTPException(400, str(e))
    return {"ok": True, "mnemonic": mnemonic}


@app.post("/api/vault/unlock")
async def vault_unlock(body: dict):
    """Unlock the local vault with a 12-word mnemonic."""
    import local_vault
    from memory import DB_PATH
    mnemonic = (body or {}).get("mnemonic", "")
    try:
        await asyncio.to_thread(local_vault.unlock_vault, DB_PATH, mnemonic)
    except (ValueError, PermissionError) as e:
        raise HTTPException(400, str(e))
    except FileNotFoundError as e:
        raise HTTPException(404, str(e))
    return {"ok": True}


@app.post("/api/vault/disable")
async def vault_disable():
    """Disable local encryption — decrypts memory.db and removes the vault."""
    import local_vault
    from memory import DB_PATH
    try:
        await asyncio.to_thread(local_vault.disable_vault, DB_PATH)
    except PermissionError as e:
        raise HTTPException(400, str(e))
    return {"ok": True}


@app.get("/conversations")
async def get_conversations():
    from notes_manager import get_conversations
    return get_conversations()

# ── Onboarding environment scan ───────────────────────────────────────────────
# Reads a few human-meaningful "signal" files per project and summarizes them
# with the LOCAL model (Ollama / llama.cpp). File CONTENTS never leave the
# device: if no local model is reachable we fall back to a deterministic,
# cloud-free heuristic. Folder names are never shipped to a cloud provider.

_SCAN_IGNORE_DIRS = {
    # Windows
    'AppData', 'Application Data', 'Local Settings', 'Cookies',
    'Recent', 'SendTo', 'Start Menu', 'NetHood', 'PrintHood',
    'Templates', 'Saved Games', 'Favorites', 'Contacts', 'Searches',
    'Links', 'OneDrive',
    # Cross-platform user folders and build junk
    'node_modules', 'venv', '.venv', '.git', 'dist', 'build',
    '__pycache__', '.cache', '.cargo', '.rustup', '.npm', '.vscode',
    'Downloads', 'Music', 'Pictures', 'Videos', 'Documents', 'Desktop',
    'Public',
    # macOS. `Library` alone holds hundreds of thousands of files, and
    # Library/CloudStorage is network-backed — walking it made the onboarding
    # scan hang indefinitely (measured: no response after 90 s), which trapped
    # the user on a step that has no forward button. `Library` is not hidden, so
    # the `startswith('.')` filter never excluded it.
    'Library', 'Applications', 'Movies', 'Pictures Library',
    'Creative Cloud Files', 'Dropbox', 'Google Drive', 'iCloud Drive',
    # Linux
    '.local', '.config', 'snap', '.steam',
}

_SCAN_EXT_SKILL = {
    '.py': 'Python', '.ts': 'TypeScript', '.tsx': 'TypeScript',
    '.js': 'JavaScript', '.jsx': 'JavaScript', '.rs': 'Rust',
    '.go': 'Go', '.cpp': 'C++', '.h': 'C++', '.java': 'Java', '.cs': 'C#',
}

# Files that mark a folder as a real project ROOT (so we list "Primnox", not
# its "src"/"backend" subfolders). A ".git" directory counts too.
_SCAN_MARKER_FILES = {
    'package.json', 'pyproject.toml', 'Cargo.toml', 'go.mod', 'setup.py',
    'pom.xml', 'build.gradle', 'requirements.txt', 'README.md', 'README.MD',
}

# Deterministic, fully-offline mapping used when no local model is available.
_SKILL_PROFILE = {
    "Python":     (["Backend Development", "Automation", "Data & ML"], ["Python ecosystem", "Scripting"]),
    "TypeScript": (["Web Development", "Frontend Engineering"],        ["Type-safe JavaScript", "UI frameworks"]),
    "JavaScript": (["Web Development", "Frontend Engineering"],        ["Browser APIs", "Node.js"]),
    "Rust":       (["Systems Programming", "Performance"],             ["Memory safety", "Native tooling"]),
    "Go":         (["Backend Services", "Cloud & Infra"],             ["Concurrency", "Microservices"]),
    "C++":        (["Systems Programming", "Performance"],             ["Low-level memory", "Native code"]),
    "Java":       (["Backend Development", "Enterprise"],             ["JVM ecosystem", "OOP design"]),
    "C#":         (["App Development", ".NET"],                        [".NET ecosystem", "Desktop / Game dev"]),
}


def _read_signal_files(proj_path, per_file=1500, max_total=4000):
    """Read a small, bounded sample of human-meaningful files from a project.

    Stays well under a few KB total so the local model can summarize fast.
    """
    import json as _json
    from pathlib import Path as _P
    snippets = []
    total = [0]
    p = _P(proj_path)

    def _add(label, text):
        if not text:
            return
        text = (" ".join(text.split()) if label == "package.json" else text.strip())
        if not text:
            return
        room = max_total - total[0]
        if room <= 0:
            return
        chunk = text[:min(per_file, room)]
        snippets.append(f"[{label}] {chunk}")
        total[0] += len(chunk)

    try:
        for fn in ("README.md", "readme.md", "README", "README.txt"):
            f = p / fn
            if f.is_file():
                _add("README", f.read_text(encoding="utf-8", errors="ignore"))
                break
        pj = p / "package.json"
        if pj.is_file():
            try:
                d = _json.loads(pj.read_text(encoding="utf-8", errors="ignore"))
                meta = " ".join(str(d.get(k, "")) for k in ("name", "description") if d.get(k))
                deps = ", ".join(list((d.get("dependencies") or {}).keys())[:15])
                _add("package.json", f"{meta} | deps: {deps}")
            except Exception:
                pass
        for fn in ("pyproject.toml", "Cargo.toml", "setup.cfg", "go.mod", "composer.json"):
            f = p / fn
            if f.is_file():
                _add(fn, f.read_text(encoding="utf-8", errors="ignore")[:600])
    except Exception:
        pass
    return "  ".join(snippets)


# Wall-clock budgets for the onboarding scan.
#
# Measured 2026-08-06: this endpoint never returned — 90 s with no response.
# It walks the whole home directory to depth 4 and only stops early once it has
# found 30 projects (which may never happen), then makes a think_local() call
# with a 90 s timeout even when no local engine is running. The onboarding step
# that calls it has no forward button, so the user was trapped on it forever.
# Both phases are now bounded; a partial profile beats an infinite spinner.
_SCAN_WALK_BUDGET_S = 5.0
_SCAN_TOTAL_BUDGET_S = 20.0


def _onboarding_scan_sync():
    import os
    import json
    import getpass
    import re as _re
    import time as _time
    from pathlib import Path
    from brain import think_local

    _scan_started = _time.monotonic()
    projects = []
    project_paths = {}
    skills = set()
    home_dir = Path.home()

    code_folders = {}  # fallback: folders with code but no project-root marker
    try:
        for root, dirs, files in os.walk(str(home_dir)):
            has_git = '.git' in dirs
            # Skip heavy / hidden folders in-place
            dirs[:] = [d for d in dirs if d not in _SCAN_IGNORE_DIRS and not d.startswith('.')]
            depth = root.count(os.sep) - str(home_dir).count(os.sep)
            if depth > 4:
                dirs.clear()
                continue
            has_code = False
            for f in files:
                ext = os.path.splitext(f)[1].lower()
                if ext in _SCAN_EXT_SKILL:
                    skills.add(_SCAN_EXT_SKILL[ext])
                    has_code = True
            if depth > 0:
                # Case-insensitive on Windows (NTFS paths are case-insensitive but
                # Python str comparison is not). Lower both sides before comparing.
                root_lower = root.lower()
                under_known = any(
                    root_lower == rp.lower() or root_lower.startswith(rp.lower() + os.sep)
                    for rp in project_paths.values()
                )
                if under_known:
                    dirs.clear()  # stop descending into this known root's subtree
                    continue
                is_root = has_git or any(m in files for m in _SCAN_MARKER_FILES)
                if is_root:
                    name = os.path.basename(root)
                    if name not in project_paths:
                        projects.append(name)
                        project_paths[name] = root
                    dirs.clear()  # don't recurse into this root's children
                elif has_code:
                    code_folders.setdefault(os.path.basename(root), root)
            if len(projects) > 30:
                break
            # Hard time budget — a home directory can be arbitrarily large, and
            # "found 30 projects" is not a bound that is guaranteed to be hit.
            if _time.monotonic() - _scan_started > _SCAN_WALK_BUDGET_S:
                log.info(f"Onboarding walk hit {_SCAN_WALK_BUDGET_S}s budget; "
                         f"continuing with {len(projects)} projects found so far.")
                break
    except Exception as e:
        log.warning(f"Onboarding scanner walk error: {e}")

    # Prefer real project roots; only fall back to loose code folders if none found.
    if not projects and code_folders:
        for nm, pth in list(code_folders.items())[:10]:
            projects.append(nm)
            project_paths[nm] = pth
    projects = projects[:10] if projects else ["Workspace Sandbox"]
    skills = list(skills)[:10] if skills else ["System Administration"]

    # ── Read real file contents and summarize ON-DEVICE ──────────────────────
    context_blocks = []
    for name in projects[:6]:
        path = project_paths.get(name)
        if not path:
            continue
        sig = _read_signal_files(path)
        if sig:
            context_blocks.append(f"### {name}\n{sig}")
    context = "\n\n".join(context_blocks)[:8000]

    llm_data = None
    used_local = False
    if context:
        prompt = (
            "You are profiling a developer from their own project files, read locally on "
            "their machine and never uploaded. Using ONLY the content below, describe them.\n\n"
            f"{context}\n\n"
            "Respond with ONLY valid JSON, no prose or code fences, exactly this shape:\n"
            '{"role": "job title in 1-3 words", '
            '"topics": ["3 things they work on"], '
            '"communication_style": ["2 style words"], '
            '"knowledge_areas": ["3 areas of expertise"]}'
        )
        # Only spend what is left of the total budget, and skip the call entirely
        # if there is no meaningful time left. The heuristic fallback below is
        # already good enough to onboard with.
        _remaining = _SCAN_TOTAL_BUDGET_S - (_time.monotonic() - _scan_started)
        raw = think_local(prompt, timeout=int(_remaining)) if _remaining >= 5 else None
        if not raw and _remaining < 5:
            log.info("Onboarding scan skipped local-LLM profiling — out of time budget.")
        if raw:
            try:
                txt = raw.strip()
                m = _re.search(r'\{.*\}', txt, _re.DOTALL)  # grab first {...} block
                if m:
                    txt = m.group(0)
                llm_data = json.loads(txt)
                used_local = True
            except Exception as e:
                log.warning(f"Onboarding local-LLM parse failed: {e}")

    # ── Deterministic, cloud-free fallback ───────────────────────────────────
    if not llm_data:
        def _dedupe(xs):
            seen, out = set(), []
            for x in xs:
                if x not in seen:
                    seen.add(x)
                    out.append(x)
            return out
        topics, knowledge = [], []
        for s in skills:
            t, k = _SKILL_PROFILE.get(s, ([], []))
            topics += t
            knowledge += k
        llm_data = {
            "role": "Developer",
            "topics": _dedupe(topics) or ["Software Development", "Tooling"],
            "communication_style": ["Direct", "Technical"],
            "knowledge_areas": _dedupe(knowledge) or ["Local Filesystem", "Version Control"],
        }

    return {
        "name": getpass.getuser(),
        "role": llm_data.get("role") or "Developer",
        "projects": projects[:5],
        "skills": skills[:5],
        "topics": (llm_data.get("topics") or [])[:4],
        "communication_style": (llm_data.get("communication_style") or [])[:3],
        "knowledge_areas": (llm_data.get("knowledge_areas") or [])[:4],
        "scan_engine": "local-llm" if used_local else "heuristic",
    }


@app.get("/api/onboarding/scan")
async def scan_onboarding():
    # Heavy I/O + a possibly-slow local model call — offload off the event loop
    # so the WebSocket and other routes stay responsive during onboarding.
    return await asyncio.to_thread(_onboarding_scan_sync)

@app.get("/settings")
async def get_settings():
    from settings_manager import load_settings
    return _sanitize_settings(load_settings())

@app.post("/generate")
async def post_generate(request: Request):
    from brain import think
    body = await request.json()
    prompt = body.get("prompt", "")
    context = body.get("context", "")
    if not prompt:
        return {"response": ""}
    # Using non-streaming think wrapper if available, or just joining the generator
    from brain import think_stream
    response_chunks = []
    for token in think_stream(prompt, context=context):
        response_chunks.append(token)
    return {"response": "".join(response_chunks)}

@app.post("/settings")
async def post_settings(request: Request):
    from settings_manager import save_settings, load_settings
    body = await request.json()
    
    old_settings = load_settings()
    
    # Merge partial updates with existing settings
    merged = {**old_settings, **body}
    
    # Do not overwrite with masked keys
    for key in _MASKED_KEYS:
        if merged.get(key) == "sk-****":
            merged[key] = old_settings.get(key, "")
            
    save_settings(merged)
    
    # Reload settings in core
    core.settings = load_settings()
    
    # Propagate to VADListener dynamically
    if hasattr(core, "vad") and core.vad:
        core.vad.settings = core.settings
        # Sensitivity: Map higher slider values (0.0 to 1.0) to lower noise thresholds (0.05 to 0.005)
        # So slider at 1.0 -> 0.005 (very sensitive), slider at 0.0 -> 0.05 (not sensitive)
        sensitivity = float(core.settings.get("vad_sensitivity", 0.5))
        core.vad.SILENCE_THRESHOLD = max(0.005, 0.05 - (sensitivity * 0.045))
        log.info(f"VAD sensitivity threshold updated dynamically to: {core.vad.SILENCE_THRESHOLD:.4f}")
        
    # If privacy shield just got enabled, kick off model loading
    if body.get("privacy_mirror_enabled") and not old_settings.get("privacy_mirror_enabled"):
        try:
            from privacy_mirror import start_model_loading
            start_model_loading()
            log.info("Privacy Shield enabled — PII model loading started")
        except Exception as e:
            log.warning(f"Could not start PII model loading: {e}")

    broadcast("settings_updated", _sanitize_settings(core.settings))
    return {"success": True}

@app.get("/voices")
async def get_voices():
    from voice_id import get_profiles
    return get_profiles()

@app.post("/voices/enroll")
async def enroll_voice(request: Request):
    from voice_id import enroll_speaker
    body = await request.json()
    wav_path = body.get("wav_path")
    name = body.get("name")
    if not wav_path or not name:
        raise HTTPException(status_code=400, detail="wav_path and name required")
    enroll_speaker(wav_path, name)
    return {"success": True}

@app.delete("/voices/{name}")
async def delete_voice(name: str):
    from voice_id import delete_profile
    delete_profile(name)
    return {"success": True}

@app.post("/clipboard/clear")
async def post_clear_clipboard():
    clear_clipboard_data()
    # Reset state to idle in the frontend
    broadcast("state", {"value": "idle"})
    return {"success": True}

@app.post("/mic/toggle")
async def toggle_mic_endpoint():
    muted = core.toggle_mic()
    return {"muted": muted}

@app.post("/incognito/toggle")
async def toggle_incognito_endpoint():
    active = core.toggle_incognito()
    return {"active": active}

@app.on_event("startup")
async def startup_event():
    global loop
    loop = asyncio.get_running_loop()
    # Register observer callback for notifications
    register_observer_callback(broadcast)
    # Start the clipboard monitor task in the background
    asyncio.create_task(start_clipboard_monitor())
    # Start auto-cleanup scheduler (runs once at startup + every 24h)
    from cleanup_manager import start_cleanup_scheduler
    await asyncio.to_thread(start_cleanup_scheduler)
    # If Privacy Shield is already enabled (persisted from a prior session), kick off
    # PII model loading now — otherwise redact_text() silently stays on the regex
    # fallback until the user toggles the setting off/on. start_model_loading() spawns
    # a daemon thread, so this never blocks startup.
    if core.settings.get("privacy_mirror_enabled", True):
        try:
            from privacy_mirror import start_model_loading
            start_model_loading()
            log.info("Privacy Shield enabled at startup — PII model loading started")
        except Exception as e:
            log.warning(f"Could not start PII model loading at startup: {e}")
    # Periodically re-sync the local vault snapshot (if enabled) so a crash
    # doesn't leave the .vault file far behind the live plaintext db.
    asyncio.create_task(_vault_sync_loop())
    log.info(f"Primnox v{APP_VERSION} Startup Complete — Event loop, clipboard monitor and cleanup scheduler initialized.")


async def _vault_sync_loop():
    import local_vault
    from memory import DB_PATH
    while True:
        await asyncio.sleep(600)  # every 10 minutes
        try:
            await asyncio.to_thread(local_vault.sync_vault, DB_PATH)
        except Exception as e:
            log.error(f"Vault periodic sync failed: {e}")


@app.on_event("shutdown")
async def shutdown_event():
    import local_vault
    from memory import DB_PATH
    try:
        await asyncio.to_thread(local_vault.lock_vault, DB_PATH)
    except Exception as e:
        log.error(f"Vault lock on shutdown failed: {e}")


@app.post("/api/cleanup")
async def trigger_cleanup():
    """Manually trigger a cleanup pass. Returns what was deleted."""
    from cleanup_manager import run_cleanup
    result = await asyncio.to_thread(run_cleanup)
    return result


@app.get("/api/storage")
async def storage_info():
    """Return storage usage for meetings, memories, TTS cache."""
    import sqlite3
    from cleanup_manager import _meetings_dir
    from memory import DB_PATH
    from pathlib import Path
    import tempfile

    info: dict = {}

    # Meetings
    meetings_dir = _meetings_dir()
    if meetings_dir.exists():
        folders = [f for f in meetings_dir.iterdir() if f.is_dir()]
        total_mb = sum(
            sum(fi.stat().st_size for fi in f.rglob("*") if fi.is_file())
            for f in folders
        ) / (1024 * 1024)
        info["meetings"] = {"count": len(folders), "size_mb": round(total_mb, 1)}
    else:
        info["meetings"] = {"count": 0, "size_mb": 0}

    # Memories
    try:
        with sqlite3.connect(DB_PATH) as conn:
            c = conn.cursor()
            c.execute("SELECT COUNT(*) FROM memories")
            mem_count = c.fetchone()[0]
        db_size = Path(DB_PATH).stat().st_size / (1024 * 1024) if Path(DB_PATH).exists() else 0
        info["memories"] = {"count": mem_count, "size_mb": round(db_size, 2)}
    except Exception:
        info["memories"] = {"count": 0, "size_mb": 0}

    return info

@app.get("/api/meetings")
async def list_meetings():
    """List all meeting recording folders with metadata."""
    from cleanup_manager import _meetings_dir
    from datetime import datetime
    import json

    meetings_dir = _meetings_dir()
    if not meetings_dir.exists():
        return {"meetings": []}

    items = []
    for folder in sorted(meetings_dir.iterdir(), key=lambda f: f.stat().st_mtime, reverse=True):
        if not folder.is_dir():
            continue
        try:
            files  = list(folder.rglob("*"))
            files  = [f for f in files if f.is_file()]
            size_b = sum(f.stat().st_size for f in files)
            mtime  = datetime.fromtimestamp(folder.stat().st_mtime)

            # Try to read a summary if one exists
            summary_text = ""
            for candidate in ("summary.txt", "summary.md", "transcript.txt"):
                p = folder / candidate
                if p.exists():
                    try:
                        summary_text = p.read_text(encoding="utf-8", errors="replace")[:300]
                    except Exception:
                        pass
                    break

            # Collect audio/video files
            media_exts  = {".mp3", ".wav", ".mp4", ".m4a", ".ogg", ".webm"}
            media_files = [f.name for f in files if f.suffix.lower() in media_exts]

            items.append({
                "name":        folder.name,
                "date":        mtime.isoformat(timespec="seconds"),
                "size_mb":     round(size_b / (1024 * 1024), 2),
                "file_count":  len(files),
                "media_files": media_files,
                "summary":     summary_text,
            })
        except Exception as e:
            log.warning(f"Could not read meeting folder {folder.name}: {e}")

    return {"meetings": items}


@app.delete("/api/meetings/{folder_name}")
async def delete_meeting(folder_name: str):
    """Delete a specific meeting folder by name."""
    import re, shutil
    from cleanup_manager import _meetings_dir

    # Sanitise — block path traversal. The regex blocks separators, but '.'/'..'
    # pass it (dot is allowed), and `_meetings_dir() / '..'` would resolve to the
    # PARENT dir (Documents/Primnox, holding the DBs) and rmtree it. Reject dotted
    # specials and require the resolved target to be a DIRECT CHILD of the meetings dir.
    if (not re.match(r'^[\w\-\. ]+$', folder_name)
            or folder_name in (".", "..") or ".." in folder_name):
        raise HTTPException(status_code=400, detail="Invalid folder name")

    base = _meetings_dir().resolve()
    target = (base / folder_name).resolve()
    if target.parent != base:
        raise HTTPException(status_code=400, detail="Invalid folder name")
    if not target.exists() or not target.is_dir():
        raise HTTPException(status_code=404, detail="Meeting not found")

    try:
        shutil.rmtree(target)
        log.info(f"Deleted meeting folder on request: {folder_name}")
        return {"deleted": folder_name}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/feedback")
async def receive_feedback(request: Request, background_tasks: BackgroundTasks):
    import time
    from pathlib import Path
    try:
        body = await request.json()
    except Exception:
        raise HTTPException(status_code=400, detail="Invalid JSON")

    category = body.get("category", "General")
    content = body.get("content", "").strip()
    contact = body.get("contact", "").strip()

    if not content:
        raise HTTPException(status_code=400, detail="Feedback content is required")

    log.info(f"User Feedback Received [{category}]: {content}")

    # ── Save locally as backup ────────────────────────────────────────────────
    feedback_file = Path(__file__).parent / "feedback.json"
    feedback_entry = {
        "timestamp": time.time(),
        "category": category,
        "content": content,
        "contact": contact or None
    }
    try:
        data = []
        if feedback_file.exists():
            with open(feedback_file, "r", encoding="utf-8") as f:
                data = json.load(f)
        data.append(feedback_entry)
        with open(feedback_file, "w", encoding="utf-8") as f:
            json.dump(data, f, indent=2)
    except Exception as e:
        log.error(f"Failed to save feedback locally: {e}")

    # ── Send to Discord in background (non-blocking) ──────────────────────────
    if FEEDBACK_DISCORD_WEBHOOK:
        background_tasks.add_task(_send_feedback_to_discord, category, content, contact)

    return {"status": "ok"}


def _send_feedback_to_discord(category: str, content: str, contact: str):
    """Send a formatted feedback embed to Discord webhook. Runs in background."""
    import requests as req
    import time

    CATEGORY_COLORS = {
        "Bug":     0xef4444,   # red
        "Feature": 0x6366f1,   # indigo
        "General": 0x3b82f6,   # blue
    }
    CATEGORY_EMOJI = {
        "Bug":     "🐛",
        "Feature": "✨",
        "General": "💬",
    }

    color = CATEGORY_COLORS.get(category, 0x3b82f6)
    emoji = CATEGORY_EMOJI.get(category, "💬")

    fields = []
    if contact:
        fields.append({"name": "Contact", "value": contact, "inline": True})
    fields.append({"name": "Category", "value": f"{emoji} {category}", "inline": True})

    payload = {
        "embeds": [{
            "title": f"{emoji} New Primnox Feedback",
            "description": content,
            "color": color,
            "fields": fields,
            "footer": {"text": "Primnox Feedback System"},
            "timestamp": __import__('datetime').datetime.utcnow().isoformat()
        }]
    }

    try:
        resp = req.post(FEEDBACK_DISCORD_WEBHOOK, json=payload, timeout=10)
        if resp.status_code not in (200, 204):
            log.warning(f"Discord webhook returned {resp.status_code}: {resp.text}")
        else:
            log.info("Feedback delivered to Discord.")
    except Exception as e:
        log.error(f"Failed to send feedback to Discord: {e}")



import threading
from profiler import run_background_profiler

@app.post("/api/profile/analyze")
async def analyze_profile():
    # Run in background to avoid blocking
    threading.Thread(target=run_background_profiler).start()
    return {"status": "started"}


from emotion_agent import run_emotion_analysis

@app.post("/api/emotion/analyze")
async def analyze_emotion():
    threading.Thread(target=run_emotion_analysis).start()
    return {"status": "started"}


# ── Reminders ──────────────────────────────────────────────────────────────────
from reminder_manager import list_reminders, add_reminder as _add_reminder

@app.get("/api/reminders")
async def get_reminders():
    """Return all pending reminders with seconds remaining."""
    return {"reminders": list_reminders()}

@app.post("/api/reminders")
async def create_reminder(request: Request):
    """Manually set a reminder: {message, delay_secs}"""
    body = await request.json()
    message = body.get("message", "reminder")
    delay = int(body.get("delay_secs", 300))
    _add_reminder(message, delay)
    return {"status": "ok", "message": message, "delay_secs": delay}


@app.delete("/api/reminders/{reminder_id}")
async def delete_reminder(reminder_id: str):
    """Cancel a pending reminder by its stable id."""
    from reminder_manager import cancel_reminder_by_id
    ok = cancel_reminder_by_id(reminder_id)
    return {"status": "ok" if ok else "not_found"}


# ── Backup ─────────────────────────────────────────────────────────────────────
@app.post("/api/backup")
async def trigger_backup(background_tasks: BackgroundTasks):
    """
    Zip memory.db + chat.db + settings.json into
    %APPDATA%/primnox_extension/backups/backup_YYYYMMDD_HHMMSS.zip
    and broadcast backup_complete / backup_failed.
    """
    background_tasks.add_task(_run_backup)
    return {"status": "started"}

def _run_backup():
    import zipfile, time as _time
    from settings_manager import get_appdata_dir
    try:
        appdata = get_appdata_dir()
        backups_dir = appdata / "backups"
        backups_dir.mkdir(parents=True, exist_ok=True)

        stamp = _time.strftime("%Y%m%d_%H%M%S")
        zip_path = backups_dir / f"backup_{stamp}.zip"

        files_to_backup = [
            appdata / "memory.db",
            appdata / "chat.db",
            appdata / "settings.json",
        ]

        with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
            for fp in files_to_backup:
                if fp.exists():
                    zf.write(fp, fp.name)

        # Keep only the 5 most recent backups
        all_backups = sorted(backups_dir.glob("backup_*.zip"), reverse=True)
        for old in all_backups[5:]:
            try:
                old.unlink()
            except Exception:
                pass

        size_kb = zip_path.stat().st_size // 1024
        log.info(f"Backup complete: {zip_path.name} ({size_kb} KB)")
        broadcast("backup_complete", {"path": str(zip_path), "size_kb": size_kb})
    except Exception as e:
        log.error(f"Backup failed: {e}")
        broadcast("backup_failed", {"error": str(e)})


# ── Skills ─────────────────────────────────────────────────────────────────────

@app.get("/api/skills")
async def get_skills():
    """Return all registered skills with name, description, and trigger words."""
    from skills.skill_router import list_skills
    return {"skills": list_skills()}


# ── Memory search ──────────────────────────────────────────────────────────────

@app.get("/api/memories")
async def get_memories_list():
    """Return all stored memories (non-stale)."""
    from memory import list_memories
    return {"memories": list_memories()}


@app.get("/api/memories/search")
async def search_memories_api(q: str = "", limit: int = 20):
    """Semantic full-text search over stored memories."""
    if not q.strip():
        from memory import list_memories
        return {"memories": list_memories()[:limit], "query": q}
    from memory import search_memories
    results = search_memories(q.strip(), limit=limit)
    return {"memories": results, "query": q}


@app.post("/api/memories")
async def create_memory_api(request: Request):
    """Manually add a memory. Body: {text, category}"""
    from memory import add_memory
    body = await request.json()
    text = body.get("text", "").strip()
    category = body.get("category", "personal")
    if not text:
        raise HTTPException(status_code=400, detail="text is required")
    ok = add_memory(text, category=category)
    return {"success": ok, "duplicate": not ok}


@app.delete("/api/memories/{key}")
async def delete_memory_api(key: str):
    """Delete a memory by its key."""
    from memory import delete_memory
    delete_memory(key)
    return {"status": "ok"}


# ── Notes search ───────────────────────────────────────────────────────────────

@app.get("/api/notes/search")
async def search_notes_api(q: str = "", limit: int = 20):
    """Full-text search over notes titles and bodies."""
    from notes_manager import search_notes, get_notes
    if not q.strip():
        notes = get_notes()[:limit]
        return {"notes": notes, "query": q}
    results = search_notes(q.strip(), limit=limit)
    return {"notes": results, "query": q}


# ── Research / web search ──────────────────────────────────────────────────────

@app.post("/api/research/deep")
async def deep_research_stream(request: Request):
    """
    Deep multi-round research with full page reading and gap analysis.
    Returns a Server-Sent Events stream of progress + final report.
    Body: { query: str, depth: 1|2|3 }
    """
    from fastapi.responses import StreamingResponse
    body  = await request.json()
    query = body.get("query", "").strip()
    depth = int(body.get("depth", 2))

    if not query:
        return {"error": "No query provided"}

    async def generate():
        try:
            from research_engine import DeepResearchEngine
            engine = DeepResearchEngine(query, depth)
            async for event in engine.run():
                yield f"data: {json.dumps(event)}\n\n"
        except Exception as e:
            log.error(f"Deep research error: {e}")
            yield f"data: {json.dumps({'type': 'error', 'text': str(e)})}\n\n"
        yield "data: [DONE]\n\n"

    return StreamingResponse(
        generate(),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"},
    )


@app.get("/api/research")
async def research_search(q: str = ""):
    """
    Run a DuckDuckGo web search and return structured results for the
    ResearchView frontend. Each result has title, url, body (snippet).
    """
    if not q.strip():
        return {"results": [], "query": q, "summary": ""}

    def _search():
        try:
            from ddgs import DDGS
            ddgs = DDGS()
            return list(ddgs.text(q.strip(), max_results=8))
        except Exception as e:
            log.error(f"Research search failed: {e}")
            return []

    results = await asyncio.to_thread(_search)

    # Ask LLM to produce a one-paragraph synthesis of the results
    summary = ""
    if results:
        def _summarize():
            from brain import think
            snippets = "\n".join(
                f"- {r.get('title', '')}: {r.get('body', '')[:200]}"
                for r in results[:5]
            )
            resp = think(
                f"Synthesize these web search results for the query '{q}' into a 2-3 sentence summary:\n{snippets}",
                system_override=(
                    "You are a research assistant. Write a concise, factual synthesis. "
                    "No preamble, no 'Here is a summary' — just the synthesis."
                )
            )
            return resp.get("choices", [{}])[0].get("message", {}).get("content", "").strip()

        try:
            summary = await asyncio.to_thread(_summarize)
        except Exception as e:
            log.warning(f"Research summarize failed: {e}")

    structured = [
        {
            "title": r.get("title", ""),
            "url": r.get("href", r.get("url", "")),
            "body": r.get("body", "")[:300],
        }
        for r in results
    ]
    return {"results": structured, "query": q, "summary": summary}

if __name__ == "__main__":
    # Install crash capture before anything else can fail, so an exception
    # during startup lands in the log instead of a stderr nobody reads.
    from logger import install_crash_handlers, capture_stdlib_logging, log_environment
    install_crash_handlers()
    capture_stdlib_logging()
    log_environment()

    loop_type = "asyncio"
    try:
        import uvloop
        asyncio.set_event_loop_policy(uvloop.EventLoopPolicy())
        loop_type = "uvloop"
        log.info("Forced high-performance event loop policy via uvloop.")
    except ImportError:
        log.info("uvloop not supported on this platform. Falling back to default asyncio loop policy.")

    config = uvicorn.Config(app=app, host="127.0.0.1", port=4009, loop=loop_type)
    server = uvicorn.Server(config)
    
    # Run uvicorn natively, it manages the loop
    server.run()

