"""Read-only previews of stored assets — one shape per family of format.

Why this is a server concern. The browser renders PDFs and images natively and
needs nothing from us; Word, Excel and PowerPoint it cannot read at all. The
libraries that *can* read them are already installed here — the sandbox writes
these formats with the very same ones — so the alternative is shipping a
second, JavaScript copy of each parser to the client. This is the cheaper half
of the job by a wide margin.

Nothing here can modify an asset. Every function reads bytes and returns a
description of them; there is no write path to reach even by accident.

Every parser is optional. A missing library downgrades one format to
`unsupported`, which the UI renders as "no preview, download it" — never as a
broken page.
"""
from __future__ import annotations

import csv
import io
import json
from pathlib import Path

# Generous enough for a real document, bounded enough that a pathological file
# cannot hand the client a hundred megabytes of JSON.
MAX_TEXT_CHARS = 200_000
MAX_ROWS_PER_SHEET = 500
MAX_COLS = 40
MAX_BLOCKS = 2_000
MAX_SLIDES = 200

TEXTUAL = {".txt", ".md", ".markdown", ".log", ".py", ".js", ".ts", ".tsx", ".jsx",
           ".json", ".html", ".htm", ".css", ".xml", ".yml", ".yaml", ".ini",
           ".toml", ".sql", ".sh", ".bat", ".ps1", ".rst", ".c", ".h", ".cpp",
           ".java", ".rs", ".go", ".rb", ".php"}
TABULAR = {".csv", ".tsv"}
# Renderable in a frame rather than read as source. Kept out of TEXTUAL, which
# still owns `.css`, `.js` and the rest — those are source when you open them,
# whereas a page is a thing to look at.
WEB = {".html", ".htm"}
IMAGES = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".svg", ".ico"}


def describe(asset: dict) -> dict:
    """What the client needs to display this asset, and nothing more."""
    name = asset.get("original_name") or ""
    suffix = Path(name).suffix.lower()
    base = {"asset_id": asset.get("id"), "name": name,
            "bytes": asset.get("bytes"), "mime": asset.get("mime")}

    # Handed straight to the browser, which already knows these. Sending the
    # bytes through here would only add a copy.
    if suffix == ".pdf" or (asset.get("mime") or "") == "application/pdf":
        return {**base, "kind": "pdf"}
    if suffix in IMAGES or (asset.get("mime") or "").startswith("image/"):
        return {**base, "kind": "image"}

    path = Path(asset.get("path") or "")
    if not path.is_file():
        return {**base, "kind": "missing"}

    try:
        if suffix in TABULAR:
            # The display name comes from the asset, never from the path: files
            # are stored content-addressed, so `path.stem` is a sha256 and the
            # viewer labelled the sheet with 64 hex characters.
            return {**base, **_delimited(path, "\t" if suffix == ".tsv" else ",",
                                         Path(name).stem)}
        if suffix == ".xlsx" or suffix == ".xlsm":
            return {**base, **_spreadsheet(path)}
        if suffix == ".docx":
            return {**base, **_word(path)}
        if suffix == ".pptx":
            # `_image_blobs` is raw picture bytes, kept out of the payload: it
            # is not JSON-serialisable, and a deck of photographs would be tens
            # of megabytes on a request that only needs to describe the deck.
            # The viewer fetches each picture from /assets/{id}/slide-image/{n}.
            deck = _slides(path)
            deck.pop("_image_blobs", None)
            return {**base, **deck}
        # Before TEXTUAL, which would otherwise catch it. A generated deck is
        # meant to be WATCHED, and `.html` sitting in TEXTUAL meant the
        # `frontend-slides` skill's whole output — an animated 16:9
        # presentation — arrived in the viewer as its own markup. The source is
        # still carried so the viewer can offer it; it is just no longer the
        # only thing on offer.
        if suffix in WEB:
            return {**base, **_web(path, suffix)}
        if suffix in TEXTUAL:
            return {**base, **_text(path, suffix)}
        if suffix in (".db", ".sqlite", ".sqlite3"):
            return {**base, **_sqlite(path)}
    except Exception as exc:
        # A file that cannot be parsed is still a file the user may download.
        # Reporting why beats rendering an empty viewer.
        return {**base, "kind": "unreadable", "error": f"{type(exc).__name__}: {exc}"}

    return {**base, "kind": "unsupported"}


# ── text ─────────────────────────────────────────────────────────────────────
def _text(path: Path, suffix: str) -> dict:
    raw = path.read_bytes()[: MAX_TEXT_CHARS * 2]
    text = raw.decode("utf-8", "replace")
    truncated = len(text) > MAX_TEXT_CHARS
    language = {".md": "markdown", ".markdown": "markdown"}.get(
        suffix, suffix.lstrip(".") or "text")
    return {"kind": "text", "language": language,
            "text": text[:MAX_TEXT_CHARS], "truncated": truncated}


# ── html ─────────────────────────────────────────────────────────────────────
def _web(path: Path, suffix: str) -> dict:
    """A page, plus its source.

    Both, because they answer different questions: the frame shows what the deck
    looks like, and the source is how you check what it actually does before
    trusting it. Neither alone is enough for a file a language model wrote.

    The bytes are NOT inlined here — the viewer points a frame at the ordinary
    download URL, so a deck with embedded images streams rather than arriving
    base64'd inside a JSON envelope.
    """
    raw = path.read_bytes()[: MAX_TEXT_CHARS * 2]
    text = raw.decode("utf-8", "replace")
    return {"kind": "web", "language": suffix.lstrip(".") or "html",
            "text": text[:MAX_TEXT_CHARS],
            "truncated": len(text) > MAX_TEXT_CHARS}


# ── csv / tsv ────────────────────────────────────────────────────────────────
def _delimited(path: Path, delimiter: str, display_name: str) -> dict:
    text = path.read_bytes().decode("utf-8", "replace")
    reader = csv.reader(io.StringIO(text), delimiter=delimiter)
    rows: list[list[str]] = []
    total = 0
    for row in reader:
        total += 1
        if len(rows) < MAX_ROWS_PER_SHEET:
            rows.append([c for c in row[:MAX_COLS]])
    header, body = (rows[0], rows[1:]) if rows else ([], [])
    return {"kind": "sheets", "sheets": [{
        "name": display_name, "header": header, "rows": body,
        "total_rows": max(total - 1, 0),
        "truncated": total - 1 > len(body),
    }]}


# ── xlsx ─────────────────────────────────────────────────────────────────────
def _spreadsheet(path: Path) -> dict:
    import openpyxl

    # Handed the bytes, not the path. Assets are stored content-addressed —
    # the file on disk is named after its sha256 and has no extension — and
    # openpyxl decides the format from the filename, so it refuses a perfectly
    # good workbook it cannot see a `.xlsx` on. python-docx and python-pptx
    # sniff the zip instead, which is why only this one needed it.
    #
    # read_only keeps a large workbook from being materialised in full, and
    # data_only asks for the last computed value rather than the formula —
    # a preview should show what the sheet says, not how it says it.
    wb = openpyxl.load_workbook(io.BytesIO(path.read_bytes()),
                                read_only=True, data_only=True)
    sheets = []
    for ws in wb.worksheets:
        rows: list[list] = []
        total = 0
        for row in ws.iter_rows(values_only=True):
            if row is None or all(v is None for v in row):
                continue
            total += 1
            if len(rows) < MAX_ROWS_PER_SHEET + 1:
                rows.append(["" if v is None else _cell(v) for v in row[:MAX_COLS]])
        header, body = (rows[0], rows[1:]) if rows else ([], [])
        sheets.append({"name": ws.title, "header": header, "rows": body,
                       "total_rows": max(total - 1, 0),
                       "truncated": max(total - 1, 0) > len(body)})
    wb.close()
    return {"kind": "sheets", "sheets": sheets}


def _cell(value) -> str:
    if isinstance(value, float) and value.is_integer():
        return str(int(value))
    return str(value)


# ── docx ─────────────────────────────────────────────────────────────────────
def _word(path: Path) -> dict:
    import docx

    document = docx.Document(str(path))
    blocks: list[dict] = []
    for para in document.paragraphs:
        text = (para.text or "").strip()
        if not text:
            continue
        style = (para.style.name or "").lower() if para.style else ""
        if style.startswith("heading"):
            level = "".join(ch for ch in style if ch.isdigit())
            blocks.append({"type": "heading", "level": int(level or 1), "text": text})
        elif style.startswith("list") or style.startswith("bullet"):
            blocks.append({"type": "bullet", "text": text})
        else:
            blocks.append({"type": "paragraph", "text": text})
        if len(blocks) >= MAX_BLOCKS:
            break

    for table in document.tables:
        rows = [[(c.text or "").strip() for c in r.cells[:MAX_COLS]]
                for r in table.rows[:MAX_ROWS_PER_SHEET]]
        if rows:
            blocks.append({"type": "table", "rows": rows})

    return {"kind": "document", "blocks": blocks[:MAX_BLOCKS]}


# ── pptx ─────────────────────────────────────────────────────────────────────
# A deck is a LAYOUT, and the old reader threw the layout away: it collected a
# title and a flat list of lines per slide, so a slide whose whole content was a
# chart and two captions came back as two captions, and a slide that was one
# full-bleed photograph came back empty. Decks generated by this app routinely
# contain images, which meant the viewer showed nothing of what was made.
#
# Every shape now carries its own box, as a FRACTION of the slide rather than in
# EMU. Fractions survive being rendered at any size, which is the whole point —
# the viewer draws into whatever width it has and the deck still looks like
# itself. Pictures are referenced by index and fetched separately, so a
# thirty-slide deck of photographs does not arrive as one enormous JSON payload.

def _emu_box(shape, w: int, h: int) -> dict | None:
    """A shape's position as fractions of the slide. None when unplaceable.

    Some shapes carry no geometry at all — placeholders inheriting from a layout
    are the common case — and guessing a position for those would scatter them
    across the slide at random. They are better dropped than misplaced.
    """
    try:
        left, top = shape.left, shape.top
        width, height = shape.width, shape.height
    except Exception:
        return None
    if None in (left, top, width, height) or not w or not h:
        return None
    return {
        "x": round(left / w, 5), "y": round(top / h, 5),
        "w": round(width / w, 5), "h": round(height / h, 5),
    }


def _run_style(frame) -> dict:
    """Size, weight and alignment of a text frame's first run.

    The first run, not an average: a text box is overwhelmingly one style in a
    generated deck, and reproducing per-run styling would mean shipping a rich
    text model to render a preview.
    """
    size, bold, align = None, False, None
    try:
        for para in frame.paragraphs:
            if para.alignment is not None:
                align = str(para.alignment).split(".")[-1].split(" ")[0].lower()
            for run in para.runs:
                if run.font.size is not None:
                    size = run.font.size.pt
                bold = bool(run.font.bold)
                raise StopIteration
    except StopIteration:
        pass
    except Exception:
        pass
    return {"size": size, "bold": bold, "align": align}


def _slides(path: Path) -> dict:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    deck = Presentation(str(path))
    sw = int(deck.slide_width or 0)
    sh = int(deck.slide_height or 0)

    slides, images = [], []
    for i, slide in enumerate(deck.slides, start=1):
        title = ""
        shapes: list[dict] = []

        for shape in slide.shapes:
            box = _emu_box(shape, sw, sh)

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    blob = shape.image.blob
                    ext = shape.image.ext or "png"
                except Exception:
                    continue                      # a linked image with no bytes
                images.append({"blob": blob, "content_type": f"image/{ext}"})
                shapes.append({"type": "image", "ref": len(images) - 1,
                               "box": box,
                               "alt": (getattr(shape, "name", "") or "").strip()})
                continue

            if not getattr(shape, "has_text_frame", False):
                continue
            text = (shape.text_frame.text or "").strip()
            if not text:
                continue

            is_title = shape == getattr(slide.shapes, "title", None)
            if is_title and not title:
                title = text
            shapes.append({"type": "text", "text": text, "box": box,
                           "title": is_title, **_run_style(shape.text_frame)})

        # Still provide the flat reading order. It is what a screen reader
        # follows, what search would index, and the fallback for any shape the
        # layout pass could not place.
        lines = [l.strip() for s in shapes if s["type"] == "text"
                 for l in s["text"].splitlines() if l.strip()]
        if not title and lines:
            title = lines[0]

        notes = ""
        try:
            if slide.has_notes_slide:
                notes = (slide.notes_slide.notes_text_frame.text or "").strip()
        except Exception:
            pass

        slides.append({"index": i, "title": title,
                       "lines": [l for l in lines if l != title],
                       "shapes": shapes, "notes": notes})
        if len(slides) >= MAX_SLIDES:
            break

    return {"kind": "slides", "slides": slides,
            "aspect": round(sw / sh, 5) if sw and sh else 16 / 9,
            "images": len(images), "_image_blobs": images}


def slide_image(asset_path: Path, ref: int) -> tuple[bytes, str] | None:
    """One picture out of a deck, by the index the preview handed out.

    Re-read rather than cached: previews are already re-derived on every open,
    and holding decks of image bytes in memory to save a parse is the wrong
    trade for a viewer someone opens occasionally.
    """
    try:
        data = _slides(asset_path)
    except Exception:
        return None
    blobs = data.get("_image_blobs") or []
    if ref < 0 or ref >= len(blobs):
        return None
    entry = blobs[ref]
    return entry["blob"], entry["content_type"]


# ── sqlite ───────────────────────────────────────────────────────────────────
def _sqlite(path: Path) -> dict:
    """Opened read-only through a URI, so previewing a database can never
    write to one — not even a journal file beside it."""
    import sqlite3

    conn = sqlite3.connect(f"file:{path}?mode=ro", uri=True)
    try:
        names = [r[0] for r in conn.execute(
            "SELECT name FROM sqlite_master WHERE type='table' "
            "AND name NOT LIKE 'sqlite_%' ORDER BY name")]
        sheets = []
        for table in names:
            # `table` comes straight from sqlite_master — i.e. from whatever
            # file the user (or a document they were sent) uploaded, not from
            # our own schema. SQLite has no placeholder syntax for
            # identifiers, only values, so the only safe way to interpolate
            # one is to quote-and-escape it ourselves: a table named
            # `x"; ATTACH DATABASE '...' AS y; --` would otherwise break out
            # of the quoted identifier and run as SQL against this
            # connection. Doubling embedded `"` is the SQL-standard escape
            # for a quoted identifier.
            quoted = table.replace('"', '""')
            cursor = conn.execute(f'SELECT * FROM "{quoted}" LIMIT {MAX_ROWS_PER_SHEET}')
            header = [d[0] for d in cursor.description][:MAX_COLS]
            rows = [["" if v is None else str(v) for v in row[:MAX_COLS]]
                    for row in cursor.fetchall()]
            total = conn.execute(f'SELECT COUNT(*) FROM "{quoted}"').fetchone()[0]
            sheets.append({"name": table, "header": header, "rows": rows,
                           "total_rows": total, "truncated": total > len(rows)})
    finally:
        conn.close()
    return {"kind": "sheets", "sheets": sheets}


def to_json(value) -> str:
    return json.dumps(value, default=str)
