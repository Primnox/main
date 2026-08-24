"""Themed document builders, shipped INTO the sandbox.

This file is copied verbatim into every execution directory as
`primnox_docs.py`, so generated code can `import primnox_docs` with no install
step and no network.

Why it exists. Left to raw python-pptx and reportlab, a model produces the
library defaults: white slides, Calibri, black Times paragraphs. Asking a model
to style them instead means asking it to write forty more lines of colour and
geometry — and the measurement that prompted this said a 7B writes short
scripts and gets long ones wrong. So the styling lives here, tested once, and
the model writes four lines.

The palettes are the app's own, by name, so a deck generated in `midnight`
matches the interface that produced it.
"""
from __future__ import annotations

# Lifted from frontend/src/styles/themes.css. Same names, same values — a
# theme is a single idea in this product, not one per output format.
THEMES: dict[str, dict] = {
    "signature": {"bg": "070707", "text": "F0EDE6", "muted": "8C8A84", "primary": "C3C0FF", "accent": "FFB695", "dark": True},
    "void":      {"bg": "000000", "text": "FFFFFF", "muted": "8C8C8C", "primary": "FFFFFF", "accent": "FF4F2B", "dark": True},
    "carbon":    {"bg": "0E1014", "text": "E3E7EE", "muted": "858B96", "primary": "5FD8FF", "accent": "FFCC66", "dark": True},
    "midnight":  {"bg": "060A1A", "text": "E7EAF7", "muted": "878CA1", "primary": "6EA8FF", "accent": "FF7AB8", "dark": True},
    "ember":     {"bg": "140F0B", "text": "F6ECE0", "muted": "908779", "primary": "F5A524", "accent": "FF6B5E", "dark": True},
    "phosphor":  {"bg": "040A06", "text": "D4E8D6", "muted": "7C8C7E", "primary": "4ADE80", "accent": "D9F99D", "dark": True},
    "paper":     {"bg": "F6F4EF", "text": "100F0E", "muted": "6B6862", "primary": "4B3FBF", "accent": "B8501F", "dark": False},
    "clinical":  {"bg": "FFFFFF", "text": "0A0A0A", "muted": "6B6B6B", "primary": "1447E6", "accent": "D81B4A", "dark": False},
    # sand's muted was 7A7266 and mono's was 6E6E6E: 3.79:1 and 4.51:1 against
    # their own backgrounds, so captions and footers on the light themes sat at
    # or under the WCAG AA floor. Darkened to 4.8:1 and 4.79:1 — the same colour
    # a shade down, which nobody will notice except the people who could not
    # read it. `test_every_theme_meets_contrast` holds the line.
    "sand":      {"bg": "ECE5D8", "text": "241D13", "muted": "6A6255", "primary": "1C6B47", "accent": "B0451A", "dark": False},
    "mono":      {"bg": "F2F1EE", "text": "121212", "muted": "6A6A6A", "primary": "121212", "accent": "3F3F3F", "dark": False},
}

DEFAULT_THEME = "midnight"

# Fonts are chosen for "present on Windows" rather than for beauty. A theme
# naming a font the machine lacks silently falls back to something arbitrary,
# which is worse than a plain font chosen on purpose.
# Bahnschrift is Microsoft's DIN — condensed, technical, and confident at large
# sizes, where Segoe UI just reads as an interface. It ships with Windows 10 and
# later, so it needs no install and cannot be silently substituted on the
# machines this runs on. Segoe UI stays for body copy: a display face set at
# 22pt across a paragraph is a headline pretending to be text.
DISPLAY_FONT = "Bahnschrift SemiBold"
BODY_FONT = "Segoe UI"
# Code must never be set in a proportional face — columns stop lining up and a
# diff becomes unreadable. Consolas ships with Windows; Cascadia Mono does not
# on older builds, so it is not the default.
MONO_FONT = "Consolas"


# "light" and "dark" are how the themes are grouped when they are described,
# so they are what a model reaches for. Measured: asked for a light deck, the
# model passed theme='light', the lookup missed, and it silently produced a
# midnight-dark deck. Treating the category as a request for its default member
# is what the caller plainly meant.
ALIASES = {"light": "paper", "dark": "midnight",
           "white": "clinical", "black": "void", "default": DEFAULT_THEME}


def theme(name: str | None = None) -> dict:
    """Resolve a theme by name, then by category, then complain.

    An unknown name still falls back rather than raising — a misspelt theme
    should cost the styling, not the document — but it says so on stdout. A
    silent fallback means the user asked for one thing, got another, and
    nothing anywhere records that it happened.
    """
    wanted = (name or DEFAULT_THEME).strip().lower()
    if wanted in THEMES:
        return THEMES[wanted]
    if wanted in ALIASES:
        return THEMES[ALIASES[wanted]]
    print(f"primnox_docs: no theme called {wanted!r}; using {DEFAULT_THEME}. "
          f"Valid themes: {', '.join(THEMES)}.")
    return THEMES[DEFAULT_THEME]


def _rgb(hex6: str):
    from pptx.dml.color import RGBColor

    return RGBColor.from_string(hex6.lstrip("#").upper())


def mix(a: str, b: str, t: float) -> str:
    """Blend two hex colours, `t` of the way from a to b."""
    a, b = a.lstrip("#"), b.lstrip("#")
    return "".join(
        f"{round(int(a[i:i + 2], 16) * (1 - t) + int(b[i:i + 2], 16) * t):02X}"
        for i in (0, 2, 4)
    )


def _luminance(hex6: str) -> float:
    hex6 = hex6.lstrip("#")
    channels = []
    for i in (0, 2, 4):
        c = int(hex6[i:i + 2], 16) / 255
        channels.append(c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4)
    return 0.2126 * channels[0] + 0.7152 * channels[1] + 0.0722 * channels[2]


def contrast_ratio(a: str, b: str) -> float:
    """WCAG 2.1 contrast ratio, 1.0 to 21.0.

    Here because a palette is not only an aesthetic choice: a caller can pass
    any two colours, and muted-on-background is exactly where a deck quietly
    becomes unreadable on a projector. 4.5:1 is the AA floor for body text;
    3.0:1 is the floor for large text.
    """
    la, lb = _luminance(a), _luminance(b)
    hi, lo = max(la, lb), min(la, lb)
    return (hi + 0.05) / (lo + 0.05)


def surfaces(t: dict) -> dict:
    """Derive the missing surface levels from a palette's bg and text.

    A design system needs more than one background. With only `bg` the sole way
    to show that a card is a distinct object is to draw a line around it — and a
    deck where every element is outlined in grey looks like a wireframe of
    itself. Two derived tints replace almost all of those borders:

      surface  a panel, barely lifted off the page
      line     a hairline, present but not competing with the content

    Derived rather than declared, so the ten built-in themes and any palette a
    caller invents all gain them without listing two more colours. Both blend
    toward `text`, so they stay correct on light themes as well as dark.
    """
    out = dict(t)
    out.setdefault("surface", mix(t["bg"], t["text"], 0.055))
    out.setdefault("line", mix(t["bg"], t["text"], 0.16))
    return out


# ── Design system ────────────────────────────────────────────────────────────
# Everything below is in POINTS, and every coordinate is a multiple of GRID.
#
# The old builder worked in inches, so its geometry landed on values like
# 64.8pt and 172.8pt. Those are not wrong to look at, but they mean the deck
# has no grid at all — nothing can be aligned to anything else on purpose, and
# a caller asking for an 8-point system cannot be satisfied. Points are the
# unit the type is already specified in, so working in them makes the grid and
# the type scale the same conversation.
GRID = 8
SLIDE_W, SLIDE_H = 960, 540          # 16:9 at 96dpi-equivalent points
MARGIN_X = 64
MARGIN_TOP = 56
MARGIN_BOTTOM = 48
CONTENT_W = SLIDE_W - 2 * MARGIN_X   # 832
BODY_TOP = 160                       # below the title block on a content slide

# H1 -> Caption. A model picks a role, never a number; that is what keeps the
# hierarchy intact across a 200-slide deck.
# `mega` is furniture, not text — the ghosted numeral on a section divider. It
# sits in the scale so that "every size on the slide came from the scale" stays
# a checkable property rather than a rule with an exception.
SCALE = {"mega": 160, "h1": 54, "h2": 40, "h3": 30, "body": 22, "caption": 16}
MIN_READABLE_PT = 16                 # never shrink past this, per spec


def leading(size: float) -> float:
    """Line spacing multiplier for a given size.

    Large type needs proportionally LESS leading than small type: at 54pt a
    1.4 multiple opens a gap wide enough to read as two separate lines, while
    at 16pt the same 1.4 is barely enough to keep a wrapped caption legible.
    PowerPoint's default is a flat single space, which is too tight for
    headlines and too tight again for body copy — everything was set solid.
    """
    if size >= SCALE["h1"]:
        return 1.05
    if size >= SCALE["h3"]:
        return 1.15
    if size >= SCALE["body"]:
        return 1.35
    return 1.45


def snap(v: float) -> int:
    """Nearest multiple of GRID. Applied to every coordinate this module emits."""
    return int(round(v / GRID) * GRID)


def _balanced_columns(n: int, most: int) -> int:
    """Column count that leaves the fewest empty cells in the last row.

    Six cards in four columns is a row of four and a row of two, and the two
    gaps read as something missing rather than as space. Three columns gives
    two full rows of three. Ties prefer more columns, because wider items on
    one line beat a taller, narrower grid.
    """
    if n <= 1:
        return 1
    best, best_cost = min(n, most), None
    for cols in range(min(n, most), 0, -1):
        rows = -(-n // cols)
        empty = cols * rows - n
        # Rows FIRST, gaps second. Minimising gaps alone picks a single column
        # for five items — perfectly gapless, and five rows tall on a slide 540
        # points high. A ragged cell costs less than a layout nobody can read.
        cost = (rows, empty)
        if best_cost is None or cost < best_cost:
            best, best_cost = cols, cost
    return best


def _card_value_size(value: str, width_pt: float, card_h: float,
                     *, lines_around: int) -> int:
    """Size for a card's headline number, given what else shares the card.

    Sizing the value against the whole card and hoping is what put a delta line
    outside its own panel: the label above it and the delta below it each take a
    caption line plus leading, and none of that was subtracted first. Here the
    surrounding lines are measured and the value gets what is actually left.
    """
    caption = SCALE["caption"] * leading(SCALE["caption"])
    available = card_h - GRID - 48 - lines_around * caption   # rule + margins
    return _fit_size(value, width_pt, max(available, MIN_READABLE_PT), start="h2")


def _fit_size(text: str, width_pt: float, height_pt: float, start: str = "body") -> int:
    """Largest size in the scale at which `text` still fits, floored at 16pt.

    Shrinking below the floor is refused rather than allowed: the spec says
    never shrink below readability, and silently producing 9pt body copy is a
    worse outcome than letting the box be tight.
    """
    order = ["h1", "h2", "h3", "body", "caption"]
    start_i = order.index(start) if start in order else order.index("body")
    for role in order[start_i:]:
        size = SCALE[role]
        if size < MIN_READABLE_PT:
            break
        chars_per_line = max(1, int(width_pt / (size * 0.5)))
        lines = max(1, -(-len(text) // chars_per_line))
        # The same leading the text will actually be set with. A fitter using a
        # flat 1.25 while the renderer used 1.45 would choose a size that then
        # overflowed by a fifth on any caption that wrapped.
        if lines * size * leading(size) <= height_pt:
            return size
    return MIN_READABLE_PT


# ── Slides ───────────────────────────────────────────────────────────────────
class Deck:
    """A themed 16:9 deck on a strict 8-point grid.

        d = Deck("volcanoes.pptx", theme="midnight", title="Types of Volcano")
        d.bullets("Shield", ["Gentle slopes", "Hawaii"])
        d.kpi("At a glance", [("Height", "4,207 m", "Mauna Kea"),
                              ("Eruptions", "33", "since 1843")])
        d.save()

    Every layout is one call. That is the constraint the whole file exists to
    satisfy: a 7B writes short scripts correctly and long ones wrong, so a deck
    it can actually produce is one where a KPI row costs a line, not thirty.

    A caller may supply its own palette instead of a theme name:

        Deck("x.pptx", palette={"bg": "0B1220", "text": "F8FAFC",
                                "primary": "22D3EE", "accent": "A3E635",
                                "muted": "94A3B8"})
    """

    #: Layout names, for callers that want to enumerate what exists.
    LAYOUTS = ("hero", "section", "bullets", "two_column", "compare", "bento",
               "kpi", "timeline", "process", "table", "chart", "code", "quote",
               "matrix", "appendix")

    def __init__(self, filename: str, theme: str | None = None,
                 title: str | None = None, subtitle: str | None = None,
                 palette: dict | None = None, footer: str | None = None,
                 display_font: str = DISPLAY_FONT,
                 body_font: str = BODY_FONT,
                 mono_font: str = MONO_FONT) -> None:
        from pptx import Presentation
        from pptx.util import Pt

        self.filename = filename
        self.t = surfaces(dict(globals()["theme"](theme)))
        if palette:
            # An explicit palette overrides the named theme key by key, so a
            # caller can change one colour without restating the rest. Without
            # this the design system in a brief ("Navy #0F172A, Cyan #22D3EE")
            # is simply unreachable — the old builder could only ever produce
            # one of ten built-in themes.
            self.t.update({k: str(v).lstrip("#").upper() for k, v in palette.items()})
            # Re-derive: a caller who changed `bg` must not keep surfaces mixed
            # from the old one, or the panels sit on a background they were
            # never blended against.
            if "bg" in palette or "text" in palette:
                for key in ("surface", "line"):
                    if key not in palette:
                        self.t.pop(key, None)
                self.t = surfaces(self.t)
        self.footer_text = footer
        self.display_font, self.body_font, self.mono_font = display_font, body_font, mono_font

        self.prs = Presentation()
        self.prs.slide_width = Pt(SLIDE_W)
        self.prs.slide_height = Pt(SLIDE_H)
        self._n = 0
        if title:
            self.hero(title, subtitle)

    # ── primitives ──────────────────────────────────────────────────────────
    def _blank(self, *, chrome: bool = True):
        """A painted blank slide, optionally with the master furniture."""
        from pptx.util import Pt

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = _rgb(self.t["bg"])

        # One consistent mark reads as design; a different flourish per slide
        # reads as noise. Full height, so it bleeds and is off-grid by intent.
        bar = slide.shapes.add_shape(1, Pt(0), Pt(0), Pt(GRID), Pt(SLIDE_H))
        bar.fill.solid()
        bar.fill.fore_color.rgb = _rgb(self.t["primary"])
        bar.line.fill.background()
        bar.shadow.inherit = False

        self._n += 1
        if chrome:
            self._chrome(slide)
        return slide

    def _chrome(self, slide) -> None:
        """Footer and page number, identical on every slide that has them.

        This is the slide master's job in PowerPoint, but a master built through
        python-pptx cannot be edited by the user afterwards without surprises,
        and drawing it costs two shapes. Consistency comes from it being one
        function rather than from the file format.
        """
        # 32pt tall, not 24: a 16pt caption at 1.45 leading needs 23pt, and the
        # frame's own insets take another 7. At 24 the footer overflowed its own
        # box on every slide in the deck.
        if self.footer_text:
            self._text(slide, MARGIN_X, SLIDE_H - 48, CONTENT_W - 80, 32,
                       self.footer_text, SCALE["caption"], self.t["muted"])
        self._text(slide, SLIDE_W - MARGIN_X - 64, SLIDE_H - 48, 64, 32,
                   str(self._n), SCALE["caption"], self.t["muted"], align="right")

    def _text(self, slide, left, top, width, height, text, size, color,
              bold=False, font=None, align=None, mono=False):
        from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
        from pptx.util import Pt

        box = slide.shapes.add_textbox(Pt(snap(left)), Pt(snap(top)),
                                       Pt(snap(width)), Pt(snap(height)))
        frame = box.text_frame
        frame.word_wrap = True
        # NONE, not the python-pptx default of SHAPE_TO_FIT_TEXT: a box that
        # silently grows past its declared geometry breaks the grid it was
        # placed on and can collide with its neighbours at render time. Sizes
        # are chosen by _fit_size before the text is placed instead.
        frame.auto_size = MSO_AUTO_SIZE.NONE
        frame.text = text
        p = frame.paragraphs[0]
        face = font or (self.mono_font if mono else self.body_font)
        # Set on the paragraph AND on every run.
        #
        # `p.font` writes the paragraph's default run properties, which
        # PowerPoint does apply — so a deck styled only that way looks right.
        # But the runs carry no properties of their own, which means anything
        # inspecting the file sees unstyled text: the deck auditor read zero
        # font sizes across a fifteen-slide deck and passed its type-scale and
        # readability checks on an empty set. Styling that cannot be verified
        # is indistinguishable from styling that is not there.
        p.line_spacing = leading(size)
        for target in (p.font, *(r.font for r in p.runs)):
            target.size = Pt(size)
            target.bold = bold
            target.name = face
            target.color.rgb = _rgb(color)
        if align == "right":
            p.alignment = PP_ALIGN.RIGHT
        elif align == "center":
            p.alignment = PP_ALIGN.CENTER
        return frame

    def _title(self, slide, text: str) -> None:
        size = _fit_size(text, CONTENT_W, 64, start="h2")
        self._text(slide, MARGIN_X, MARGIN_TOP, CONTENT_W, 64, text, size,
                   self.t["primary"], bold=True, font=self.display_font)

    def _card(self, slide, left, top, width, height, *, fill=None, border=None,
              rule=None):
        """A panel. Text goes INSIDE its own frame, never in a box laid on top —
        two stacked shapes are an overlap, and an overlap is a defect.

        Filled by default rather than outlined. A grey 1pt border around every
        element makes a deck look like a wireframe of itself; a barely-lifted
        surface says "this is an object" without drawing attention to the
        saying. `rule` adds a short accent bar above the panel — one deliberate
        mark of colour per card, which is what stops a filled panel reading as
        an empty box.
        """
        from pptx.enum.text import MSO_ANCHOR
        from pptx.util import Pt

        left, top, width, height = snap(left), snap(top), snap(width), snap(height)
        if rule:
            bar = slide.shapes.add_shape(1, Pt(left), Pt(top), Pt(40), Pt(GRID))
            bar.fill.solid()
            bar.fill.fore_color.rgb = _rgb(rule)
            bar.line.fill.background()
            bar.shadow.inherit = False
            top += GRID
            height -= GRID

        shape = slide.shapes.add_shape(1, Pt(left), Pt(top), Pt(width), Pt(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(fill or self.t["surface"])
        if border:
            shape.line.color.rgb = _rgb(border)
            shape.line.width = Pt(1)
        else:
            shape.line.fill.background()
        shape.shadow.inherit = False
        frame = shape.text_frame
        frame.word_wrap = True
        frame.margin_left = frame.margin_right = Pt(24)
        frame.margin_top = frame.margin_bottom = Pt(24)
        # Top, not the autoshape default of middle. Centred vertically, a card
        # with two lines sits its content lower than a card with three, so a row
        # of metrics has its numbers on four different baselines — the one thing
        # a dashboard row exists to make comparable.
        frame.vertical_anchor = MSO_ANCHOR.TOP
        return shape

    def _lines(self, frame, items, size, color, *, first=True, bold=False,
               space=12, prefix="", align="left"):
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Pt

        # Alignment is set explicitly on EVERY paragraph, never inherited.
        # An autoshape's text defaults to centred, so the first line of a card
        # came out centred while every line added after it came out left — the
        # KPI values and the matrix headings were centred above their own
        # left-aligned body copy. Nothing in the file is wrong by any check;
        # it just looks broken.
        alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                     "right": PP_ALIGN.RIGHT}[align]
        for i, line in enumerate(items):
            p = frame.paragraphs[0] if (first and i == 0) else frame.add_paragraph()
            p.text = f"{prefix}{line}"
            p.alignment = alignment
            p.line_spacing = leading(size)
            for target in (p.font, *(r.font for r in p.runs)):
                target.size = Pt(size)
                target.bold = bold
                target.name = self.body_font
                target.color.rgb = _rgb(color)
            p.space_after = Pt(space)

    # ── layouts ─────────────────────────────────────────────────────────────
    def hero(self, title: str, subtitle: str | None = None) -> "Deck":
        """Opening or full-bleed statement slide. Headline centred by design."""
        s = self._blank(chrome=False)
        size = _fit_size(title, CONTENT_W, 96, start="h1")
        self._text(s, MARGIN_X, 232, CONTENT_W, 96, title, size,
                   self.t["text"], bold=True, font=self.display_font)
        if subtitle:
            self._text(s, MARGIN_X, 344, CONTENT_W, 40, subtitle,
                       SCALE["body"], self.t["muted"])
        return self

    def section(self, title: str, subtitle: str | None = None,
                number: str | None = None) -> "Deck":
        """A divider. Distinct from `hero` so a deck can open once and still
        break into parts without reusing the cover layout."""
        s = self._blank()
        top = 216
        # The numeral gets its own column rather than sitting behind the title.
        # Ghosting it in the surface tint is not enough on its own: two text
        # boxes on the same ground is an overlap however faint one of them is,
        # and at 160pt it collides with a descender long before a reader would
        # call it texture.
        text_w = CONTENT_W - 352 if number else CONTENT_W
        if number:
            self._text(s, MARGIN_X + CONTENT_W - 320, 96, 320, 208, number,
                       SCALE["mega"], self.t["surface"], bold=True,
                       font=self.display_font, align="right")
            self._text(s, MARGIN_X, top - 64, 200, 48, number,
                       SCALE["h3"], self.t["accent"], bold=True, font=self.display_font)
        size = _fit_size(title, text_w, 72, start="h1")
        self._text(s, MARGIN_X, top, text_w, 72, title, size,
                   self.t["text"], bold=True, font=self.display_font)
        if subtitle:
            self._text(s, MARGIN_X, top + 88, text_w, 40, subtitle,
                       SCALE["body"], self.t["muted"])
        return self

    def bullets(self, title: str, items: list[str] | None = None) -> "Deck":
        s = self._blank()
        self._title(s, title)
        if items:
            from pptx.enum.text import MSO_AUTO_SIZE
            from pptx.util import Pt

            height = SLIDE_H - BODY_TOP - 64
            box = s.shapes.add_textbox(Pt(MARGIN_X), Pt(BODY_TOP),
                                       Pt(CONTENT_W), Pt(snap(height)))
            frame = box.text_frame
            frame.word_wrap = True
            frame.auto_size = MSO_AUTO_SIZE.NONE
            longest = max(items, key=len)
            size = _fit_size(longest, CONTENT_W - 40, height / max(len(items), 1))
            self._lines(frame, items, size, self.t["text"], prefix="—   ")
        return self

    def two_column(self, title: str, left: list[str], right: list[str],
                   left_title: str | None = None,
                   right_title: str | None = None) -> "Deck":
        s = self._blank()
        self._title(s, title)
        gap = 32
        col_w = (CONTENT_W - gap) // 2
        for x, heading, items in ((MARGIN_X, left_title, left),
                                  (MARGIN_X + col_w + gap, right_title, right)):
            top = BODY_TOP
            if heading:
                self._text(s, x, top, col_w, 48, heading, SCALE["h3"],
                           self.t["accent"], bold=True, font=self.display_font)
                top += 64
            from pptx.enum.text import MSO_AUTO_SIZE
            from pptx.util import Pt
            height = SLIDE_H - top - 64
            box = s.shapes.add_textbox(Pt(snap(x)), Pt(snap(top)),
                                       Pt(snap(col_w)), Pt(snap(height)))
            frame = box.text_frame
            frame.word_wrap = True
            frame.auto_size = MSO_AUTO_SIZE.NONE
            if items:
                size = _fit_size(max(items, key=len), col_w - 32,
                                 height / max(len(items), 1))
                self._lines(frame, items, size, self.t["text"], prefix="—   ")
        return self

    def compare(self, title: str, left_title: str, left: list[str],
                right_title: str, right: list[str]) -> "Deck":
        """Two panelled columns. `two_column` is bare text; this is a contrast,
        so the panels carry the comparison."""
        s = self._blank()
        self._title(s, title)
        gap = 32
        col_w = (CONTENT_W - gap) // 2
        height = SLIDE_H - BODY_TOP - 64
        for x, heading, items, colour in (
                (MARGIN_X, left_title, left, self.t["primary"]),
                (MARGIN_X + col_w + gap, right_title, right, self.t["accent"])):
            card = self._card(s, x, BODY_TOP, col_w, height, rule=colour)
            frame = card.text_frame
            self._lines(frame, [heading], SCALE["h3"], colour, bold=True, space=16)
            if items:
                size = _fit_size(max(items, key=len), col_w - 64,
                                 (height - 80) / max(len(items), 1))
                self._lines(frame, items, size, self.t["text"], first=False,
                            prefix="—   ")
        return self

    def bento(self, title: str, cards: list) -> "Deck":
        """A grid of panels. Each card is (label, value) or (label, value, note)."""
        s = self._blank()
        self._title(s, title)
        cards = list(cards)[:6]
        cols = _balanced_columns(len(cards), most=3)
        rows = -(-len(cards) // cols)
        gap = 24
        cw = (CONTENT_W - gap * (cols - 1)) // cols
        ch = (SLIDE_H - BODY_TOP - 64 - gap * (rows - 1)) // rows
        for i, card in enumerate(cards):
            label, value, *rest = list(card) + [""]
            note = rest[0] if rest else ""
            x = MARGIN_X + (i % cols) * (cw + gap)
            y = BODY_TOP + (i // cols) * (ch + gap)
            shape = self._card(s, x, y, cw, ch, rule=self.t["line"])
            frame = shape.text_frame
            self._lines(frame, [str(label)], SCALE["caption"], self.t["muted"], space=8)
            self._lines(frame, [str(value)],
                        _card_value_size(str(value), cw - 64, ch,
                                         lines_around=2 if note else 1),
                        self.t["text"], first=False, bold=True, space=8)
            if note:
                self._lines(frame, [str(note)], SCALE["caption"],
                            self.t["accent"], first=False, space=0)
        return self

    def kpi(self, title: str, metrics: list) -> "Deck":
        """A dense metric row. Up to 8, as the spec's dashboards require."""
        s = self._blank()
        self._title(s, title)
        metrics = list(metrics)[:8]
        cols = _balanced_columns(len(metrics), most=4)
        rows = -(-len(metrics) // cols)
        gap = 24
        cw = (CONTENT_W - gap * (cols - 1)) // cols
        ch = min(176, (SLIDE_H - BODY_TOP - 64 - gap * (rows - 1)) // rows)
        for i, metric in enumerate(metrics):
            label, value, *rest = list(metric) + [""]
            delta = rest[0] if rest else ""
            x = MARGIN_X + (i % cols) * (cw + gap)
            y = BODY_TOP + (i // cols) * (ch + gap)
            shape = self._card(s, x, y, cw, ch, rule=self.t["line"])
            frame = shape.text_frame
            self._lines(frame, [str(value)],
                        _card_value_size(str(value), cw - 64, ch,
                                         lines_around=2 if delta else 1),
                        self.t["text"], bold=True, space=8)
            self._lines(frame, [str(label)], SCALE["caption"], self.t["muted"],
                        first=False, space=4)
            if delta:
                self._lines(frame, [str(delta)], SCALE["caption"],
                            self.t["accent"], first=False, space=0)
        return self

    def timeline(self, title: str, events: list) -> "Deck":
        """Horizontal milestones. Each event is (when, what)."""
        from pptx.util import Pt

        s = self._blank()
        self._title(s, title)
        events = list(events)[:6]
        if not events:
            return self

        axis_y = 296
        rule = s.shapes.add_shape(1, Pt(MARGIN_X), Pt(axis_y), Pt(CONTENT_W), Pt(GRID))
        rule.fill.solid()
        rule.fill.fore_color.rgb = _rgb(self.t["muted"])
        rule.line.fill.background()
        rule.shadow.inherit = False

        slot = CONTENT_W // len(events)
        for i, event in enumerate(events):
            when, what = (list(event) + [""])[:2]
            x = MARGIN_X + i * slot

            # A tick joining each label to the axis. Without it the dates float
            # unattached and the slide is three captions near a line rather than
            # a timeline — legible to a checker, meaningless to a reader.
            #
            # It ABUTS the rule rather than sitting on it: touching edges are
            # not an overlap, so the marker reads as attached without tripping
            # the overlap check, and no auditor rule had to be loosened to
            # accommodate a decoration.
            tick_y = axis_y - GRID if i % 2 == 0 else axis_y + GRID
            tick = s.shapes.add_shape(1, Pt(snap(x)), Pt(tick_y), Pt(GRID), Pt(GRID))
            tick.fill.solid()
            tick.fill.fore_color.rgb = _rgb(self.t["accent"])
            tick.line.fill.background()
            tick.shadow.inherit = False
            # Alternating above/below keeps six labels legible where a single
            # row would collide; the axis is the constant that ties them.
            above = i % 2 == 0
            # 48pt for a 30pt label, not 32: a box sized to the type leaves no
            # room for the frame's own insets, and PowerPoint then clips the
            # descenders. The date and its caption stack without overlapping
            # the axis in either direction.
            self._text(s, x, axis_y - 112 if above else axis_y + 24, slot - 16, 48,
                       str(when), SCALE["h3"], self.t["accent"], bold=True,
                       font=self.display_font)
            self._text(s, x, axis_y - 64 if above else axis_y + 72, slot - 16, 56,
                       str(what), SCALE["caption"], self.t["text"])
        return self

    def process(self, title: str, steps: list[str]) -> "Deck":
        """Numbered stages, left to right."""
        s = self._blank()
        self._title(s, title)
        steps = list(steps)[:6]
        if not steps:
            return self
        gap = 24
        cw = (CONTENT_W - gap * (len(steps) - 1)) // len(steps)
        for i, step in enumerate(steps):
            x = MARGIN_X + i * (cw + gap)
            shape = self._card(s, x, 216, cw, 152, rule=self.t["primary"])
            frame = shape.text_frame
            self._lines(frame, [f"{i + 1}"], SCALE["h3"], self.t["accent"],
                        bold=True, space=8)
            self._lines(frame, [str(step)],
                        _fit_size(str(step), cw - 64, 72), self.t["text"],
                        first=False, space=0)
        return self

    def table(self, title: str, rows: list[list], header: bool = True) -> "Deck":
        """A real table. The spec asks for 15+ row comparisons, so this paginates
        rather than shrinking type past the readability floor."""
        from pptx.util import Pt

        rows = [list(r) for r in rows if r]
        if not rows:
            return self
        head, body = (rows[0], rows[1:]) if header else (None, rows)

        # 11 body rows is what fits at 16pt inside the content area. Beyond that
        # the deck gains a slide; it never gains 9pt text.
        per_slide = 11
        pages = [body[i:i + per_slide] for i in range(0, len(body), per_slide)] or [[]]
        for page_no, page in enumerate(pages):
            s = self._blank()
            caption = title if len(pages) == 1 else f"{title} ({page_no + 1}/{len(pages)})"
            self._title(s, caption)

            data = ([head] if head else []) + page
            n_rows, n_cols = len(data), max(len(r) for r in data)
            height = snap(min(SLIDE_H - BODY_TOP - 64, 40 * n_rows))
            frame = s.shapes.add_table(n_rows, n_cols, Pt(MARGIN_X), Pt(BODY_TOP),
                                       Pt(CONTENT_W), Pt(height))
            tbl = frame.table
            tbl.first_row = bool(head)
            for r, row in enumerate(data):
                for c in range(n_cols):
                    cell = tbl.cell(r, c)
                    cell.text = str(row[c]) if c < len(row) else ""
                    para = cell.text_frame.paragraphs[0]
                    para.font.size = Pt(SCALE["caption"])
                    para.font.name = self.body_font
                    is_head = head is not None and r == 0
                    para.font.bold = is_head
                    para.font.color.rgb = _rgb(self.t["primary"] if is_head
                                               else self.t["text"])
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = _rgb(self.t["bg"])
        return self

    def chart(self, title: str, categories: list, series: dict,
              kind: str = "bar") -> "Deck":
        """A native, editable chart — not a picture of one."""
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
        from pptx.util import Pt

        kinds = {
            "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "hbar": XL_CHART_TYPE.BAR_CLUSTERED,
            "line": XL_CHART_TYPE.LINE_MARKERS,
            "area": XL_CHART_TYPE.AREA,
            "pie": XL_CHART_TYPE.PIE,
            "doughnut": XL_CHART_TYPE.DOUGHNUT,
            "scatter": XL_CHART_TYPE.XY_SCATTER,
        }
        s = self._blank()
        self._title(s, title)

        data = CategoryChartData()
        data.categories = list(categories)
        for name, values in series.items():
            data.add_series(str(name), tuple(values))

        height = snap(SLIDE_H - BODY_TOP - 64)
        frame = s.shapes.add_chart(kinds.get(kind, kinds["bar"]), Pt(MARGIN_X),
                                   Pt(BODY_TOP), Pt(CONTENT_W), Pt(height), data)
        chart = frame.chart
        chart.has_title = False
        if len(series) > 1:
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
            chart.legend.font.size = Pt(SCALE["caption"])
            chart.legend.font.color.rgb = _rgb(self.t["muted"])
        try:
            for axis in (chart.category_axis, chart.value_axis):
                axis.tick_labels.font.size = Pt(SCALE["caption"])
                axis.tick_labels.font.color.rgb = _rgb(self.t["muted"])
        except Exception:
            # Pie and doughnut have no axes. Styling them is optional; failing
            # the whole deck over it is not.
            pass
        return self

    def code(self, title: str, source: str, caption: str | None = None) -> "Deck":
        """Monospaced source on a panel. Never wrapped — wrapped code is wrong."""
        from pptx.util import Pt

        s = self._blank()
        self._title(s, title)
        height = snap(SLIDE_H - BODY_TOP - (96 if caption else 64))
        panel = self._card(s, MARGIN_X, BODY_TOP, CONTENT_W, height,
                           rule=self.t["primary"])
        frame = panel.text_frame
        frame.word_wrap = False
        lines = source.replace("\t", "    ").split("\n")
        size = SCALE["caption"] if len(lines) > 14 else SCALE["body"]
        for i, line in enumerate(lines):
            p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
            p.text = line or " "
            p.font.size = Pt(size)
            p.font.name = self.mono_font
            p.font.color.rgb = _rgb(self.t["text"])
            p.space_after = Pt(0)
        if caption:
            self._text(s, MARGIN_X, SLIDE_H - 96, CONTENT_W, 32, caption,
                       SCALE["caption"], self.t["muted"])
        return self

    def quote(self, text: str, attribution: str | None = None) -> "Deck":
        s = self._blank()
        size = _fit_size(text, CONTENT_W - 64, 160, start="h2")
        self._text(s, MARGIN_X, 184, CONTENT_W, 160, f"“{text}”",
                   size, self.t["text"], font=self.display_font)
        if attribution:
            self._text(s, MARGIN_X, 376, CONTENT_W, 40, f"— {attribution}",
                       SCALE["body"], self.t["accent"])
        return self

    def matrix(self, title: str, quadrants: list) -> "Deck":
        """A 2x2. Each quadrant is (heading, [items]) — SWOT, risk, effort/impact."""
        s = self._blank()
        self._title(s, title)
        quadrants = list(quadrants)[:4]
        gap = 24
        cw = (CONTENT_W - gap) // 2
        ch = (SLIDE_H - BODY_TOP - 64 - gap) // 2
        accents = [self.t["primary"], self.t["accent"],
                   self.t["muted"], self.t["text"]]
        for i, quadrant in enumerate(quadrants):
            heading, items = (list(quadrant) + [[]])[:2]
            x = MARGIN_X + (i % 2) * (cw + gap)
            y = BODY_TOP + (i // 2) * (ch + gap)
            card = self._card(s, x, y, cw, ch, rule=accents[i % 4])
            frame = card.text_frame
            self._lines(frame, [str(heading)], SCALE["h3"], accents[i % 4],
                        bold=True, space=12)
            if items:
                self._lines(frame, [str(v) for v in items], SCALE["caption"],
                            self.t["text"], first=False, space=6, prefix="—   ")
        return self

    def appendix(self, title: str, rows: list) -> "Deck":
        """Dense reference material: two columns of label/value pairs."""
        s = self._blank()
        self._title(s, title)
        rows = [tuple(r) for r in rows][:16]
        gap = 32
        col_w = (CONTENT_W - gap) // 2
        half = -(-len(rows) // 2)
        for col, chunk in enumerate((rows[:half], rows[half:])):
            if not chunk:
                continue
            from pptx.enum.text import MSO_AUTO_SIZE
            from pptx.util import Pt
            x = MARGIN_X + col * (col_w + gap)
            height = snap(SLIDE_H - BODY_TOP - 64)
            box = s.shapes.add_textbox(Pt(snap(x)), Pt(BODY_TOP),
                                       Pt(snap(col_w)), Pt(height))
            frame = box.text_frame
            frame.word_wrap = True
            frame.auto_size = MSO_AUTO_SIZE.NONE
            self._lines(frame, [f"{k}   {v}" for k, v in chunk],
                        SCALE["caption"], self.t["text"], space=8)
        return self

    # ── compatibility ───────────────────────────────────────────────────────
    # The previous API. Kept because the themed-documents skill and existing
    # generated scripts call them by these names; renaming would break every
    # deck the model already knows how to write.
    def cover(self, title: str, subtitle: str | None = None) -> "Deck":
        return self.hero(title, subtitle)

    def slide(self, title: str, bullets: list[str] | None = None) -> "Deck":
        return self.bullets(title, bullets)

    def save(self, filename: str | None = None) -> str:
        target = filename or self.filename
        self.prs.save(target)
        return target


# ── PDF ──────────────────────────────────────────────────────────────────────
class Report:
    """A themed PDF.

        r = Report("deep_sea.pdf", theme="paper", title="Bioluminescence")
        r.heading("Why it glows"); r.text("…"); r.bullets(["a", "b"]); r.save()
    """

    def __init__(self, filename: str, theme: str | None = None,
                 title: str | None = None, subtitle: str | None = None) -> None:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.units import mm
        from reportlab.platypus import SimpleDocTemplate

        self.filename = filename
        self.t = globals()["theme"](theme)
        self.story: list = []
        self._doc = SimpleDocTemplate(
            filename, pagesize=A4,
            leftMargin=22 * mm, rightMargin=22 * mm,
            topMargin=22 * mm, bottomMargin=20 * mm,
            title=title or filename,
        )
        if title:
            self.title(title, subtitle)

    def _colour(self, hex6: str):
        from reportlab.lib.colors import HexColor

        return HexColor(f"#{hex6}")

    def _style(self, name, size, colour, leading=None, bold=False, space=6):
        from reportlab.lib.styles import ParagraphStyle

        return ParagraphStyle(
            name, fontName="Helvetica-Bold" if bold else "Helvetica",
            fontSize=size, leading=leading or size * 1.45,
            textColor=self._colour(colour), spaceAfter=space,
        )

    def title(self, text: str, subtitle: str | None = None) -> "Report":
        from reportlab.platypus import Paragraph, Spacer

        self.story.append(Paragraph(text, self._style("t", 26, self.t["text"], bold=True, space=4)))
        if subtitle:
            self.story.append(Paragraph(subtitle, self._style("st", 12, self.t["muted"])))
        self.story.append(Spacer(1, 14))
        return self

    def heading(self, text: str) -> "Report":
        from reportlab.platypus import Paragraph, Spacer

        self.story.append(Spacer(1, 10))
        self.story.append(Paragraph(text, self._style("h", 15, self.t["primary"], bold=True, space=4)))
        return self

    def text(self, body: str) -> "Report":
        from reportlab.platypus import Paragraph

        self.story.append(Paragraph(body, self._style("p", 10.5, self.t["text"], space=8)))
        return self

    def bullets(self, items: list[str]) -> "Report":
        from reportlab.platypus import Paragraph

        style = self._style("b", 10.5, self.t["text"], space=4)
        for item in items:
            self.story.append(Paragraph(f"&mdash;&nbsp;&nbsp;{item}", style))
        return self

    def table(self, rows: list[list], header: bool = True) -> "Report":
        from reportlab.platypus import Table, TableStyle

        if not rows:
            return self
        t = Table([[str(c) for c in r] for r in rows], hAlign="LEFT")
        commands = [
            ("FONTNAME", (0, 0), (-1, -1), "Helvetica"),
            ("FONTSIZE", (0, 0), (-1, -1), 9),
            ("TEXTCOLOR", (0, 0), (-1, -1), self._colour(self.t["text"])),
            ("LINEBELOW", (0, 0), (-1, -1), 0.25, self._colour(self.t["muted"])),
            ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
            ("LEFTPADDING", (0, 0), (-1, -1), 4),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
        ]
        if header:
            commands += [("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                         ("TEXTCOLOR", (0, 0), (-1, 0), self._colour(self.t["primary"]))]
        t.setStyle(TableStyle(commands))
        self.story.append(t)
        return self

    def _paint(self, canvas, doc):
        """The page's own background, drawn before anything else.

        A dark theme needs it — reportlab's page is white, and light text on
        white is an invisible document. Drawn on every page rather than once,
        because a second page would otherwise arrive blank-white."""
        from reportlab.lib.units import mm

        canvas.saveState()
        canvas.setFillColor(self._colour(self.t["bg"]))
        canvas.rect(0, 0, doc.pagesize[0], doc.pagesize[1], stroke=0, fill=1)
        canvas.setFillColor(self._colour(self.t["primary"]))
        canvas.rect(0, 0, 3 * mm, doc.pagesize[1], stroke=0, fill=1)
        canvas.restoreState()

    def save(self, filename: str | None = None) -> str:
        if filename:
            self._doc.filename = filename
        self._doc.build(self.story, onFirstPage=self._paint, onLaterPages=self._paint)
        return filename or self.filename


# ── Word ─────────────────────────────────────────────────────────────────────
class Doc:
    """A themed Word document.

        d = Doc("brief.docx", theme="sand", title="Urban Heat Islands")
        d.heading("Causes"); d.text("…"); d.save()
    """

    def __init__(self, filename: str, theme: str | None = None,
                 title: str | None = None) -> None:
        import docx

        self.filename = filename
        self.t = globals()["theme"](theme)
        self.document = docx.Document()
        self._paint_page()
        if title:
            self.title(title)

    def _paint_page(self) -> None:
        """Word ignores a background colour unless it is told to display one,
        which is a document setting rather than a property of the colour."""
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement

        root = self.document.element
        background = OxmlElement("w:background")
        background.set(qn("w:color"), self.t["bg"])
        root.insert(0, background)

        settings = self.document.settings.element
        display = OxmlElement("w:displayBackgroundShape")
        settings.append(display)

    def _colour(self):
        from docx.shared import RGBColor

        return RGBColor.from_string

    def _write(self, text: str, size: int, colour: str, bold=False, font=BODY_FONT):
        from docx.shared import Pt, RGBColor

        p = self.document.add_paragraph()
        run = p.add_run(text)
        run.bold = bold
        run.font.size = Pt(size)
        run.font.name = font
        run.font.color.rgb = RGBColor.from_string(colour)
        return p

    def title(self, text: str) -> "Doc":
        self._write(text, 24, self.t["text"], bold=True, font=DISPLAY_FONT)
        return self

    def heading(self, text: str) -> "Doc":
        self._write(text, 14, self.t["primary"], bold=True, font=DISPLAY_FONT)
        return self

    def text(self, body: str) -> "Doc":
        self._write(body, 11, self.t["text"])
        return self

    def bullets(self, items: list[str]) -> "Doc":
        for item in items:
            self._write(f"—  {item}", 11, self.t["text"])
        return self

    def save(self, filename: str | None = None) -> str:
        target = filename or self.filename
        self.document.save(target)
        return target


# ── Charts ───────────────────────────────────────────────────────────────────
def chart_style(name: str | None = None) -> dict:
    """rcParams for matplotlib, so a chart matches the document it goes in.

        import matplotlib; matplotlib.rcParams.update(chart_style("midnight"))
    """
    t = theme(name)
    fg, bg = f"#{t['text']}", f"#{t['bg']}"
    return {
        "figure.facecolor": bg, "axes.facecolor": bg, "savefig.facecolor": bg,
        "text.color": fg, "axes.labelcolor": fg, "axes.edgecolor": f"#{t['muted']}",
        "xtick.color": f"#{t['muted']}", "ytick.color": f"#{t['muted']}",
        "grid.color": f"#{t['muted']}", "grid.alpha": 0.25,
        "axes.grid": True, "axes.spines.top": False, "axes.spines.right": False,
        "axes.titlecolor": fg, "font.size": 10,
        "axes.prop_cycle": __import__("cycler").cycler(
            color=[f"#{t['primary']}", f"#{t['accent']}", f"#{t['muted']}"]),
    }
