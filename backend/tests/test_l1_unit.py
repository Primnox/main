"""L1 — Unit tests. Fast, deterministic, no UI.

Each class here corresponds to one component the spec names: the Context
Service, the event sequencer, workspace versioning, and sandbox permissions.
"""
from __future__ import annotations

import threading

import pytest
from conftest import wait_until

from primnox2.assets import service as assets
from primnox2.chat import turns
from primnox2.context import service as context
from primnox2.kernel.events import bus
from primnox2.sandbox import permissions as sandbox_perms
from primnox2.sandbox import snapshots
from primnox2.storage import db
from primnox2.workspaces import service as workspaces


def _completed_turn(conversation_id: str, user_text: str, reply: str) -> str:
    turn = turns.create_turn(conversation_id, user_text)
    turns.complete(turn["turn_id"], reply)
    return turn["turn_id"]


# ── Context Service ──────────────────────────────────────────────────────────
class TestContextService:
    def test_token_budget_respected(self, conversation):
        for i in range(40):
            _completed_turn(conversation, f"question {i} " + "x" * 400,
                            f"answer {i} " + "y" * 400)
        bundle = context.build(conversation, "the newest question", budget=1500)
        assert bundle.tokens <= bundle.budget, "context exceeded its own budget"
        assert bundle.dropped_turns > 0, "nothing was dropped despite a tight budget"

    def test_ordering_preserved(self, conversation):
        for i in range(5):
            _completed_turn(conversation, f"q{i}", f"a{i}")
        bundle = context.build(conversation, "latest", budget=100_000)

        body = [m for m in bundle.messages if m["role"] != "system"]
        assert body[-1]["content"] == "latest", "the current prompt must come last"

        # History must reach the model forwards, even though selection walked
        # backwards to fit the budget.
        indices = [int(m["content"][1:]) for m in body if m["content"][:1] in ("q", "a")
                   and m["content"][1:].isdigit()]
        assert indices == sorted(indices), f"history reached the model out of order: {indices}"

    def test_current_prompt_never_dropped(self, conversation):
        for i in range(30):
            _completed_turn(conversation, "x" * 500, "y" * 500)
        bundle = context.build(conversation, "THE ACTUAL QUESTION", budget=400)
        assert bundle.messages[-1]["content"].startswith("THE ACTUAL QUESTION")

    def test_oversized_prompt_truncated_not_rejected(self, conversation):
        bundle = context.build(conversation, "q" * 100_000, budget=500)
        assert bundle.tokens <= bundle.budget
        assert any("truncat" in n for n in bundle.notes)

    def test_asset_references_intact(self, conversation):
        turn = turns.create_turn(conversation, "summarise it")
        tid = turn["turn_id"]
        asset = assets.ingest_bytes(b"The capital of France is Paris.", "geo.txt",
                                    conversation_id=conversation, turn_id=tid)
        wait_until(lambda: assets.get(asset["id"])["status"] == "ready",
                   what="asset ingestion")

        bundle = context.build(conversation, "what does it say?", turn_id=tid, budget=50_000)
        assert asset["id"] in bundle.asset_ids
        joined = "\n".join(m["content"] for m in bundle.messages)
        assert "capital of France" in joined, "the asset's text never reached the model"
        assert "geo.txt" in joined, "the asset was not named for citation"

    def test_abandoned_turns_do_not_pollute_history(self, conversation):
        """A turn stopped before its first token must not leave its question
        in the context.

        Found live: two stopped turns produced three consecutive `user`
        messages in the prompt, and the model answered the FIRST one — so a
        user who pressed stop and then asked something else got a reply to the
        question they had abandoned.
        """
        _completed_turn(conversation, "a real question", "a real answer")

        # Two turns stopped before producing anything.
        for text in ("abandoned question one", "abandoned question two"):
            abandoned = turns.create_turn(conversation, text)
            turns.finish_cancelled(abandoned["turn_id"], "")

        bundle = context.build(conversation, "the question I actually asked", budget=50_000)
        contents = [m["content"] for m in bundle.messages]
        assert "abandoned question one" not in contents
        assert "abandoned question two" not in contents

        # And no two user messages ever sit next to each other.
        roles = [m["role"] for m in bundle.messages if m["role"] != "system"]
        assert not any(a == "user" and b == "user" for a, b in zip(roles, roles[1:])), \
            f"consecutive user messages in the prompt: {roles}"

    def test_cancelled_turn_with_partial_text_is_kept(self, conversation):
        """Partial output is a real exchange — it must survive into history."""
        stopped = turns.create_turn(conversation, "half answered")
        turns.finish_cancelled(stopped["turn_id"], "here is the first half")

        bundle = context.build(conversation, "carry on", budget=50_000)
        contents = "\n".join(m["content"] for m in bundle.messages)
        assert "half answered" in contents
        assert "here is the first half" in contents

    def test_unreadable_asset_is_named_not_silently_empty(self, conversation):
        """V1 sent an empty string for a scanned PDF. The model must instead be
        told the document exists and could not be read."""
        turn = turns.create_turn(conversation, "read it")
        tid = turn["turn_id"]
        asset = assets.ingest_bytes(b"\x89PNG\r\n\x1a\n" + b"\x00" * 40, "scan.png",
                                    conversation_id=conversation, turn_id=tid)
        wait_until(lambda: assets.get(asset["id"])["status"] == "ready",
                   what="image ingestion")

        bundle = context.build(conversation, "what is in it?", turn_id=tid, budget=50_000)
        joined = "\n".join(m["content"] for m in bundle.messages)
        assert "scan.png" in joined
        assert "OCR" in joined or "no text" in joined


# ── Event sequencer ──────────────────────────────────────────────────────────
class TestEventSequencer:
    def test_monotonic_under_concurrency(self, conversation):
        """Generate 1,2,3 — never 1,3,2. Sequence monotonicity is sacred."""
        seen: list[int] = []
        lock = threading.Lock()

        def emit_many():
            for _ in range(25):
                e = bus.emit("token", {"text": "x"}, conversation_id=conversation)
                with lock:
                    seen.append(e["sequence"])

        threads = [threading.Thread(target=emit_many) for _ in range(4)]
        for t in threads:
            t.start()
        for t in threads:
            t.join()

        assert len(seen) == len(set(seen)), "two events were given the same sequence"
        ordered = sorted(seen)
        assert ordered == list(range(ordered[0], ordered[0] + len(ordered))), \
            "concurrent emission left a gap"

    def test_ambient_consumes_no_sequence(self):
        """CRS §11.1.4 — ambient events are ephemeral and must not advance the
        durable counter, or the island's chatter would inflate every cursor."""
        before = bus.head()
        for _ in range(5):
            bus.emit("now_playing", {"track": "x"}, scope="ambient")
        assert bus.head() == before

    def test_ambient_never_persisted(self):
        """CRS §11.1.2."""
        bus.emit("flow_state", {"state": "deep"}, scope="ambient")
        rows = db.connect().execute(
            "SELECT COUNT(*) AS n FROM events WHERE kind='flow_state'").fetchone()
        assert rows["n"] == 0


# ── Workspaces ───────────────────────────────────────────────────────────────
class TestWorkspaceVersioning:
    def test_create_modify_version_diff(self, conversation):
        ws = workspaces.create("python", "calc", {"main.py": "v1", "util.py": "helper"},
                               conversation_id=conversation)
        wid = ws["workspace_id"]

        workspaces.update(wid, {"main.py": "v2"}, conversation_id=conversation)
        assert workspaces.read_files(wid)["main.py"] == "v2"
        assert workspaces.read_files(wid)["util.py"] == "helper", "untouched file was lost"

        # Every prior version stays readable — versions are immutable.
        assert workspaces.read_files(wid, 1)["main.py"] == "v1"

        d = workspaces.diff(wid, 1, 2)
        assert d["modified"] == ["main.py"] and not d["created"] and not d["deleted"]

    def test_history_is_never_destroyed(self, conversation):
        ws = workspaces.create("markdown", "notes", {"a.md": "one"},
                               conversation_id=conversation)
        wid = ws["workspace_id"]
        workspaces.update(wid, {"a.md": "two"}, conversation_id=conversation)
        workspaces.revert(wid, 1, conversation_id=conversation)

        versions = [v["version"] for v in workspaces.versions(wid)]
        assert versions == [1, 2, 3], "revert rewrote history instead of appending to it"
        assert workspaces.read_files(wid)["a.md"] == "one"
        assert workspaces.read_files(wid, 2)["a.md"] == "two", "the reverted version was destroyed"

    def test_noop_edit_creates_no_version(self, conversation):
        ws = workspaces.create("doc", "same", {"a.txt": "identical"},
                               conversation_id=conversation)
        wid = ws["workspace_id"]
        result = workspaces.update(wid, {"a.txt": "identical"})
        assert result["unchanged"] and result["version"] == 1

    def test_survives_origin_turn_deletion(self, conversation):
        """CRS §2.5 — deleting the chat must not delete the artifact."""
        turn = turns.create_turn(conversation, "build it")
        ws = workspaces.create("html", "page", {"index.html": "<p>hi</p>"},
                               origin_turn_id=turn["turn_id"], conversation_id=conversation)
        with db.tx() as c:
            c.execute("DELETE FROM turns WHERE id=?", (turn["turn_id"],))

        surviving = workspaces.get(ws["workspace_id"])
        assert surviving is not None, "the workspace died with its turn"
        assert surviving["files"]["index.html"] == "<p>hi</p>"
        assert surviving["origin_turn_id"] is None

    def test_current_version_points_at_a_real_version(self, conversation):
        ws = workspaces.create("python", "v", {"a.py": "1"}, conversation_id=conversation)
        workspaces.update(ws["workspace_id"], {"a.py": "2"})
        row = db.connect().execute("SELECT current_version FROM workspaces WHERE id=?",
                                   (ws["workspace_id"],)).fetchone()
        exists = db.connect().execute(
            "SELECT 1 FROM workspace_versions WHERE workspace_id=? AND version=?",
            (ws["workspace_id"], row["current_version"])).fetchone()
        assert exists is not None


# ── Sandbox permissions ──────────────────────────────────────────────────────
class TestSandboxPermissions:
    """Every permission in the spec's table becomes a test."""

    def test_workspace_write_allowed(self):
        m = sandbox_perms.manifest_for("python", sandbox_perms.SAFE)
        assert m.filesystem["workspace"] == sandbox_perms.ALLOW
        assert sandbox_perms.validate(m) == []

    def test_documents_prompts_rather_than_allows(self):
        m = sandbox_perms.manifest_for("shell", sandbox_perms.LIMITED)
        assert m.filesystem["documents"] == sandbox_perms.ASK

    def test_registry_denied_by_default(self):
        for tier in (sandbox_perms.SAFE, sandbox_perms.LIMITED):
            m = sandbox_perms.manifest_for("python", tier)
            assert m.filesystem["registry"] == sandbox_perms.DENY

    def test_system_denied_by_default(self):
        for tier in (sandbox_perms.SAFE, sandbox_perms.LIMITED):
            m = sandbox_perms.manifest_for("python", tier)
            assert m.filesystem["system"] == sandbox_perms.DENY

    def test_privilege_cannot_be_laundered_through_a_safe_manifest(self):
        """The tier is a ceiling. Granting system access on a safe manifest
        must be rejected, not quietly honoured."""
        for scope in ("system", "registry"):
            m = sandbox_perms.manifest_for("python", sandbox_perms.SAFE,
                                           filesystem={scope: sandbox_perms.ALLOW})
            assert sandbox_perms.validate(m), f"{scope} was grantable at the safe tier"

    def test_network_requires_at_least_limited(self):
        m = sandbox_perms.manifest_for("python", sandbox_perms.SAFE,
                                       network=sandbox_perms.BALANCED)
        assert sandbox_perms.validate(m)

    def test_workspace_cannot_be_denied(self):
        m = sandbox_perms.manifest_for("python", sandbox_perms.SAFE,
                                       filesystem={"workspace": sandbox_perms.DENY})
        assert sandbox_perms.validate(m)

    def test_resource_limits_bounded(self):
        assert sandbox_perms.validate(sandbox_perms.manifest_for("python", timeout_s=99_999))
        assert sandbox_perms.validate(sandbox_perms.manifest_for("python", memory_mb=1))

    def test_safe_tier_approval_is_reusable_shell_is_not(self):
        safe = sandbox_perms.manifest_for("python", sandbox_perms.SAFE)
        limited = sandbox_perms.manifest_for("shell", sandbox_perms.LIMITED)
        assert sandbox_perms.is_reusable(safe)
        assert not sandbox_perms.is_reusable(limited), \
            "a shell approval must be asked for every time"


class TestRetry:
    """CRS §5.2.3 — a retry creates a new turn referencing the failed one.

    The button existed in the UI with no handler behind it for a while, which
    is worse than no button: it tells the user an action is available and then
    silently does nothing.
    """

    def test_retry_creates_a_new_turn_referencing_the_failed_one(self, conversation):
        original = turns.create_turn(conversation, "this one will fail")
        turns.fail(original["turn_id"], "provider_unreachable", "nope", retryable=True)

        retried = turns.retry_turn(original["turn_id"])
        assert retried["turn_id"] != original["turn_id"], "retry reopened the old turn"

        row = db.connect().execute(
            "SELECT retry_of_turn_id, status FROM turns WHERE id=?",
            (retried["turn_id"],)).fetchone()
        assert row["retry_of_turn_id"] == original["turn_id"], "the link back was not recorded"
        assert row["status"] == "queued"

    def test_failed_turn_stays_failed(self, conversation):
        """The original is evidence of what happened and must not be rewritten."""
        original = turns.create_turn(conversation, "keep the record")
        turns.fail(original["turn_id"], "provider_auth", "bad key", retryable=False)
        turns.retry_turn(original["turn_id"])

        row = db.connect().execute("SELECT status, error_code FROM turns WHERE id=?",
                                   (original["turn_id"],)).fetchone()
        assert row["status"] == "failed" and row["error_code"] == "provider_auth"

    def test_retry_carries_the_original_text(self, conversation):
        original = turns.create_turn(conversation, "the exact original wording")
        turns.fail(original["turn_id"], "internal", "boom", retryable=True)
        retried = turns.retry_turn(original["turn_id"])

        text = db.connect().execute(
            "SELECT text FROM messages WHERE turn_id=? AND role='user'",
            (retried["turn_id"],)).fetchone()["text"]
        assert text == "the exact original wording"


class TestUsageAccounting:
    """CRS §3.6 — `turn.completed` carries `usage`.

    It went unnoticed that this was always `{}` until someone asked what a turn
    actually costs. An always-empty required field is a contract that exists
    only on paper.
    """

    def test_completed_turn_records_usage(self, conversation, events, scripted):
        from conftest import run_turn, wait_for_turn

        scripted("A measured reply.")
        tid = run_turn(conversation, "how much did this cost?")
        assert wait_for_turn(tid) == "completed"

        completed = events.of_kind("turn.completed", tid)
        assert completed, "no turn.completed event"
        usage = completed[0]["payload"]["usage"]
        assert usage.get("input_tokens", 0) > 0, "input tokens were not recorded"
        assert usage.get("output_tokens", 0) > 0, "output tokens were not recorded"
        assert usage.get("model_calls") == 1

    def test_tool_turn_sums_every_model_call(self, conversation, events, scripted):
        """A tool turn makes several provider calls. Reporting only the last
        would understate the turn's real cost."""
        from conftest import run_turn, wait_for_turn

        scripted(
            '<tool name="search_assets">\n{"query": "anything"}\n</tool>',
            "Here is what I found.",
        )
        tid = run_turn(conversation, "search for something")
        assert wait_for_turn(tid) == "completed"

        usage = events.of_kind("turn.completed", tid)[0]["payload"]["usage"]
        assert usage["model_calls"] == 2, "usage did not cover every model call"
        assert usage["tool_calls"] == 1


class TestErrorClassification:
    """CRS §10.2 — `retryable` must be truthful.

    A retry button on an error that can never succeed is worse than no button:
    it converts a clear failure into a loop the user runs until they give up.
    """

    @pytest.mark.parametrize("message,code,retryable", [
        ("HTTP Error 401: Unauthorized", "provider_auth", False),
        ("HTTP Error 403: Forbidden", "provider_auth", False),
        ("HTTP Error 429: Too Many Requests", "provider_rate_limited", True),
        ("HTTP Error 402: Payment Required", "provider_quota", False),
        ("HTTP Error 404: model_not_found", "model_unavailable", False),
        ("HTTP Error 502: Bad Gateway", "provider_unreachable", True),
        ("HTTP Error 503: Service Unavailable", "provider_unreachable", True),
        ("Expecting value: line 1 column 1 (char 0)", "provider_unreachable", True),
        ("<urlopen error [Errno 11001] getaddrinfo failed>", "provider_unreachable", True),
        ("maximum context length is 8192 tokens", "context_overflow", False),
    ])
    def test_classification(self, message, code, retryable):
        from primnox2.kernel.scheduler import _classify

        got_code, got_message, got_retryable = _classify(RuntimeError(message))
        assert got_code == code, f"{message!r} classified as {got_code}"
        assert got_retryable is retryable, f"{message!r} retryable={got_retryable}"
        assert got_message and not got_message.startswith(("Traceback", "<")), \
            "a raw exception reached the user-facing message"

    def test_unknown_errors_are_not_silently_swallowed(self):
        from primnox2.kernel.scheduler import _classify

        code, message, _ = _classify(RuntimeError("something nobody predicted"))
        assert code == "internal"
        assert message, "an unclassified failure produced no message at all"


class TestSnapshots:
    def test_detects_created_modified_deleted(self, tmp_path):
        (tmp_path / "keep.txt").write_text("same")
        (tmp_path / "change.txt").write_text("before")
        (tmp_path / "gone.txt").write_text("bye")
        before = snapshots.snapshot(tmp_path)

        (tmp_path / "change.txt").write_text("after")
        (tmp_path / "gone.txt").unlink()
        (tmp_path / "new.txt").write_text("hello")
        d = snapshots.diff(before, snapshots.snapshot(tmp_path))

        assert d["created"] == ["new.txt"]
        assert d["modified"] == ["change.txt"]
        assert d["deleted"] == ["gone.txt"]

    def test_rewriting_identical_bytes_is_not_a_change(self, tmp_path):
        """Content hashing, not mtime: a fast execution can rewrite a file
        without its timestamp moving."""
        f = tmp_path / "a.txt"
        f.write_text("identical")
        before = snapshots.snapshot(tmp_path)
        f.write_text("identical")
        assert snapshots.is_empty(snapshots.diff(before, snapshots.snapshot(tmp_path)))

    def test_plumbing_never_appears_as_a_result(self, tmp_path):
        before = snapshots.snapshot(tmp_path)
        for name in ("main.py", "__ac_wrapper.py", "__ac_stdout.txt", "package.json"):
            (tmp_path / name).write_text("plumbing")
        assert snapshots.is_empty(snapshots.diff(before, snapshots.snapshot(tmp_path)))


# ── Permission broker ────────────────────────────────────────────────────────
def _status(turn_id: str) -> str | None:
    row = db.connect().execute("SELECT status FROM turns WHERE id=?", (turn_id,)).fetchone()
    return row["status"] if row else None



class TestPermissionBroker:
    """A parked turn's question must survive the client that asked it.

    Found by reloading the page while a prompt was on screen: the turn came
    back as `awaiting_input` with nothing to answer, and stayed there until the
    request timed out ten minutes later.
    """

    def _ask(self, conversation_id: str, turn_id: str, monkeypatch,
             request_id: str = "req_test_1", park: bool = True):
        from primnox2.tools import permissions as broker_mod

        monkeypatch.setattr(broker_mod, "AUTO_APPROVE", "off")
        if park:
            # A turn only reaches a tool call by running, and `awaiting_input`
            # is only legal from there (§5.2).
            turns.set_status(turn_id, "building_context")
            turns.set_status(turn_id, "tool_running")
        result: dict = {}

        def ask():
            result["choice"] = broker_mod.broker.request(
                request_id=request_id, action="run_python",
                detail="Execute Python code in an isolated sandbox.",
                turn_id=turn_id, conversation_id=conversation_id, reusable=True,
                timeout_s=10,
            )

        thread = threading.Thread(target=ask, daemon=True)
        thread.start()
        # Registered first, parked second — the question is answerable before
        # it is advertised, never the other way round.
        wait_until(lambda: request_id in broker_mod.broker.pending_ids(),
                   what="the question to be registered")
        wait_until(lambda: _status(turn_id) == "awaiting_input",
                   what="the turn to park on the question")
        return broker_mod, thread, result

    def test_pending_question_is_readable_after_the_asker_is_gone(
            self, conversation, monkeypatch):
        turn_id = turns.create_turn(conversation, "run something")["turn_id"]
        broker_mod, thread, result = self._ask(conversation, turn_id, monkeypatch)

        question = broker_mod.broker.pending_for_turn(turn_id)
        assert question is not None, "a parked turn exposed no question to answer"
        assert question["id"] == "req_test_1"
        assert question["action"] == "run_python"
        assert [o["id"] for o in question["options"]] == \
            ["allow_once", "allow_turn", "deny"], "the answer buttons were lost"

        broker_mod.broker.resolve("req_test_1", "allow_once")
        thread.join(timeout=5)
        assert result["choice"] == "allow_once"
        assert broker_mod.broker.pending_for_turn(turn_id) is None, \
            "an answered question is still being offered"

    def test_history_carries_the_question_a_turn_is_parked_on(
            self, conversation, monkeypatch):
        """Opening a conversation is a state read (§3.3.3), so the state read
        is the only place a reconnecting client can learn what is being asked."""
        import asyncio

        from primnox2 import app as app_module

        turn_id = turns.create_turn(conversation, "run something")["turn_id"]
        broker_mod, thread, _ = self._ask(conversation, turn_id, monkeypatch)

        rows = asyncio.run(app_module.history(conversation))["turns"]
        row = next(r for r in rows if r["turn_id"] == turn_id)
        assert row["status"] == "awaiting_input"
        assert row.get("permission"), \
            "history rebuilt a parked turn with no question attached"
        assert row["permission"]["id"] == "req_test_1"

        broker_mod.broker.resolve("req_test_1", "deny")
        thread.join(timeout=5)

        rows = asyncio.run(app_module.history(conversation))["turns"]
        row = next(r for r in rows if r["turn_id"] == turn_id)
        assert not row.get("permission"), \
            "an answered question is still shown as pending"

    def test_answering_returns_the_turn_to_the_state_it_left(
            self, conversation, monkeypatch):
        """`awaiting_input` is a detour, not a destination — the lifecycle has
        to continue where it paused rather than jump forward."""
        turn_id = turns.create_turn(conversation, "run something")["turn_id"]
        broker_mod, thread, _ = self._ask(conversation, turn_id, monkeypatch)

        broker_mod.broker.resolve("req_test_1", "allow_once")
        thread.join(timeout=5)
        assert _status(turn_id) == "tool_running", \
            "answering a question left the turn somewhere it never was"

    def test_an_allowance_for_the_turn_is_not_asked_twice(
            self, conversation, monkeypatch):
        """'Allow for this turn' has to mean it. Asking again for the same
        action in the same turn would make the option a lie."""
        turn_id = turns.create_turn(conversation, "run something")["turn_id"]
        broker_mod, thread, _ = self._ask(conversation, turn_id, monkeypatch)

        broker_mod.broker.resolve("req_test_1", "allow_turn")
        thread.join(timeout=5)

        # The second call must not block, so a timeout here is the failure.
        second: dict = {}
        again = threading.Thread(target=lambda: second.update(
            choice=broker_mod.broker.request(
                request_id="req_test_2", action="run_python", detail="again",
                turn_id=turn_id, conversation_id=conversation, reusable=True,
                timeout_s=10)), daemon=True)
        again.start()
        again.join(timeout=3)
        assert not again.is_alive(), "a granted turn-wide allowance asked again"
        assert second["choice"] == "allow_turn"

        broker_mod.broker.forget_turn(turn_id)

    def test_a_denial_is_not_remembered_as_an_allowance(
            self, conversation, monkeypatch):
        turn_id = turns.create_turn(conversation, "run something")["turn_id"]
        broker_mod, thread, _ = self._ask(conversation, turn_id, monkeypatch)

        broker_mod.broker.resolve("req_test_1", "deny")
        thread.join(timeout=5)

        _, second_thread, second = self._ask(conversation, turn_id, monkeypatch,
                                             request_id="req_test_3", park=False)
        assert second_thread.is_alive(), "a refusal was reused as consent"
        broker_mod.broker.resolve("req_test_3", "deny")
        second_thread.join(timeout=5)
        assert second["choice"] == "deny"

    def test_cancelling_a_turn_answers_the_question_it_is_parked_on(
            self, conversation, monkeypatch):
        """CRS §9.2. Without this the worker sits on the question for the full
        timeout and the cancelled turn never terminates."""
        turn_id = turns.create_turn(conversation, "run something")["turn_id"]
        broker_mod, thread, result = self._ask(conversation, turn_id, monkeypatch)

        assert broker_mod.broker.cancel_for_turn(turn_id) == 1
        thread.join(timeout=5)
        assert not thread.is_alive(), "cancelling left the worker blocked on a prompt"
        assert result["choice"] == "deny", "a cancelled turn's question was granted"
        assert broker_mod.broker.pending_for_turn(turn_id) is None

    def test_every_silent_grant_is_still_recorded(self, conversation, events,
                                                  monkeypatch):
        """Auto-approval is defensible only because it leaves a record.

        The turn-wide reuse path returned without emitting anything, so a turn
        that ran the same tool three times produced one `permission.request`
        for three grants — the log described the decision, not the runs.
        """
        from primnox2.tools import permissions as broker_mod

        monkeypatch.setattr(broker_mod, "AUTO_APPROVE", "all")
        turn_id = turns.create_turn(conversation, "run three things")["turn_id"]

        for _ in range(3):
            choice = broker_mod.broker.request(
                request_id=f"req_{_}", action="run_python", detail="python",
                turn_id=turn_id, conversation_id=conversation, reusable=True,
            )
            assert choice == "allow_auto"

        asked = events.of_kind("permission.request", turn_id)
        answered = events.of_kind("permission.resolved", turn_id)
        assert len(asked) == 3, f"three grants left {len(asked)} records"
        assert len(answered) == 3
        assert len({e["payload"]["job_id"] for e in asked}) == 3, \
            "two grants shared an id, so one cannot be told from the other"

        broker_mod.broker.forget_turn(turn_id)


# ── Emulated tool protocol ───────────────────────────────────────────────────
class TestToolCallParsing:
    """Every string here was emitted by qwen2.5:7b against the real prompt.

    The canonical grammar alone scored 0/5 on this model — and all five
    failures named the right tool and carried valid JSON. They were refused
    over their punctuation.
    """

    from primnox2.tools import runtime as _rt

    @pytest.mark.parametrize("text", [
        '<tool name="run_python">{"code": "print(137 * 449)"}</tool>',
        'run_python({"code": "print(137 * 449)"})',
        '<run_python> {"code": "print(137 * 449)"} </run_python>',
    ])
    def test_accepts_the_shapes_models_actually_emit(self, text):
        from primnox2.tools import runtime

        call = runtime.parse_call(text)
        assert call is not None, f"a valid call was refused: {text}"
        assert call["name"] == "run_python"
        assert call["arguments"]["code"] == "print(137 * 449)"

    @pytest.mark.parametrize("text", [
        "Just prose mentioning run_python and nothing else.",
        "I would use run_python(the function) here.",
        "Nothing tool-shaped at all.",
    ])
    def test_prose_is_not_mistaken_for_a_call(self, text):
        from primnox2.tools import runtime

        assert runtime.parse_call(text) is None, f"prose parsed as a call: {text}"

    def test_an_unknown_name_is_not_a_call(self):
        """The variants are anchored to registered tools, which is the only
        thing keeping `anything({...})` from being executable."""
        from primnox2.tools import runtime

        assert runtime.parse_call('delete_everything({"path": "C:/"})') is None

    def test_the_variant_forms_are_kept_out_of_the_token_stream(self):
        """Otherwise the user watches the call type itself out before the
        runtime quietly runs it."""
        from primnox2.tools import runtime

        f = runtime.StreamFilter()
        raw = 'Let me compute. run_python({"code": "print(137 * 449)"}) Done.'
        shown = "".join(f.feed(ch) for ch in raw) + f.flush()
        assert "run_python" not in shown
        assert "print(" not in shown
        assert "Let me compute." in shown and "Done." in shown

    def test_a_mention_in_prose_does_not_swallow_the_rest(self):
        """The suppression opener is `name({`, not `name(` — otherwise a reply
        that merely mentions run_python(...) is eaten to the end."""
        from primnox2.tools import runtime

        f = runtime.StreamFilter()
        raw = "You could call run_python(code) yourself. Here is the rest."
        shown = "".join(f.feed(ch) for ch in raw) + f.flush()
        assert "Here is the rest." in shown

    def test_an_orphan_closer_with_no_opener_is_stripped_not_shown(self):
        """The defect this closes, found 2026-08-30 live: a degraded
        free-tier turn's entire visible reply was the four characters
        `</tool>` — a closer this filter never opened (so `self._closing`
        was never armed), which walked straight through as if it were
        ordinary prose. A model under strain emitting a stray fragment of
        its own call syntax is a model-quality problem this filter cannot
        prevent; showing that fragment raw to the user is the filter's own
        job to prevent, exactly as much as a matched pair is."""
        from primnox2.tools import runtime

        f = runtime.StreamFilter()
        shown = "".join(f.feed(ch) for ch in "</tool>") + f.flush()
        assert shown == ""

    def test_an_orphan_closer_mid_reply_is_stripped_and_prose_survives(self):
        from primnox2.tools import runtime

        f = runtime.StreamFilter()
        raw = "Saved that for you. </tool> Anything else?"
        shown = "".join(f.feed(ch) for ch in raw) + f.flush()
        assert "</tool>" not in shown
        assert "Saved that for you." in shown
        assert "Anything else?" in shown


# ── Streaming watchdog ────────────────────────────────────────────────────────
class TestBoundedStream:
    """The defect this closes, found 2026-08-30 live: a turn sat in
    `thinking` for 250+ seconds with no resolution, and cancelling it
    (`{"ok": true}` from the API) did nothing. `urlopen(..., timeout=120)`
    bounds each individual socket read, not the request as a whole — a
    provider that trickles even one byte occasionally resets that per-read
    clock indefinitely without ever finishing. `_bounded_stream` puts a
    genuine absolute ceiling on total time, independent of per-read
    timeouts, by force-closing the response from a watchdog thread."""

    class _HangingResponse:
        """A response whose iterator blocks forever until `.close()` is
        called — the same shape a real blocked socket read has: nothing
        else in the process can unblock it except closing the resource."""
        def __init__(self):
            self._closed = threading.Event()

        def __iter__(self):
            return self

        def __next__(self):
            self._closed.wait()          # blocks "forever" until closed
            raise StopIteration          # what a real close() surfaces as EOF

        def close(self):
            self._closed.set()

    def test_a_response_that_never_sends_anything_is_force_closed_and_raises(self):
        from primnox2.models import gateway

        resp = self._HangingResponse()
        with pytest.raises(TimeoutError):
            list(gateway._bounded_stream(resp, deadline_s=0.2))

    def test_a_normal_response_is_unaffected_and_the_watchdog_never_fires(self):
        """The deadline must not cost a healthy, ordinary-speed response
        anything — this is a safety net for a hang, not a rate limiter."""
        from primnox2.models import gateway

        lines = [b"line one\n", b"line two\n", b"line three\n"]
        out = list(gateway._bounded_stream(iter(lines), deadline_s=5.0))
        assert out == lines

    def test_the_watchdog_thread_is_cancelled_on_normal_completion(self):
        """A completed stream must not leave a pending timer around — that
        would be a resource leak on every single successful turn."""
        from primnox2.models import gateway

        before = threading.active_count()
        list(gateway._bounded_stream(iter([b"one\n"]), deadline_s=5.0))
        # The Timer thread is daemonic and cancelled in `finally`; give it
        # one scheduler tick to actually unwind before counting.
        import time as _time
        _time.sleep(0.05)
        assert threading.active_count() <= before


# ── Themed documents ─────────────────────────────────────────────────────────
class TestDocumentThemes:
    """Styling lives in the sandbox helper so the model does not have to write
    it. A 7B writes short scripts and gets long ones wrong; forty lines of
    colour and geometry per deck is exactly the shape of code it fails at."""

    def test_a_deck_carries_the_theme_it_was_given(self, tmp_path):
        from pptx import Presentation

        from primnox2.sandbox import doc_themes

        target = tmp_path / "d.pptx"
        deck = doc_themes.Deck(str(target), theme="midnight", title="Cover")
        deck.slide("A heading", ["first", "second"])
        deck.save()

        p = Presentation(str(target))
        assert len(p.slides) == 2
        assert round(p.slide_width / p.slide_height, 2) == 1.78, "not 16:9"
        body = p.slides[1]
        assert str(body.background.fill.fore_color.rgb) == "060A1A", \
            "the slide background is not the theme's"
        # Found by its text, not by its index. The slide also carries master
        # furniture (footer, page number), and which shape happens to be drawn
        # first is not a contract — asserting on position made this test fail
        # the moment the deck gained a page number, for no defect.
        heading = next(s for s in body.shapes
                       if s.has_text_frame and s.text_frame.text == "A heading")
        assert str(heading.text_frame.paragraphs[0].font.color.rgb) == "6EA8FF"

    def test_every_theme_produces_a_readable_document(self, tmp_path):
        """A palette that renders text the same colour as its background is a
        theme that silently destroys the document."""
        from primnox2.sandbox import doc_themes

        for name, t in doc_themes.THEMES.items():
            assert t["bg"].lower() != t["text"].lower(), f"{name}: text invisible on bg"
            assert t["bg"].lower() != t["primary"].lower(), f"{name}: headings invisible"

    def test_an_unknown_theme_falls_back_rather_than_raising(self):
        from primnox2.sandbox import doc_themes

        assert doc_themes.theme("not-a-theme") == doc_themes.THEMES[doc_themes.DEFAULT_THEME]
        assert doc_themes.theme(None)["bg"]

    def test_a_pdf_and_a_word_file_are_produced_and_parseable(self, tmp_path):
        import docx
        from pypdf import PdfReader

        from primnox2.sandbox import doc_themes

        pdf = tmp_path / "r.pdf"
        report = doc_themes.Report(str(pdf), theme="paper", title="Title")
        report.heading("Section")
        report.text("Some body text that is long enough to be worth reading.")
        report.bullets(["one", "two"])
        report.table([["a", "b"], ["1", "2"]])
        report.save()
        assert len(PdfReader(str(pdf)).pages) >= 1

        docx_path = tmp_path / "w.docx"
        doc_themes.Doc(str(docx_path), theme="sand", title="Title").heading("H").text("Body").save()
        d = docx.Document(str(docx_path))
        assert [p.text for p in d.paragraphs if p.text.strip()] == ["Title", "H", "Body"]

    def test_the_helper_is_installed_for_python_and_hidden_from_results(self, tmp_path):
        """It is ours, not the user's — surfacing it would attach a library to
        every document as though the model had written one."""
        from primnox2.sandbox import snapshots, workspace

        workspace.install_helpers(tmp_path, "python")
        assert (tmp_path / "primnox_docs.py").is_file()

        before = snapshots.snapshot(tmp_path)
        (tmp_path / "real_output.txt").write_text("a result")
        diff = snapshots.diff(before, snapshots.snapshot(tmp_path))
        assert diff["created"] == ["real_output.txt"]

    def test_the_helper_is_not_installed_for_other_runtimes(self, tmp_path):
        from primnox2.sandbox import workspace

        workspace.install_helpers(tmp_path, "node")
        assert not (tmp_path / "primnox_docs.py").exists()


# ── Skills ───────────────────────────────────────────────────────────────────
class TestSkills:
    """Capability instructions the prompt does not pay for until they are used.

    Teaching the model about themed documents inline cost ~209 tokens on every
    turn, including the ones that will never make a document. Four more
    capabilities taught that way and the preamble outweighs most replies.
    """

    def test_the_index_is_cheap_and_the_body_is_not_in_the_prompt(self):
        from primnox2.skills import loader
        from primnox2.tools import runtime

        prompt = runtime.system_prompt()
        skill = loader.get("themed-documents")
        assert skill is not None

        assert "themed-documents" in prompt, "the model cannot ask for what it cannot see"
        # 400 held while two skills shipped. Four more (data-analysis,
        # interactive-apps, memory-and-recall, running-commands) + the design
        # system put the index at 709 chars — roughly 177 tokens on every turn.
        # Still an order of magnitude under inlining all bodies. Two paths to
        # slides now: code generation (35-40% on 0.5B) and design routing
        # (70-85% target). Both live in system prompt index.
        assert len(loader.index()) < 800, "the always-present cost grew"
        assert "from primnox_docs import Deck" not in prompt, \
            "the whole skill is inlined again — the point was that it is not"

    def test_a_request_that_needs_it_selects_it(self):
        from primnox2.skills import loader

        for asked in ("make me a deck about volcanoes",
                      "write a PDF briefing on urban heat",
                      "can you produce a word document",
                      "plot a chart of the results"):
            assert [s.name for s in loader.select(asked)] == ["themed-documents"], \
                f"not selected for: {asked}"

    def test_an_unrelated_request_selects_nothing(self):
        from primnox2.skills import loader

        for asked in ("what is a race condition?",
                      "explain recursion in two sentences",
                      "what is 137 * 449"):
            assert loader.select(asked) == [], f"wrongly selected for: {asked}"

    def test_the_body_carries_what_the_model_needs(self):
        """The deck is inline. Everything else is one `read_skill` away, and the
        body has to name the file or it may as well not ship."""
        from primnox2.skills import loader

        skill = loader.get("themed-documents")
        for expected in ("primnox_docs", "Deck", "midnight", "paper",
                         "layouts.md", "pdf-and-word.md"):
            assert expected in skill.body, f"the skill never mentions {expected}"
        for name, expected in (("layouts.md", ("kpi", "table", "chart", "palette")),
                               ("pdf-and-word.md", ("Report", "Doc", "chart_style"))):
            text = skill.read_asset(name)
            assert text, f"{name} is named by the body but does not ship"
            for token in expected:
                assert token in text, f"{name} never mentions {token}"

    def test_the_smallest_deck_comes_first(self):
        """A small model reads the top of the body and copies whatever is there,
        so the opening has to be a complete runnable deck and the instruction to
        send it as a tool call. Measured on qwen2.5:0.5b, 20 runs each: the
        previous 5k body produced 0 decks, this one produced 8. The difference
        is almost entirely that the directive and the template come first."""
        from primnox2.skills import loader

        body = loader.get("themed-documents").body
        assert "run_python" in body[:200], \
            "the reply-with-a-tool-call instruction is not the first thing said"
        head = body[:900]
        for expected in ("from primnox_docs import Deck", "Deck(", "bullets(",
                         "print(d.save())"):
            assert expected in head, \
                f"the opening example is missing {expected}:\n{head}"

    def test_the_layout_file_warns_about_the_silent_failure(self):
        """`kpi('Metrics', ['Revenue', 'Users'])` unpacks each string into
        characters and saves a deck of nonsense with no error anywhere. Measured:
        two cards reading `e | R | v` and `s | U | e`. The warning lives in
        layouts.md alongside kpi and friends; the body keeps only the happy path."""
        from primnox2.skills import loader

        skill = loader.get("themed-documents")
        layouts = skill.read_asset("layouts.md")
        assert "kpi('Metrics', ['Revenue', 'Users'])" in layouts, \
            "the wrong-shaped call is not shown, so it cannot be recognised"
        assert "tuples" in layouts, \
            "the tuple vs. string hazard is documented where it matters"

    def test_reading_an_unknown_skill_lists_the_real_ones(self):
        """A dead end should say where the road is."""
        from primnox2.tools import runtime
        from primnox2.tools.registry import ToolContext

        result = runtime.execute("read_skill", {"name": "nonsense"}, ToolContext())
        assert result["status"] == "error"
        assert "themed-documents" in result["output"]

    def test_reading_a_real_skill_returns_its_instructions(self):
        from primnox2.tools import runtime
        from primnox2.tools.registry import ToolContext

        result = runtime.execute("read_skill", {"name": "themed-documents"},
                                 ToolContext())
        assert result["status"] == "success"
        assert "from primnox_docs import Deck" in result["output"]

    def test_the_skill_says_which_tool_runs_the_code(self):
        """It showed Python and never said how to execute it. The model wrote
        correct code, filed it with create_workspace, and announced a deck that
        did not exist."""
        from primnox2.skills import loader

        body = loader.get("themed-documents").body
        assert "run_python" in body, "the skill never names the tool that runs it"

    def test_saving_a_script_does_not_read_as_running_it(self):
        from primnox2.tools import builtins

        note = builtins._saved_not_run({"main.py": "print(1)"})
        assert "NOTHING WAS EXECUTED" in note
        assert "run_python" in note

    def test_saving_a_binary_name_to_a_workspace_is_called_out(self):
        """A workspace holds text, so `deck.pptx` saved into one is not a deck."""
        from primnox2.tools import builtins

        note = builtins._saved_not_run({"deck.pptx": "..."})
        assert "deck.pptx" in note and "not a real document" in note

    def test_a_category_word_resolves_to_a_theme_of_that_kind(self):
        """Measured: asked for a light deck, the model passed theme='light',
        the lookup missed, and it silently produced a dark one. 'light' and
        'dark' are how the themes are described, so they are what gets passed."""
        from primnox2.sandbox import doc_themes

        assert doc_themes.theme("light")["dark"] is False
        assert doc_themes.theme("dark")["dark"] is True
        assert doc_themes.theme("light")["bg"] == doc_themes.THEMES["paper"]["bg"]

    def test_an_unknown_theme_says_so_rather_than_falling_back_in_silence(self, capsys):
        from primnox2.sandbox import doc_themes

        doc_themes.theme("chartreuse")
        printed = capsys.readouterr().out
        assert "no theme called 'chartreuse'" in printed
        assert "midnight" in printed, "it does not say what it used instead"
