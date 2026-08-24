"""The auditor is itself tested, by planting each violation it claims to catch.

An auditor that has never caught a known-bad input is a rubber stamp. Every
check below builds a deck containing exactly one deliberate defect and asserts
that check fires — and the clean-deck test asserts none of them fire on a deck
built to the system, which is the half that catches an auditor that simply
reports everything.
"""
from __future__ import annotations

import pytest

pptx = pytest.importorskip("pptx")

from pptx import Presentation                      # noqa: E402
from pptx.dml.color import RGBColor                # noqa: E402
from pptx.util import Pt                           # noqa: E402

from deck_audit import DesignSystem, audit, summarise   # noqa: E402

SYSTEM = DesignSystem.torture_test()
NAVY = RGBColor(0x0F, 0x17, 0x2A)
TEXT = RGBColor(0xF8, 0xFA, 0xFC)


def _deck(tmp_path, build) -> str:
    prs = Presentation()
    prs.slide_width = Pt(960)
    prs.slide_height = Pt(540)
    build(prs)
    path = tmp_path / "deck.pptx"
    prs.save(str(path))
    return str(path)


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _text(slide, left, top, width, height, text, size=22, color=TEXT, font="Inter"):
    box = slide.shapes.add_textbox(Pt(left), Pt(top), Pt(width), Pt(height))
    tf = box.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.name = font
    run.font.color.rgb = color
    return box


def _conforming_slide(prs, title="Section Title", body="A compliant line of body copy."):
    slide = _blank(prs)
    _text(slide, 64, 64, 832, 64, title, size=40)
    _text(slide, 64, 160, 832, 64, body, size=22)
    return slide


def _checks(findings):
    return {f.check for f in findings}


# ── the clean case ───────────────────────────────────────────────────────────
def test_a_conforming_deck_reports_nothing(tmp_path):
    """The half that catches an auditor which flags everything.

    Three slides, not four: the rhythm rule allows a run of three, and a
    fixture that tripped it would be testing the wrong thing here.
    """
    def build(prs):
        for i in range(3):
            _conforming_slide(prs, f"Section {i}", f"Body copy for section {i}.")

    findings = audit(_deck(tmp_path, build), SYSTEM)
    assert findings == [], "\n".join(str(f) for f in findings)


def test_catches_a_monotonous_run(tmp_path):
    """Nothing per-slide is wrong with six identical layouts in a row, and it is
    still a bad deck. This is the only check that looks at the deck rather than
    at a slide."""
    def build(prs):
        for i in range(6):
            _conforming_slide(prs, f"Section {i}", f"Body copy {i}.")

    assert "monotonous-run" in _checks(audit(_deck(tmp_path, build), SYSTEM))


def test_a_varied_deck_does_not_trip_the_rhythm_rule(tmp_path):
    def build(prs):
        for i in range(6):
            slide = _blank(prs)
            _text(slide, 64, 64, 832, 64, f"Section {i}", size=40)
            # Alternating shape counts stand in for alternating layouts.
            if i % 2:
                _text(slide, 64, 160, 400, 64, "Left column", size=22)
                _text(slide, 496, 160, 400, 64, "Right column", size=22)
            else:
                _text(slide, 64, 160, 832, 64, "One column", size=22)

    assert "monotonous-run" not in _checks(audit(_deck(tmp_path, build), SYSTEM))


# ── one planted defect per check ─────────────────────────────────────────────
def test_catches_color_drift(tmp_path):
    def build(prs):
        slide = _conforming_slide(prs)
        _text(slide, 64, 288, 400, 48, "Off-palette", color=RGBColor(0xFF, 0x00, 0xFF))

    assert "color-drift" in _checks(audit(_deck(tmp_path, build), SYSTEM))


def test_catches_font_drift(tmp_path):
    def build(prs):
        slide = _conforming_slide(prs)
        _text(slide, 64, 288, 400, 48, "Wrong face", font="Comic Sans MS")

    assert "font-drift" in _checks(audit(_deck(tmp_path, build), SYSTEM))


def test_catches_a_size_off_the_type_scale(tmp_path):
    def build(prs):
        slide = _conforming_slide(prs)
        _text(slide, 64, 288, 400, 48, "Odd size", size=23)

    assert "type-scale" in _checks(audit(_deck(tmp_path, build), SYSTEM))


def test_catches_text_below_the_readability_floor(tmp_path):
    """The prompt says never shrink below readability, so 9pt is a defect even
    though shrinking would technically resolve an overflow."""
    def build(prs):
        slide = _conforming_slide(prs)
        _text(slide, 64, 288, 400, 48, "Too small to read", size=9)

    assert "below-readable" in _checks(audit(_deck(tmp_path, build), SYSTEM))


def test_catches_a_missing_title(tmp_path):
    def build(prs):
        slide = _blank(prs)
        # Body-sized type only, low on the slide: nothing reads as a heading.
        _text(slide, 64, 400, 400, 48, "Body only, no heading.", size=22)
        _text(slide, 64, 464, 400, 48, "A second body line.", size=22)

    assert "missing-title" in _checks(audit(_deck(tmp_path, build), SYSTEM))


def test_a_centred_hero_headline_counts_as_a_title(tmp_path):
    """Hero and cover slides put the headline in the middle by design. Judging
    titles by position alone reported every one of them as untitled."""
    def build(prs):
        slide = _blank(prs)
        _text(slide, 64, 232, 832, 80, "The Headline", size=54)
        _text(slide, 64, 328, 832, 48, "A supporting subtitle.", size=22)

    assert "missing-title" not in _checks(audit(_deck(tmp_path, build), SYSTEM))


def test_catches_overlapping_objects(tmp_path):
    def build(prs):
        slide = _conforming_slide(prs)
        _text(slide, 64, 288, 400, 64, "First block")
        _text(slide, 96, 304, 400, 64, "Second block on top of it")

    assert "overlap" in _checks(audit(_deck(tmp_path, build), SYSTEM))


def test_catches_grid_misalignment(tmp_path):
    def build(prs):
        slide = _conforming_slide(prs)
        _text(slide, 67, 291, 400, 48, "Off grid")   # 67 and 291 are not multiples of 8

    assert "grid-misalign" in _checks(audit(_deck(tmp_path, build), SYSTEM))


def test_catches_accidental_duplicate_slides(tmp_path):
    def build(prs):
        _conforming_slide(prs, "Identical", "Same body.")
        _conforming_slide(prs, "Identical", "Same body.")

    assert "duplicate-slide" in _checks(audit(_deck(tmp_path, build), SYSTEM))


def test_catches_text_overflow(tmp_path):
    """auto_size NONE: nothing adjusts, so the text is genuinely clipped."""
    from pptx.enum.text import MSO_AUTO_SIZE

    def build(prs):
        slide = _conforming_slide(prs)
        box = _text(slide, 64, 288, 160, 24, "x", size=22)
        box.text_frame.auto_size = MSO_AUTO_SIZE.NONE
        box.text_frame.paragraphs[0].runs[0].text = "overflowing " * 40

    assert "text-overflow" in _checks(audit(_deck(tmp_path, build), SYSTEM))


def test_reports_autogrow_separately_from_overflow(tmp_path):
    """SHAPE_TO_FIT_TEXT does not clip - it silently expands the box past the
    geometry every other check reads, which is a different problem and must not
    be filed as an overflow."""
    from pptx.enum.text import MSO_AUTO_SIZE

    def build(prs):
        slide = _conforming_slide(prs)
        box = _text(slide, 64, 288, 160, 24, "x", size=22)
        box.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        box.text_frame.paragraphs[0].runs[0].text = "expanding " * 40

    checks = _checks(audit(_deck(tmp_path, build), SYSTEM))
    assert "autogrow-geometry" in checks
    assert "text-overflow" not in checks


def test_shrink_to_fit_is_not_reported_as_overflow(tmp_path):
    """TEXT_TO_FIT_SHAPE shrinks the type; the readability floor check owns
    that case, so reporting it here too would double-count one defect."""
    from pptx.enum.text import MSO_AUTO_SIZE

    def build(prs):
        slide = _conforming_slide(prs)
        box = _text(slide, 64, 288, 160, 24, "x", size=22)
        box.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        box.text_frame.paragraphs[0].runs[0].text = "shrinking " * 40

    assert "text-overflow" not in _checks(audit(_deck(tmp_path, build), SYSTEM))


def test_does_not_invent_overflow_for_ordinary_text(tmp_path):
    """The estimator errs toward silence on purpose: an auditor that cries wolf
    gets ignored, and then the whole checklist is worthless."""
    def build(prs):
        slide = _blank(prs)
        _text(slide, 64, 64, 832, 64, "A Perfectly Ordinary Heading", size=40)
        _text(slide, 64, 160, 832, 96,
              "A sentence of body copy that comfortably fits inside its box.", size=22)

    assert "text-overflow" not in _checks(audit(_deck(tmp_path, build), SYSTEM))


def test_catches_too_many_font_families(tmp_path):
    def build(prs):
        slide = _conforming_slide(prs)
        for i, face in enumerate(["Georgia", "Verdana", "Tahoma", "Courier New"]):
            _text(slide, 64, 288 + i * 56, 400, 48, f"Face {i}", font=face)

    assert "font-count" in _checks(audit(_deck(tmp_path, build), SYSTEM))


# ── summary ──────────────────────────────────────────────────────────────────
def test_summary_separates_errors_from_warnings(tmp_path):
    def build(prs):
        slide = _conforming_slide(prs)
        _text(slide, 67, 291, 400, 48, "Off grid but legal colour", size=22)

    findings = audit(_deck(tmp_path, build), SYSTEM)
    summary = summarise(findings, 1)
    assert summary["warnings"] >= 1
    assert summary["clean"] == (summary["errors"] == 0)
