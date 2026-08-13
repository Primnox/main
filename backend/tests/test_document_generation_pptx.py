"""PPTX generation quality — is the deck Primnox produces actually usable?

"The file opens" is the floor. A deck can open perfectly and still be
unshippable: text spilling out of its box, two shapes stacked on top of each
other, a logo squashed to 3:1, a slide hanging off the canvas. None of that
is visible to a test that only checks `len(prs.slides)`.

So this file is built in two layers.

**Layer 1 — checkers** (`deck_report` and the functions above it). Real
geometric analysis of a .pptx: shape rectangles in EMU, estimated text
extent, image aspect ratios. Pure functions over a `python-pptx`
Presentation, running in milliseconds.

**Layer 2 — subjects.** The checkers are pointed at two kinds of deck:

* decks built right here with `python-pptx`, deterministically, *including
  deliberately broken ones*. These run by default and are what prove the
  checkers detect the defects they claim to — a checker that never fires is
  worse than no checker at all.
* real model-generated decks, under `@pytest.mark.live` (see
  `test_document_generation_live.py`), which cost an API call and a sandbox
  run and are excluded from the default run.

The text-extent estimator is explicitly a heuristic: PowerPoint's real
layout engine isn't reproducible here without rendering. It is calibrated to
catch *gross* overflow — a paragraph two or three times too tall for its box
— not to adjudicate the last few points. The tests below pin both halves of
that: it fires on gross overflow, and it stays quiet on comfortable text.
"""

import math
from pathlib import Path

import pytest

pytest.importorskip("pptx")
pytest.importorskip("PIL")

from pptx import Presentation  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa: E402
from pptx.enum.text import MSO_AUTO_SIZE  # noqa: E402
from pptx.util import Emu, Inches, Pt  # noqa: E402

EMU_PER_PT = 12700
EMU_PER_INCH = 914400

# Body-text default for runs that inherit their size from the layout.
# python-pptx reports None there, and resolving the real inherited value
# means walking layout -> master -> theme. 18pt is PowerPoint's body default,
# and erring small makes the overflow checker *less* trigger-happy — the
# right direction for a heuristic.
DEFAULT_FONT_PT = 18.0

# Average glyph advance as a fraction of font size for a proportional sans
# face, calibrated against Calibri/Arial lowercase-heavy prose.
AVG_GLYPH_WIDTH_EM = 0.5

LINE_SPACING = 1.2


# ─────────────────────────────────────────────────────────────────────────
# Layer 1: checkers
# ─────────────────────────────────────────────────────────────────────────

def _rect(shape):
    """(left, top, width, height) in EMU, or None when the shape inherits its
    geometry from the layout — python-pptx reports None, and an inheriting
    placeholder has no geometry of its own to judge."""
    vals = (shape.left, shape.top, shape.width, shape.height)
    if any(v is None for v in vals):
        return None
    return tuple(int(v) for v in vals)


def _font_pt(paragraph) -> float:
    for run in paragraph.runs:
        if run.font.size is not None:
            return run.font.size.pt
    if paragraph.font.size is not None:
        return paragraph.font.size.pt
    return DEFAULT_FONT_PT


def estimated_text_height_emu(shape) -> int:
    """Roughly how tall the shape's text wants to be, in EMU: wrap each
    paragraph at the shape's usable width, then stack the resulting lines."""
    tf = shape.text_frame
    rect = _rect(shape)
    if rect is None:
        return 0
    usable_emu = rect[2] - (tf.margin_left or 0) - (tf.margin_right or 0)
    usable_pt = max(1.0, usable_emu / EMU_PER_PT)

    total_pt = 0.0
    for para in tf.paragraphs:
        size_pt = _font_pt(para)
        chars_per_line = max(1, int(usable_pt / (size_pt * AVG_GLYPH_WIDTH_EM)))
        lines = max(1, math.ceil(len(para.text or "") / chars_per_line))
        total_pt += lines * size_pt * LINE_SPACING

    insets = (tf.margin_top or 0) + (tf.margin_bottom or 0)
    return int(total_pt * EMU_PER_PT) + insets


def overflow_ratio(shape) -> float:
    """estimated text height / shape height. Above 1.0 means it spills.

    Zero for shapes that cannot overflow: no text frame, empty text, or an
    autofit mode where PowerPoint grows the shape around its text.
    """
    if not shape.has_text_frame or not shape.text_frame.text.strip():
        return 0.0
    if shape.text_frame.auto_size == MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT:
        return 0.0
    rect = _rect(shape)
    if rect is None or rect[3] <= 0:
        return 0.0
    return estimated_text_height_emu(shape) / rect[3]


def overflowing_shapes(slide, tolerance: float = 1.15) -> list:
    """[(shape_name, ratio)] for every shape whose text grossly exceeds it."""
    return [(s.name, round(overflow_ratio(s), 2))
            for s in slide.shapes if overflow_ratio(s) > tolerance]


def _is_visible(shape) -> bool:
    """An empty placeholder occupies a rectangle but paints nothing, so it
    cannot visually collide with anything. Excluding it stops the overlap
    check firing on layout scaffolding the deck simply didn't fill in."""
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return True
    if shape.has_text_frame:
        return bool(shape.text_frame.text.strip())
    return True


def _intersection_area(a, b) -> int:
    ax, ay, aw, ah = a
    bx, by, bw, bh = b
    dx = min(ax + aw, bx + bw) - max(ax, bx)
    dy = min(ay + ah, by + bh) - max(ay, by)
    return dx * dy if dx > 0 and dy > 0 else 0


def overlapping_pairs(slide, min_fraction: float = 0.05) -> list:
    """[(name_a, name_b, fraction)] for shape pairs whose rectangles collide.

    `fraction` is of the *smaller* shape's area, so a caption dumped on top
    of a full-bleed background counts as a heavy collision rather than a
    negligible one. Thresholded rather than zero because real decks routinely
    share a pixel or two of border.
    """
    shapes = [(s, _rect(s)) for s in slide.shapes if _is_visible(s)]
    shapes = [(s, r) for s, r in shapes if r and r[2] > 0 and r[3] > 0]
    out = []
    for i in range(len(shapes)):
        for j in range(i + 1, len(shapes)):
            (sa, ra), (sb, rb) = shapes[i], shapes[j]
            area = _intersection_area(ra, rb)
            if not area:
                continue
            fraction = area / min(ra[2] * ra[3], rb[2] * rb[3])
            if fraction > min_fraction:
                out.append((sa.name, sb.name, round(fraction, 3)))
    return out


def out_of_bounds_shapes(prs, slide, tolerance_emu: int = EMU_PER_INCH // 100) -> list:
    """[(name, reason)] for shapes hanging off the slide canvas."""
    out = []
    for shape in slide.shapes:
        rect = _rect(shape)
        if rect is None:
            continue
        left, top, width, height = rect
        if left < -tolerance_emu:
            out.append((shape.name, f"left={left}"))
        if top < -tolerance_emu:
            out.append((shape.name, f"top={top}"))
        if left + width > prs.slide_width + tolerance_emu:
            out.append((shape.name, f"right={left + width} > {prs.slide_width}"))
        if top + height > prs.slide_height + tolerance_emu:
            out.append((shape.name, f"bottom={top + height} > {prs.slide_height}"))
    return out


def picture_aspect_error(picture) -> float:
    """Relative error between the placed rectangle's aspect ratio and the
    source image's, accounting for crop. 0.0 is faithful; 0.25 means it has
    been squashed or stretched by a quarter."""
    src_w, src_h = picture.image.size
    keep_w = 1.0 - (picture.crop_left or 0) - (picture.crop_right or 0)
    keep_h = 1.0 - (picture.crop_top or 0) - (picture.crop_bottom or 0)
    if keep_w <= 0 or keep_h <= 0 or src_h <= 0:
        return 0.0
    source_ratio = (src_w * keep_w) / (src_h * keep_h)
    placed_ratio = picture.width / picture.height
    return abs(placed_ratio - source_ratio) / source_ratio


def distorted_pictures(slide, tolerance: float = 0.02) -> list:
    """[(name, error)] for images placed at the wrong aspect ratio."""
    return [(s.name, round(picture_aspect_error(s), 4))
            for s in slide.shapes
            if s.shape_type == MSO_SHAPE_TYPE.PICTURE
            and picture_aspect_error(s) > tolerance]


def slide_text(slide) -> str:
    chunks = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            chunks.append(shape.text_frame.text)
        if getattr(shape, "has_table", False) and shape.has_table:
            for row in shape.table.rows:
                chunks.extend(c.text for c in row.cells)
    return "\n".join(c for c in chunks if c)


def deck_report(path) -> dict:
    """Every checker, over a whole deck. The single entry point the live
    tests use, so a real generated deck is judged by exactly the rules the
    fixtures below prove correct."""
    prs = Presentation(str(path))
    slides = []
    for index, slide in enumerate(prs.slides, start=1):
        slides.append({
            "index": index,
            "text": slide_text(slide),
            "shape_count": len(slide.shapes),
            "overflow": overflowing_shapes(slide),
            "overlap": overlapping_pairs(slide),
            "out_of_bounds": out_of_bounds_shapes(prs, slide),
            "distorted_images": distorted_pictures(slide),
        })
    return {"slide_count": len(prs.slides),
            "slide_width": prs.slide_width,
            "slide_height": prs.slide_height,
            "slides": slides}


def assert_deck_is_presentable(report: dict):
    """One assertion with a readable failure, so a live run says *what* is
    wrong with the deck rather than just that something is."""
    problems = []
    for s in report["slides"]:
        for name, ratio in s["overflow"]:
            problems.append(f"slide {s['index']}: '{name}' text overflows its box ({ratio}x)")
        for a, b, frac in s["overlap"]:
            problems.append(f"slide {s['index']}: '{a}' and '{b}' overlap ({frac:.0%})")
        for name, why in s["out_of_bounds"]:
            problems.append(f"slide {s['index']}: '{name}' is off-canvas ({why})")
        for name, err in s["distorted_images"]:
            problems.append(f"slide {s['index']}: image '{name}' distorted by {err:.1%}")
    assert not problems, "deck is not presentable:\n  " + "\n  ".join(problems)


# ─────────────────────────────────────────────────────────────────────────
# Deck builders — every fixture is code, no committed binaries
# ─────────────────────────────────────────────────────────────────────────

AI_SLIDES = [
    ("What Machine Learning Actually Does",
     "Finds patterns in data instead of following written rules."),
    ("Supervised Learning",
     "Labelled examples in, a predictor out. The workhorse of production ML."),
    ("Unsupervised Learning",
     "No labels. Structure has to be discovered, not confirmed."),
    ("Neural Networks",
     "Layers of weighted sums, trained by gradient descent."),
    ("Transformers",
     "Attention replaced recurrence and made scale pay off."),
    ("Large Language Models",
     "Next-token prediction, at a scale where it starts to generalise."),
    ("Fine-Tuning",
     "Adapting a general model to one narrow job, cheaply."),
    ("Retrieval Augmentation",
     "Give the model the document rather than hoping it memorised it."),
    ("Evaluation",
     "If you cannot measure the output, you cannot improve it."),
    ("Where It Breaks",
     "Confident wrong answers, stale knowledge, and cost at scale."),
]


def _add_content_slide(prs, title, body):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    for left, top, width, height, text, size in (
        (0.6, 0.5, 8.8, 1.2, title, 32),
        (0.6, 2.2, 8.8, 3.0, body, 18),
    ):
        box = slide.shapes.add_textbox(Inches(left), Inches(top),
                                       Inches(width), Inches(height))
        box.text_frame.word_wrap = True
        # Fixed height on purpose. python-pptx's add_textbox defaults to
        # spAutoFit (see TestPythonPptxDefaults), under which the box grows to
        # fit and the overflow checker correctly reports 0.0 — which would
        # make the false-positive guard below pass vacuously.
        box.text_frame.auto_size = MSO_AUTO_SIZE.NONE
        box.text_frame.text = text
        box.text_frame.paragraphs[0].runs[0].font.size = Pt(size)
    return slide


def build_clean_deck(path: Path, slides=AI_SLIDES) -> Path:
    """A deck that should pass every checker: blank layout, hand-placed
    non-overlapping boxes, text sized to fit."""
    prs = Presentation()
    for title, body in slides:
        _add_content_slide(prs, title, body)
    prs.save(str(path))
    return path


def build_overflowing_deck(path: Path) -> Path:
    """One slide, one small box, far too much text at far too large a size."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(0.8))
    box.text_frame.word_wrap = True
    box.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    box.text_frame.text = (
        "Retrieval augmented generation gives the model the document instead of "
        "hoping it memorised it during pre-training, which matters most when the "
        "corpus changes weekly and the answer has to cite a source the reader can "
        "open and verify before acting on it."
    )
    box.text_frame.paragraphs[0].runs[0].font.size = Pt(28)
    prs.save(str(path))
    return path


def build_overlapping_deck(path: Path) -> Path:
    """Two filled text boxes sitting almost squarely on top of one another."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for left in (Inches(1.0), Inches(1.4)):
        box = slide.shapes.add_textbox(left, Inches(1.0), Inches(4), Inches(2))
        box.text_frame.text = "Quarterly results"
        box.text_frame.paragraphs[0].runs[0].font.size = Pt(14)
    prs.save(str(path))
    return path


def build_offslide_deck(path: Path) -> Path:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(prs.slide_width - Inches(1), Inches(1),
                                   Inches(4), Inches(1))
    box.text_frame.text = "Hanging off the right edge"
    box.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
    prs.save(str(path))
    return path


def make_image(path: Path, width: int, height: int, colour=(30, 90, 160)) -> Path:
    from PIL import Image
    Image.new("RGB", (width, height), colour).save(str(path))
    return path


# ─────────────────────────────────────────────────────────────────────────
# 18. Basic PPTX — valid file, correct slide count, real text
# ─────────────────────────────────────────────────────────────────────────

class TestDeckStructureChecks:
    """Proves the structural assertions the live deck test leans on actually
    tell a good deck from a bad one."""

    def test_a_ten_slide_deck_reads_back_correctly(self, tmp_path):
        report = deck_report(build_clean_deck(tmp_path / "ai.pptx"))
        assert report["slide_count"] == 10
        for slide, (title, body) in zip(report["slides"], AI_SLIDES):
            assert title in slide["text"]
            assert body in slide["text"]

    def test_every_slide_carries_substantive_text(self, tmp_path):
        report = deck_report(build_clean_deck(tmp_path / "ai.pptx"))
        for slide in report["slides"]:
            assert len(slide["text"].split()) >= 5, (
                f"slide {slide['index']} is nearly empty: {slide['text']!r}")

    def test_an_empty_slide_is_detected(self, tmp_path):
        """The check above has to fail on a deck padded with blank slides —
        a common way a generated deck reaches its slide count dishonestly."""
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        path = tmp_path / "blank.pptx"
        prs.save(str(path))
        assert deck_report(path)["slides"][0]["text"].strip() == ""

    def test_a_truncated_file_is_not_a_valid_pptx(self, tmp_path):
        """A half-written deck is the signature of a generation run cut off
        mid-flight; opening it must raise, not silently succeed."""
        raw = build_clean_deck(tmp_path / "ai.pptx").read_bytes()
        broken = tmp_path / "broken.pptx"
        broken.write_bytes(raw[: len(raw) // 2])
        with pytest.raises(Exception):
            Presentation(str(broken))

    def test_an_empty_file_is_not_a_valid_pptx(self, tmp_path):
        """The shape a failed generation leaves behind when the script
        created the file and then died before writing anything."""
        empty = tmp_path / "empty.pptx"
        empty.write_bytes(b"")
        with pytest.raises(Exception):
            Presentation(str(empty))

    def test_a_clean_deck_passes_the_full_presentability_gate(self, tmp_path):
        assert_deck_is_presentable(deck_report(build_clean_deck(tmp_path / "ai.pptx")))


# ─────────────────────────────────────────────────────────────────────────
# 19. Layout stress — overflow and overlap, computed from geometry
# ─────────────────────────────────────────────────────────────────────────

class TestOverflowDetection:
    def test_comfortable_text_does_not_register_as_overflow(self, tmp_path):
        """The false-positive guard. If this fires, the estimator is too
        aggressive and every real deck would look broken."""
        report = deck_report(build_clean_deck(tmp_path / "ai.pptx"))
        for slide in report["slides"]:
            assert slide["overflow"] == [], f"slide {slide['index']}: {slide['overflow']}"

    def test_gross_overflow_is_caught(self, tmp_path):
        overflow = deck_report(build_overflowing_deck(tmp_path / "over.pptx"))["slides"][0]["overflow"]
        assert overflow, "a 28pt paragraph crammed into a 3x0.8in box went undetected"
        assert overflow[0][1] > 2.0, f"expected a large ratio, got {overflow}"

    def test_the_ratio_scales_with_the_amount_of_text(self, tmp_path):
        """Confirms the estimator measures something, rather than returning a
        constant that happens to trip the threshold."""
        prs = Presentation()
        ratios = []
        for repeats in (1, 8):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
            box.text_frame.word_wrap = True
            box.text_frame.auto_size = MSO_AUTO_SIZE.NONE
            box.text_frame.text = "attention is all you need. " * repeats
            box.text_frame.paragraphs[0].runs[0].font.size = Pt(18)
            ratios.append(overflow_ratio(box))
        assert ratios[1] > ratios[0] * 4

    def test_a_dense_slide_of_bullets_is_measured_paragraph_by_paragraph(self, tmp_path):
        """The realistic layout-stress case: not one huge paragraph but
        fifteen bullets, each individually short. Height has to accumulate."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(0.6), Inches(1.0), Inches(8.8), Inches(1.5))
        box.text_frame.word_wrap = True
        box.text_frame.auto_size = MSO_AUTO_SIZE.NONE
        box.text_frame.text = "Bullet number 1 about model evaluation"
        box.text_frame.paragraphs[0].runs[0].font.size = Pt(18)
        for n in range(2, 16):
            para = box.text_frame.add_paragraph()
            para.text = f"Bullet number {n} about model evaluation"
            para.font.size = Pt(18)
        assert overflow_ratio(box) > 2.0, (
            "fifteen 18pt bullets in a 1.5in box must read as overflow")

    def test_autofit_shapes_are_exempt(self, tmp_path):
        """SHAPE_TO_FIT_TEXT means PowerPoint grows the box, so exceeding the
        authored height is not a defect."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(0.5))
        box.text_frame.word_wrap = True
        box.text_frame.auto_size = MSO_AUTO_SIZE.NONE
        box.text_frame.text = "a very long line of text " * 12
        box.text_frame.paragraphs[0].runs[0].font.size = Pt(24)
        assert overflow_ratio(box) > 1.15
        box.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        assert overflow_ratio(box) == 0.0

    def test_empty_shapes_score_zero(self, tmp_path):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        empty = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(0.2))
        assert overflow_ratio(empty) == 0.0

    def test_a_narrow_column_overflows_sooner_than_a_wide_one(self, tmp_path):
        """Wrapping must depend on width, not just character count — the
        difference between a two-column layout working and not."""
        prs = Presentation()
        text = "Retrieval augmented generation grounds the answer in a source document."
        ratios = []
        for width in (Inches(1.5), Inches(8.0)):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            box = slide.shapes.add_textbox(Inches(0.5), Inches(1), width, Inches(0.6))
            box.text_frame.word_wrap = True
            box.text_frame.auto_size = MSO_AUTO_SIZE.NONE
            box.text_frame.text = text
            box.text_frame.paragraphs[0].runs[0].font.size = Pt(18)
            ratios.append(overflow_ratio(box))
        assert ratios[0] > ratios[1] * 3


class TestPythonPptxDefaults:
    """Library defaults that quietly change what a generated deck does.

    These aren't Primnox bugs — they're properties of the tool the pptx skill
    builds with, and each one is a trap a generation script can fall into
    without producing any error.
    """

    def test_add_textbox_defaults_to_autofit(self, tmp_path):
        """`add_textbox` emits `<a:spAutoFit/>`, so a box silently grows to
        fit its text instead of honouring the height it was given. Good for
        never clipping, bad when the layout depended on that height — and it
        means an overflow check that ignores autofit sees nothing anywhere."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(0.4))
        assert box.text_frame.auto_size == MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

    def test_add_textbox_defaults_to_no_word_wrap(self, tmp_path):
        """The other half of the same trap: `wrap="none"`, so long text runs
        off the slide in a single line rather than wrapping. A generator has
        to set `word_wrap = True` explicitly."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(0.4))
        assert box.text_frame.word_wrap is False

    def test_run_font_size_is_none_until_set(self, tmp_path):
        """Why DEFAULT_FONT_PT exists: an unset run reports None rather than
        the size it will actually render at."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(0.4))
        box.text_frame.text = "hello"
        assert box.text_frame.paragraphs[0].runs[0].font.size is None

    def test_layout_placeholders_report_real_geometry(self, tmp_path):
        """The bounds and overlap checks depend on placeholders from a real
        layout having usable coordinates rather than inheriting None."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        assert slide.shapes, "expected title + content placeholders"
        assert all(_rect(s) is not None for s in slide.shapes)

    def test_default_deck_is_four_by_three(self, tmp_path):
        """python-pptx's default template is 10x7.5in (4:3), not 16:9. A deck
        authored against 13.33in coordinates would run off this canvas — the
        bounds checker is what catches that."""
        prs = Presentation()
        assert prs.slide_width == Inches(10)
        assert prs.slide_height == Inches(7.5)


class TestOverlapDetection:
    def test_a_tidy_deck_has_no_collisions(self, tmp_path):
        report = deck_report(build_clean_deck(tmp_path / "ai.pptx"))
        for slide in report["slides"]:
            assert slide["overlap"] == [], f"slide {slide['index']}: {slide['overlap']}"

    def test_stacked_shapes_are_caught(self, tmp_path):
        overlap = deck_report(build_overlapping_deck(tmp_path / "clash.pptx"))["slides"][0]["overlap"]
        assert overlap, "two boxes sharing most of their area went undetected"
        assert overlap[0][2] > 0.5

    def test_touching_edges_are_not_an_overlap(self, tmp_path):
        """Boxes flush against each other are good layout, not a collision."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        for left in (Inches(1), Inches(4)):
            box = slide.shapes.add_textbox(left, Inches(1), Inches(3), Inches(1))
            box.text_frame.text = "column"
        assert overlapping_pairs(slide) == []

    def test_a_hairline_overlap_is_tolerated(self, tmp_path):
        """A couple of thousand EMU (well under a point) is a rounding
        artefact, not a defect."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        a = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
        a.text_frame.text = "left"
        b = slide.shapes.add_textbox(Emu(int(Inches(4)) - 2000), Inches(1),
                                     Inches(3), Inches(1))
        b.text_frame.text = "right"
        assert overlapping_pairs(slide) == []

    def test_a_caption_dropped_on_a_background_is_caught(self, tmp_path):
        """Scored against the *smaller* shape, so a small label buried in a
        full-bleed panel reads as a full collision rather than 4% of the
        panel — which is what it looks like to a reader."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        panel = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(10), Inches(7.5))
        panel.text_frame.text = "background panel"
        label = slide.shapes.add_textbox(Inches(4), Inches(3), Inches(2), Inches(0.5))
        label.text_frame.text = "caption"
        pairs = overlapping_pairs(slide)
        assert pairs and pairs[0][2] > 0.95

    def test_empty_placeholders_do_not_count_as_collisions(self, tmp_path):
        """A layout's unfilled title placeholder paints nothing; flagging it
        would make every template-based deck look broken."""
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[1])  # title + content, unfilled
        assert overlapping_pairs(prs.slides[0]) == []


# ─────────────────────────────────────────────────────────────────────────
# 20. Images — aspect ratio preserved, and inside the slide
# ─────────────────────────────────────────────────────────────────────────

class TestImagePlacement:
    def test_a_faithfully_placed_image_reports_no_distortion(self, tmp_path):
        img = make_image(tmp_path / "wide.png", 1600, 900)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(str(img), Inches(1), Inches(1),
                                 width=Inches(4), height=Inches(2.25))  # 16:9
        path = tmp_path / "ok.pptx"
        prs.save(str(path))
        assert deck_report(path)["slides"][0]["distorted_images"] == []

    def test_a_squashed_image_is_caught(self, tmp_path):
        """The classic: a 16:9 photo forced into a square hole."""
        img = make_image(tmp_path / "wide.png", 1600, 900)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(str(img), Inches(1), Inches(1),
                                 width=Inches(3), height=Inches(3))
        distorted = deck_report(_save(prs, tmp_path / "squashed.pptx"))["slides"][0]["distorted_images"]
        assert distorted, "a 16:9 image placed 1:1 went undetected"
        # 16:9 (1.778) rendered at 1:1 is a 44% relative error.
        assert distorted[0][1] > 0.4

    def test_a_stretched_portrait_is_caught(self, tmp_path):
        img = make_image(tmp_path / "tall.png", 600, 1200)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        pic = slide.shapes.add_picture(str(img), Inches(1), Inches(0.5),
                                       width=Inches(5), height=Inches(2))
        assert picture_aspect_error(pic) > 3.0

    def test_a_two_percent_nudge_stays_within_tolerance(self, tmp_path):
        """Generators round to whole points constantly; the checker must not
        fire on that or it is useless on real output."""
        img = make_image(tmp_path / "wide.png", 1600, 900)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        pic = slide.shapes.add_picture(str(img), Inches(1), Inches(1),
                                       width=Inches(4), height=Emu(int(Inches(2.25) * 1.01)))
        assert picture_aspect_error(pic) < 0.02
        assert distorted_pictures(slide) == []

    def test_width_only_placement_preserves_ratio(self, tmp_path):
        """python-pptx derives the height when only width is given — the safe
        idiom, worth pinning so a generator can rely on it."""
        img = make_image(tmp_path / "tall.png", 600, 1000)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        pic = slide.shapes.add_picture(str(img), Inches(1), Inches(0.5), width=Inches(2))
        assert picture_aspect_error(pic) < 0.01

    def test_cropping_is_accounted_for(self, tmp_path):
        """A cropped image legitimately has a different placed ratio; naive
        ratio-checking would call every cropped photo distorted."""
        img = make_image(tmp_path / "square.png", 1000, 1000)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        pic = slide.shapes.add_picture(str(img), Inches(1), Inches(1),
                                       width=Inches(4), height=Inches(2))
        assert picture_aspect_error(pic) > 0.9  # 2:1 placement of a 1:1 source
        pic.crop_top = 0.25
        pic.crop_bottom = 0.25          # visible source is now 1000x500 = 2:1
        assert picture_aspect_error(pic) < 0.01

    def test_an_image_inside_the_canvas_is_in_bounds(self, tmp_path):
        img = make_image(tmp_path / "i.png", 800, 600)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(str(img), Inches(1), Inches(1), width=Inches(4))
        assert deck_report(_save(prs, tmp_path / "in.pptx"))["slides"][0]["out_of_bounds"] == []

    def test_an_image_running_off_the_edge_is_caught(self, tmp_path):
        img = make_image(tmp_path / "i.png", 800, 600)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(str(img), prs.slide_width - Inches(1),
                                 Inches(1), width=Inches(4))
        problems = deck_report(_save(prs, tmp_path / "out.pptx"))["slides"][0]["out_of_bounds"]
        assert problems and "right=" in problems[0][1]

    def test_a_full_bleed_image_is_in_bounds(self, tmp_path):
        """Edge-to-edge is a design choice, not an error — the bounds check
        must not punish it."""
        img = make_image(tmp_path / "i.png", 1600, 1200)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(str(img), 0, 0,
                                 width=prs.slide_width, height=prs.slide_height)
        assert deck_report(_save(prs, tmp_path / "bleed.pptx"))["slides"][0]["out_of_bounds"] == []

    def test_negative_coordinates_are_caught(self, tmp_path):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Emu(-EMU_PER_INCH), Inches(1), Inches(3), Inches(1))
        box.text_frame.text = "off the left edge"
        assert deck_report(_save(prs, tmp_path / "neg.pptx"))["slides"][0]["out_of_bounds"]

    def test_offslide_text_is_caught_too(self, tmp_path):
        report = deck_report(build_offslide_deck(tmp_path / "off.pptx"))
        assert report["slides"][0]["out_of_bounds"]


def _save(prs, path: Path) -> Path:
    prs.save(str(path))
    return path


# ─────────────────────────────────────────────────────────────────────────
# 21. Editing — the original deck has to survive
# ─────────────────────────────────────────────────────────────────────────

def edit_deck(source: Path, dest: Path, new_title: str, extra_slides) -> Path:
    """Open, retitle slide 1, append slides, save to a NEW file. The exact
    open/modify/save round trip the pptx skill's edit path depends on."""
    prs = Presentation(str(source))
    for shape in prs.slides[0].shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            shape.text_frame.paragraphs[0].runs[0].text = new_title
            break
    for title, body in extra_slides:
        _add_content_slide(prs, title, body)
    prs.save(str(dest))
    return dest


NEW_SLIDES = [
    ("Agents", "A loop of model calls with tools attached, bounded by a step budget."),
    ("Cost Control", "Cache the prompt, cap the steps, and measure tokens per task."),
]


class TestEditingAnExistingDeck:
    @pytest.fixture
    def edited(self, tmp_path):
        source = build_clean_deck(tmp_path / "original.pptx")
        dest = edit_deck(source, tmp_path / "edited.pptx",
                         "How Machine Learning Actually Works", NEW_SLIDES)
        return source, dest

    def test_the_edited_file_still_opens(self, edited):
        _, dest = edited
        assert deck_report(dest)["slide_count"] == len(AI_SLIDES) + len(NEW_SLIDES)

    def test_the_title_actually_changed(self, edited):
        _, dest = edited
        first = deck_report(dest)["slides"][0]["text"]
        assert "How Machine Learning Actually Works" in first
        assert "What Machine Learning Actually Does" not in first

    def test_every_original_slide_keeps_its_text(self, edited):
        """The failure this guards: an "edit" implemented as a rebuild, which
        silently drops whatever the rebuild didn't know how to reproduce."""
        _, dest = edited
        slides = deck_report(dest)["slides"]
        for (title, body), slide in zip(AI_SLIDES[1:], slides[1:len(AI_SLIDES)]):
            assert title in slide["text"], f"original slide lost: {title}"
            assert body in slide["text"], f"original body lost under {title}"

    def test_the_new_slides_landed_at_the_end(self, edited):
        _, dest = edited
        slides = deck_report(dest)["slides"]
        for (title, body), slide in zip(NEW_SLIDES, slides[len(AI_SLIDES):]):
            assert title in slide["text"]
            assert body in slide["text"]

    def test_original_shape_geometry_is_preserved_exactly(self, edited):
        """"Text still present" is necessary but not sufficient — a rebuild
        can keep the words and lose every position. EMU-exact comparison."""
        source, dest = edited
        before = Presentation(str(source))
        after = Presentation(str(dest))
        for index in range(1, len(AI_SLIDES)):  # slide 1's title was edited
            rects_before = [_rect(s) for s in before.slides[index].shapes]
            rects_after = [_rect(s) for s in after.slides[index].shapes]
            assert rects_before == rects_after, f"slide {index + 1} shapes moved"

    def test_the_original_file_is_untouched(self, edited):
        """Editing must write a new file, not quietly mutate the input the
        user handed over."""
        source, _ = edited
        report = deck_report(source)
        assert report["slide_count"] == len(AI_SLIDES)
        assert "What Machine Learning Actually Does" in report["slides"][0]["text"]

    def test_the_edited_deck_is_still_presentable(self, edited):
        _, dest = edited
        assert_deck_is_presentable(deck_report(dest))

    def test_the_media_part_count_does_not_balloon(self, tmp_path):
        """Repeated open/save duplicating embedded media is a real OOXML
        bloat bug — the file still opens, it just doubles in size each pass."""
        img = make_image(tmp_path / "logo.png", 800, 600)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(str(img), Inches(1), Inches(1), width=Inches(3))
        current = _save(prs, tmp_path / "media0.pptx")

        import zipfile
        counts = []
        for step in range(3):
            with zipfile.ZipFile(current) as z:
                counts.append(len([n for n in z.namelist() if n.startswith("ppt/media/")]))
            current = edit_deck(current, tmp_path / f"media{step + 1}.pptx", "Pass", [])
        assert counts == [1, 1, 1], f"embedded media duplicated across saves: {counts}"

    def test_the_edit_survives_a_second_round_trip(self, tmp_path, edited):
        """Repeated open/save is where OOXML corruption usually shows up."""
        _, dest = edited
        again = edit_deck(dest, tmp_path / "twice.pptx", "Third Pass", [])
        report = deck_report(again)
        assert report["slide_count"] == len(AI_SLIDES) + len(NEW_SLIDES)
        assert "Third Pass" in report["slides"][0]["text"]
        assert NEW_SLIDES[0][0] in report["slides"][len(AI_SLIDES)]["text"]
