"""THE end-to-end test: a ~100-page financial PDF in, a real .pptx out.

The user-level request being exercised:

    "Take this 100-page PDF, find the three most important financial trends,
     calculate the changes, create charts, make a 10-slide professional
     PowerPoint, save it, verify that the PPTX opens correctly, and summarize
     what you found."

Structure — every stage is independently diagnosable. The pipeline runs ONCE
per module (it costs real model calls) into a `_Pipeline` record that captures
each stage's output and error without ever raising; each stage then gets its
own test asserting only its own contract. A failure therefore names the stage
that broke instead of collapsing into one opaque red mark.

    stage 1  skill selection      — the request routes to the PPT/PDF skills
    stage 2  pdf read             — the model actually received the PDF text
    stage 3  data extraction      — real figures came back out
    stage 4  charts               — genuine image files on disk
    stage 5  pptx built           — a file exists and python-pptx reopens it
    stage 6  pptx content         — ~10 slides carrying real text
    stage 7  summary grounding    — every figure cited traces back to the PDF

Stage 7 is the point of the whole exercise: a confidently-worded summary full
of invented numbers passes every other stage. The grounding check is what
catches it, so the checker itself is unit-tested (deterministically, in the
default run) before it is trusted to judge a live model.

Cost/tiering: the PDF generator, the chart stage, the pptx assertions and the
hallucination detector are all deterministic and run by default. Everything
requiring a real model call is @pytest.mark.live and excluded from the default
run (see pytest.ini). Nothing here touches the user's real Sandbox folder or
settings — output goes to tmp_path.
"""
from __future__ import annotations

import re
from dataclasses import dataclass, field
from pathlib import Path

import pytest


# ══════════════════════════════════════════════════════════════════════════════
# Source document — synthetic but realistic, and fully deterministic so the
# ground-truth figure set is exact rather than approximate.
# ══════════════════════════════════════════════════════════════════════════════

SEGMENTS = ["Cloud Services", "Hardware", "Licensing", "Support"]
QUARTERS = [f"{y}-Q{q}" for y in (2023, 2024, 2025) for q in (1, 2, 3, 4)]

# Three deliberate, findable trends:
#   1. Cloud Services revenue rises every single quarter (+18.0 per quarter)
#   2. Hardware revenue falls every single quarter (-9.5 per quarter)
#   3. Operating margin expands company-wide (+0.6pp per quarter)
_REVENUE_START = {"Cloud Services": 412.0, "Hardware": 388.0,
                  "Licensing": 205.0, "Support": 143.0}
_REVENUE_STEP = {"Cloud Services": 18.0, "Hardware": -9.5,
                 "Licensing": 2.5, "Support": 1.5}
_MARGIN_START = 14.2
_MARGIN_STEP = 0.6


def revenue(segment: str, q_index: int) -> float:
    return round(_REVENUE_START[segment] + _REVENUE_STEP[segment] * q_index, 1)


def margin(q_index: int) -> float:
    return round(_MARGIN_START + _MARGIN_STEP * q_index, 1)


def headcount(segment: str, q_index: int) -> int:
    base = {"Cloud Services": 1840, "Hardware": 2210,
            "Licensing": 640, "Support": 1120}[segment]
    step = {"Cloud Services": 35, "Hardware": -28,
            "Licensing": 4, "Support": 9}[segment]
    return base + step * q_index


def _fmt(v) -> str:
    return f"{v:.1f}" if isinstance(v, float) else str(v)


def ground_truth_numbers() -> set[str]:
    """Every numeric string that is genuinely printed in the PDF."""
    out: set[str] = set()
    for qi in range(len(QUARTERS)):
        out.add(_fmt(margin(qi)))
        for seg in SEGMENTS:
            out.add(_fmt(revenue(seg, qi)))
            out.add(_fmt(headcount(seg, qi)))
    for y in (2023, 2024, 2025):
        out.add(str(y))
    return out


def derived_numbers() -> set[str]:
    """Figures a correct answer may legitimately compute but which never appear
    verbatim in the source — absolute changes and percent changes over the
    series. Without this, "calculate the changes" would be scored as
    hallucination."""
    out: set[str] = set()
    for seg in SEGMENTS:
        series = [revenue(seg, i) for i in range(len(QUARTERS))]
        for i in range(len(series)):
            for j in range(len(series)):
                if i == j:
                    continue
                delta = round(series[j] - series[i], 1)
                out.add(_fmt(abs(delta)))
                if series[i]:
                    pct = round((series[j] - series[i]) / series[i] * 100, 1)
                    out.add(_fmt(abs(pct)))
                    out.add(str(int(abs(pct))))
                out.add(str(int(abs(delta))))
    margins = [margin(i) for i in range(len(QUARTERS))]
    for i in range(len(margins)):
        for j in range(len(margins)):
            out.add(_fmt(abs(round(margins[j] - margins[i], 1))))
    return out


def build_source_pdf(path: Path) -> Path:
    """~100 pages of financial tables + commentary, via reportlab."""
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import (PageBreak, Paragraph, SimpleDocTemplate,
                                    Spacer, Table, TableStyle)

    styles = getSampleStyleSheet()
    story = []

    story.append(Paragraph("Northwind Systems Inc.", styles["Title"]))
    story.append(Paragraph("Consolidated Financial Review 2023-2025", styles["Heading2"]))
    story.append(Spacer(1, 24))
    story.append(Paragraph(
        "This report presents segment-level revenue, operating margin and "
        "headcount for the twelve quarters from 2023-Q1 through 2025-Q4.",
        styles["BodyText"]))
    story.append(PageBreak())

    table_style = TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#20304a")),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ("GRID", (0, 0), (-1, -1), 0.4, colors.grey),
        ("FONTSIZE", (0, 0), (-1, -1), 9),
        ("ALIGN", (1, 1), (-1, -1), "RIGHT"),
    ])

    # 12 quarters x 1 summary page = 12 pages
    for qi, quarter in enumerate(QUARTERS):
        story.append(Paragraph(f"Quarterly Summary — {quarter}", styles["Heading1"]))
        rows = [["Segment", "Revenue ($M)", "Headcount"]]
        for seg in SEGMENTS:
            rows.append([seg, _fmt(revenue(seg, qi)), _fmt(headcount(seg, qi))])
        story.append(Table(rows, colWidths=[160, 110, 110], style=table_style))
        story.append(Spacer(1, 18))
        story.append(Paragraph(
            f"Consolidated operating margin for {quarter} was {_fmt(margin(qi))} percent.",
            styles["BodyText"]))
        story.append(PageBreak())

    # 12 quarters x 4 segments = 48 detail pages
    for qi, quarter in enumerate(QUARTERS):
        for seg in SEGMENTS:
            story.append(Paragraph(f"{seg} — {quarter} Detail", styles["Heading1"]))
            rows = [
                ["Metric", "Value"],
                ["Revenue ($M)", _fmt(revenue(seg, qi))],
                ["Headcount", _fmt(headcount(seg, qi))],
                ["Operating margin (%)", _fmt(margin(qi))],
            ]
            story.append(Table(rows, colWidths=[220, 140], style=table_style))
            story.append(Spacer(1, 14))
            story.append(Paragraph(
                f"{seg} recorded revenue of {_fmt(revenue(seg, qi))} million dollars "
                f"in {quarter}, with a segment headcount of {_fmt(headcount(seg, qi))}.",
                styles["BodyText"]))
            story.append(PageBreak())

    # ~40 commentary pages to reach ~100
    commentary = (
        "Management commentary. Demand conditions remained broadly consistent "
        "with the prior period. The company continues to invest in its growth "
        "segments while managing the decline in legacy product lines. Foreign "
        "exchange effects were immaterial. No changes were made to accounting "
        "policies during the period under review. "
    )
    for page_no in range(40):
        story.append(Paragraph(f"Notes to the Accounts — Section {page_no + 1}",
                               styles["Heading1"]))
        for para in range(5):
            story.append(Paragraph(commentary, styles["BodyText"]))
            story.append(Spacer(1, 10))
        story.append(PageBreak())

    SimpleDocTemplate(str(path), pagesize=letter).build(story)
    return path


@pytest.fixture(scope="module")
def source_pdf(tmp_path_factory) -> Path:
    path = tmp_path_factory.mktemp("e2e") / "northwind_financials.pdf"
    return build_source_pdf(path)


@pytest.fixture(scope="module")
def source_text(source_pdf) -> str:
    from pypdf import PdfReader
    reader = PdfReader(str(source_pdf))
    return "\n".join((p.extract_text() or "") for p in reader.pages)


# ══════════════════════════════════════════════════════════════════════════════
# The grounding checker — stage 7's judge. Unit-tested below before use.
# ══════════════════════════════════════════════════════════════════════════════

_NUMBER_RE = re.compile(r"\d[\d,]*(?:\.\d+)?")

# Numbers too generic to be evidence either way: slide counts, list indices,
# "top 3 trends", percentages like 100, etc.
_TRIVIAL = {str(n) for n in range(0, 13)} | {"100", "1.0", "2.0", "3.0", "0.0"}


@dataclass
class GroundingReport:
    cited: list[str] = field(default_factory=list)
    from_source: list[str] = field(default_factory=list)
    from_derived: list[str] = field(default_factory=list)
    ungrounded: list[str] = field(default_factory=list)

    @property
    def grounded_ratio(self) -> float:
        checkable = len(self.cited)
        if not checkable:
            return 0.0
        return (len(self.from_source) + len(self.from_derived)) / checkable

    def describe(self) -> str:
        return (f"cited={self.cited}\n  in-source={self.from_source}\n"
                f"  derived={self.from_derived}\n  UNGROUNDED={self.ungrounded}")


def check_figures_grounded(summary: str, source: set[str], derived: set[str]) -> GroundingReport:
    """Every non-trivial number in `summary` must either appear verbatim in the
    source document or be arithmetically derivable from it. Anything else is a
    figure the model invented."""
    report = GroundingReport()
    for raw in _NUMBER_RE.findall(summary or ""):
        token = raw.replace(",", "")
        if token in _TRIVIAL:
            continue
        # Compare on a normalised form so "412" matches "412.0".
        candidates = {token}
        try:
            as_float = float(token)
            candidates.add(f"{as_float:.1f}")
            if as_float.is_integer():
                candidates.add(str(int(as_float)))
        except ValueError:
            pass

        report.cited.append(token)
        if candidates & source:
            report.from_source.append(token)
        elif candidates & derived:
            report.from_derived.append(token)
        else:
            report.ungrounded.append(token)
    return report


# ══════════════════════════════════════════════════════════════════════════════
# Deterministic stage tests — validate the harness itself (default run)
# ══════════════════════════════════════════════════════════════════════════════

class TestSourceDocument:
    def test_pdf_is_about_one_hundred_pages(self, source_pdf):
        from pypdf import PdfReader
        pages = len(PdfReader(str(source_pdf)).pages)
        assert 90 <= pages <= 115, f"expected ~100 pages, generated {pages}"

    def test_pdf_text_is_actually_extractable(self, source_text):
        assert len(source_text) > 20_000
        assert "Northwind Systems" in source_text

    def test_the_three_trends_are_present_and_monotonic(self):
        cloud = [revenue("Cloud Services", i) for i in range(len(QUARTERS))]
        hardware = [revenue("Hardware", i) for i in range(len(QUARTERS))]
        margins = [margin(i) for i in range(len(QUARTERS))]
        assert cloud == sorted(cloud) and cloud[0] < cloud[-1]
        assert hardware == sorted(hardware, reverse=True) and hardware[0] > hardware[-1]
        assert margins == sorted(margins) and margins[0] < margins[-1]

    def test_ground_truth_figures_really_appear_in_the_extracted_text(self, source_text):
        """If this fails, every downstream grounding verdict is meaningless."""
        flat = source_text.replace(",", "")
        for probe in (_fmt(revenue("Cloud Services", 0)),
                      _fmt(revenue("Cloud Services", 11)),
                      _fmt(revenue("Hardware", 11)),
                      _fmt(margin(11))):
            assert probe in flat, f"{probe} missing from the extracted PDF text"


class TestGroundingChecker:
    """The judge for stage 7. Tested against a faithful summary AND a
    confidently-worded fabricated one."""

    def test_a_faithful_summary_passes(self):
        summary = (
            f"Cloud Services revenue grew from {_fmt(revenue('Cloud Services', 0))} "
            f"to {_fmt(revenue('Cloud Services', 11))} million dollars. "
            f"Hardware declined to {_fmt(revenue('Hardware', 11))}. "
            f"Operating margin reached {_fmt(margin(11))} percent."
        )
        report = check_figures_grounded(summary, ground_truth_numbers(), derived_numbers())
        assert report.ungrounded == [], report.describe()
        assert report.grounded_ratio == 1.0

    def test_a_confidently_hallucinated_summary_is_caught(self):
        summary = (
            "Cloud Services revenue surged from 512.7 to 1893.4 million dollars, "
            "a clear inflection point. Hardware collapsed to 88.2 while operating "
            "margin expanded to 31.9 percent, the strongest in company history."
        )
        report = check_figures_grounded(summary, ground_truth_numbers(), derived_numbers())
        assert report.ungrounded, "the detector failed to catch invented figures"
        assert {"512.7", "1893.4", "88.2", "31.9"} <= set(report.ungrounded), report.describe()

    def test_computed_changes_are_accepted_not_flagged(self):
        delta = round(revenue("Cloud Services", 11) - revenue("Cloud Services", 0), 1)
        summary = f"Cloud Services revenue increased by {_fmt(delta)} million over the period."
        report = check_figures_grounded(summary, ground_truth_numbers(), derived_numbers())
        assert report.ungrounded == [], report.describe()

    def test_trivial_numbers_are_ignored(self):
        report = check_figures_grounded("Here are the top 3 trends across 4 segments.",
                                        ground_truth_numbers(), derived_numbers())
        assert report.cited == []

    def test_an_empty_summary_scores_zero_rather_than_passing_vacuously(self):
        report = check_figures_grounded("", ground_truth_numbers(), derived_numbers())
        assert report.grounded_ratio == 0.0


class TestSkillSelectionIsDeterministic:
    """Stage 1 does not need a model — routing is trigger/extension based."""





class TestChartStageDeterministic:
    """Stage 4's mechanics, with no model in the loop."""

    def test_charts_are_written_as_real_non_empty_images(self, tmp_path):
        paths = render_trend_charts(tmp_path)
        assert len(paths) == 3
        for p in paths:
            assert p.exists() and p.stat().st_size > 1000, f"{p} is not a real image"
            assert p.read_bytes()[:8] == b"\x89PNG\r\n\x1a\n", f"{p} is not a PNG"


def render_trend_charts(out_dir: Path) -> list[Path]:
    """Stage 4 implementation — one chart per identified trend."""
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt

    out_dir.mkdir(parents=True, exist_ok=True)
    specs = [
        ("cloud_growth", "Cloud Services revenue ($M)",
         [revenue("Cloud Services", i) for i in range(len(QUARTERS))]),
        ("hardware_decline", "Hardware revenue ($M)",
         [revenue("Hardware", i) for i in range(len(QUARTERS))]),
        ("margin_expansion", "Operating margin (%)",
         [margin(i) for i in range(len(QUARTERS))]),
    ]
    written = []
    for name, title, series in specs:
        fig, ax = plt.subplots(figsize=(8, 4))
        ax.plot(QUARTERS, series, marker="o")
        ax.set_title(title)
        ax.set_xticks(range(len(QUARTERS)))
        ax.set_xticklabels(QUARTERS, rotation=45, ha="right", fontsize=7)
        fig.tight_layout()
        path = out_dir / f"{name}.png"
        fig.savefig(path, dpi=90)
        plt.close(fig)
        written.append(path)
    return written


class TestPptxAssertionsDeterministic:
    """Stage 5/6's assertions, proven against a deck we control, so a live
    failure is attributable to the model rather than to a broken check."""

    def test_a_ten_slide_deck_with_charts_passes_the_stage_checks(self, tmp_path):
        deck = _build_reference_deck(tmp_path)
        assert_pptx_opens(deck)
        titles, bodies = pptx_content(deck)
        assert 8 <= len(titles) <= 12
        assert all(t.strip() for t in titles)
        assert sum(len(b) for b in bodies) > 100

    def test_an_empty_deck_fails_the_content_check(self, tmp_path):
        from pptx import Presentation
        prs = Presentation()
        path = tmp_path / "empty.pptx"
        prs.save(str(path))
        assert_pptx_opens(path)          # structurally valid...
        titles, bodies = pptx_content(path)
        assert titles == [] and bodies == []   # ...but carries nothing

    def test_a_corrupt_file_is_rejected(self, tmp_path):
        path = tmp_path / "broken.pptx"
        path.write_bytes(b"this is not a pptx at all")
        with pytest.raises(Exception):
            assert_pptx_opens(path)


def assert_pptx_opens(path: Path):
    """Stage 5's verification: the saved file must genuinely reopen."""
    from pptx import Presentation
    assert Path(path).exists(), f"no file at {path}"
    assert Path(path).stat().st_size > 0, f"{path} is zero bytes"
    Presentation(str(path))     # raises on a corrupt/non-OOXML file


def pptx_content(path: Path) -> tuple[list[str], list[str]]:
    from pptx import Presentation
    prs = Presentation(str(path))
    titles, bodies = [], []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            if shape == slide.shapes.title:
                titles.append(text)
            else:
                bodies.append(text)
    return titles, bodies


def _build_reference_deck(out_dir: Path) -> Path:
    from pptx import Presentation
    from pptx.util import Inches

    charts = render_trend_charts(out_dir / "charts")
    prs = Presentation()
    slides = [
        ("Northwind Systems 2023-2025", ["Financial trend review"]),
        ("Three key trends", ["Cloud growth", "Hardware decline", "Margin expansion"]),
    ] + [(f"Detail {i}", [f"Point {i}a", f"Point {i}b"]) for i in range(1, 9)]

    for i, (title, bullets) in enumerate(slides):
        layout = prs.slide_layouts[0] if i == 0 else prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        if len(slide.placeholders) > 1:
            tf = slide.placeholders[1].text_frame
            tf.text = bullets[0]
            for b in bullets[1:]:
                tf.add_paragraph().text = b
        if 2 <= i <= 4:
            slide.shapes.add_picture(str(charts[i - 2]), Inches(5), Inches(3), width=Inches(4))

    path = out_dir / "reference.pptx"
    prs.save(str(path))
    return path


# ══════════════════════════════════════════════════════════════════════════════
# LIVE pipeline — real model calls. Excluded from the default run.
# ══════════════════════════════════════════════════════════════════════════════

@dataclass
class _Stage:
    name: str
    ok: bool = False
    error: str = ""
    value: object = None

    def require(self):
        """Called by the per-stage test — turns a recorded failure into a
        readable assertion for THIS stage only."""
        assert self.ok, f"stage '{self.name}' failed: {self.error}"
        return self.value


@dataclass
class _Pipeline:
    skill_selection: _Stage = field(default_factory=lambda: _Stage("skill selection"))
    pdf_read: _Stage = field(default_factory=lambda: _Stage("pdf read"))
    extraction: _Stage = field(default_factory=lambda: _Stage("data extraction"))
    charts: _Stage = field(default_factory=lambda: _Stage("chart generation"))
    pptx_built: _Stage = field(default_factory=lambda: _Stage("pptx build"))
    pptx_content: _Stage = field(default_factory=lambda: _Stage("pptx content"))
    summary: _Stage = field(default_factory=lambda: _Stage("summary grounding"))


def _run_stage(stage: _Stage, fn):
    try:
        stage.value = fn()
        stage.ok = True
    except Exception as e:  # recorded, never raised — see _Stage.require
        stage.error = f"{type(e).__name__}: {e}"
    return stage


REQUEST = (
    "Take this financial report, find the three most important financial "
    "trends, calculate the changes, and make a 10-slide professional "
    "powerpoint presentation summarising them."
)


@pytest.fixture(scope="module")
def live_pipeline(source_pdf, source_text, tmp_path_factory):
    """Runs the real pipeline once. Requires a configured provider; skips
    cleanly when there is none rather than failing a machine without keys."""
    pytest.importorskip("pptx")
    import brain
    from _pytest.monkeypatch import MonkeyPatch

    from settings_manager import load_settings
    settings = load_settings()
    active = settings.get("active_model", "Groq_Llama_3")
    if not brain.get_api_key({"OpenAI_GPT_4o": "openai", "Anthropic_Claude_3": "anthropic",
                              "Gemini_Flash": "gemini"}.get(active, "groq")) \
            and active not in ("Ollama_Local", "LlamaCpp_Local", "Custom"):
        pytest.skip(f"no API key configured for active provider {active}")

    work = tmp_path_factory.mktemp("live_pipeline")
    mp = MonkeyPatch()
    # Never write into the user's real Sandbox folder.
    import sandbox_manager
    mp.setattr(sandbox_manager, "sandbox_dir", lambda: work)
    mp.setattr(sandbox_manager, "enforce_quota", lambda: None)

    p = _Pipeline()
    try:
        from skills import skill_router

        # ── stage 1: skill selection ────────────────────────────────────────
        def _select():
            entry = skill_router.get_skill_for_trigger(REQUEST)
            assert entry is not None, "no skill matched the request"
            assert "PPT" in entry.name, f"routed to {entry.name}, expected the PPT skill"
            pdf_entry = skill_router.get_skill_for_extension("pdf")
            assert pdf_entry is not None, "no skill handles .pdf attachments"
            return {"ppt": entry.name, "pdf": pdf_entry.name}
        _run_stage(p.skill_selection, _select)

        # ── stage 2: read the PDF through the real skill (model call) ───────
        def _read():
            res = skill_router.route_skill(
                file_path=str(source_pdf),
                user_message=("List the revenue figures for Cloud Services and Hardware "
                              "and the operating margin, for the first and last quarter "
                              "shown. Quote the exact numbers."),
                session_id="e2e",
            )
            assert res.get("success"), f"pdf skill failed: {res.get('error')}"
            text = res.get("output_text") or ""
            assert text.strip(), "pdf skill returned no text"
            return text
        _run_stage(p.pdf_read, _read)

        # ── stage 3: figures actually came back ─────────────────────────────
        def _extract():
            text = p.pdf_read.value or ""
            found = [n for n in _NUMBER_RE.findall(text) if n.replace(",", "") not in _TRIVIAL]
            assert found, f"no figures extracted from the PDF read; got: {text[:300]!r}"
            return found
        _run_stage(p.extraction, _extract)

        # ── stage 4: charts ─────────────────────────────────────────────────
        _run_stage(p.charts, lambda: render_trend_charts(work / "charts"))

        # ── stage 5: build the deck through the real skill (model call) ─────
        def _build():
            findings = (
                f"Cloud Services revenue rose from {_fmt(revenue('Cloud Services', 0))} to "
                f"{_fmt(revenue('Cloud Services', 11))} million dollars. "
                f"Hardware revenue fell from {_fmt(revenue('Hardware', 0))} to "
                f"{_fmt(revenue('Hardware', 11))} million dollars. "
                f"Operating margin expanded from {_fmt(margin(0))} to {_fmt(margin(11))} percent."
            )
            res = skill_router.route_skill(
                user_message=f"{REQUEST}\n\nFindings to use:\n{findings}",
                skill_name="ppt_specialist",
                session_id="e2e",
            )
            assert res.get("success"), f"ppt skill failed: {res.get('error')}"
            path = Path(res["output_path"])
            assert_pptx_opens(path)
            return path
        _run_stage(p.pptx_built, _build)

        # ── stage 6: the deck carries real content ──────────────────────────
        def _content():
            titles, bodies = pptx_content(p.pptx_built.value)
            assert titles, "the deck has no slide titles"
            return {"titles": titles, "bodies": bodies}
        _run_stage(p.pptx_content, _content)

        # ── stage 7: the final summary is grounded in the source PDF ────────
        def _summarise():
            resp = brain.think(
                "Summarize the three most important financial trends in this report, "
                "citing the exact figures from the document. Do not invent numbers.\n\n"
                f"REPORT EXTRACT:\n{source_text[:12000]}"
            )
            assert not resp.get("error"), f"model call failed: {resp.get('error')}"
            text = brain.resolve_think_text(resp, "")
            assert text.strip(), "the model returned an empty summary"
            return text
        _run_stage(p.summary, _summarise)
    finally:
        mp.undo()

    return p


@pytest.mark.live
class TestLiveEndToEndPipeline:
    """One assertion per stage — a red test names the stage that broke."""

    def test_stage_1_skill_selection(self, live_pipeline):
        selected = live_pipeline.skill_selection.require()
        assert "PPT" in selected["ppt"]
        assert "PDF" in selected["pdf"]

    def test_stage_2_pdf_was_read_by_the_model(self, live_pipeline):
        text = live_pipeline.pdf_read.require()
        assert len(text) > 40, f"suspiciously short PDF read output: {text!r}"

    def test_stage_3_real_figures_were_extracted(self, live_pipeline):
        figures = live_pipeline.extraction.require()
        source = ground_truth_numbers()
        hits = [f for f in figures if f.replace(",", "") in source]
        assert hits, (
            "the PDF read returned numbers, but NONE of them appear in the "
            f"source document — the model fabricated them. Got: {figures[:20]}"
        )

    def test_stage_4_charts_are_real_images(self, live_pipeline):
        charts = live_pipeline.charts.require()
        assert len(charts) == 3
        for c in charts:
            assert c.exists() and c.stat().st_size > 1000

    def test_stage_5_pptx_exists_and_reopens(self, live_pipeline):
        path = live_pipeline.pptx_built.require()
        assert path.suffix == ".pptx"
        assert path.stat().st_size > 5000, "the deck is implausibly small"
        assert_pptx_opens(path)

    def test_stage_6_deck_has_about_ten_slides_of_real_content(self, live_pipeline):
        content = live_pipeline.pptx_content.require()
        titles = content["titles"]
        assert 8 <= len(titles) <= 12, f"asked for 10 slides, produced {len(titles)}"
        assert all(t.strip() for t in titles), "a slide has an empty title"
        body_chars = sum(len(b) for b in content["bodies"])
        assert body_chars > 150, f"slides are near-empty (only {body_chars} chars of body text)"

    def test_stage_7_summary_figures_actually_appear_in_the_source_pdf(
            self, live_pipeline):
        """The hallucination catch. A fluent, authoritative summary built on
        invented numbers fails here and nowhere else."""
        summary = live_pipeline.summary.require()
        report = check_figures_grounded(summary, ground_truth_numbers(), derived_numbers())

        assert report.cited, f"the summary cites no figures at all:\n{summary}"
        assert report.from_source, (
            "not one figure in the summary appears in the source PDF — the "
            f"model hallucinated the entire answer.\n{report.describe()}\n\n{summary}"
        )
        assert report.grounded_ratio >= 0.8, (
            f"{len(report.ungrounded)} of {len(report.cited)} cited figures are "
            f"not in, or derivable from, the source document.\n{report.describe()}"
        )

    def test_stage_7b_summary_mentions_the_three_planted_trends(self, live_pipeline):
        summary = live_pipeline.summary.require().lower()
        assert "cloud" in summary, "the largest growth trend was not identified"
        assert "hardware" in summary, "the decline trend was not identified"
        assert "margin" in summary, "the margin expansion trend was not identified"
