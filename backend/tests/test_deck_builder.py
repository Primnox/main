"""The deck builder, graded by the auditor rather than by opinion.

Every layout is exercised and the whole deck is then audited against its own
design system. That loop is the point: "proper layouts and a grid" is a claim
that can be checked, and checking it is cheaper than looking at slides.
"""
from __future__ import annotations

import pytest

pytest.importorskip("pptx")

from deck_audit import DesignSystem, audit, slide_count, summarise   # noqa: E402

from primnox2.sandbox import doc_themes                              # noqa: E402


def _system(theme_name: str = "midnight") -> DesignSystem:
    """The deck's own theme, expressed as a contract to grade it against.

    `surfaces()` is applied here too: the panel and hairline tints are derived
    tokens, as much part of the system as the five declared colours. Grading
    against the raw theme would report every filled card as colour drift.
    """
    t = doc_themes.surfaces(dict(doc_themes.theme(theme_name)))
    return DesignSystem(
        palette={t[k].upper() for k in
                 ("bg", "text", "muted", "primary", "accent", "surface", "line")},
        font_sizes=set(doc_themes.SCALE.values()),
        fonts={doc_themes.DISPLAY_FONT, doc_themes.BODY_FONT, doc_themes.MONO_FONT},
        grid_pt=doc_themes.GRID,
        min_readable_pt=doc_themes.MIN_READABLE_PT,
    )


def _full_deck(path: str) -> str:
    d = doc_themes.Deck(path, theme="midnight", footer="Primnox")
    d.hero("Every Layout", "One call each")
    d.section("Part One", "The basics", number="01")
    d.bullets("Bullets", ["First point", "Second point that runs a good deal longer",
                          "Third"])
    d.two_column("Two Column", ["Alpha", "Beta"], ["Gamma", "Delta"],
                 left_title="Left", right_title="Right")
    d.compare("Compare", "Before", ["Slow", "Manual"], "After", ["Fast", "Automatic"])
    d.bento("Bento", [("Users", "12,400", "+8%"), ("Revenue", "$4.2M", "+12%"),
                      ("Churn", "1.9%", "-0.4pt"), ("NPS", "62", "+5")])
    d.kpi("KPI Dashboard", [("Uptime", "99.98%", "30d"), ("p95", "180ms", "-12ms"),
                            ("Errors", "0.02%", ""), ("Deploys", "48", "+9"),
                            ("MTTR", "6m", "-2m"), ("Cost", "$1.2k", "-8%")])
    d.timeline("Timeline", [("2021", "Founded"), ("2023", "Series A"),
                            ("2025", "Profitable")])
    d.process("Process", ["Collect", "Extract", "Index", "Query"])
    d.table("Table", [["Region", "Q1", "Q2"], ["EMEA", "1.2", "1.6"],
                      ["APAC", "0.9", "1.1"], ["AMER", "2.1", "2.4"]])
    d.chart("Chart", ["Q1", "Q2", "Q3"], {"Revenue": (1.2, 1.6, 2.0)}, kind="bar")
    d.code("Code", "def hello():\n    return 1", caption="hello.py")
    d.quote("Simplicity is the ultimate sophistication.", "Attributed")
    d.matrix("Matrix", [("Strengths", ["Fast"]), ("Weaknesses", ["New"]),
                        ("Opportunities", ["Growth"]), ("Threats", ["Rivals"])])
    d.appendix("Appendix", [("Nodes", "2,459"), ("Edges", "3,207"),
                            ("Communities", "116")])
    return d.save()


# ── the claim, graded ────────────────────────────────────────────────────────
def test_every_layout_builds(tmp_path):
    path = _full_deck(str(tmp_path / "all.pptx"))
    assert slide_count(path) == 15, "one slide per layout"


def test_the_full_deck_passes_its_own_design_system(tmp_path):
    path = _full_deck(str(tmp_path / "all.pptx"))
    findings = audit(path, _system())
    errors = [f for f in findings if f.severity == "error"]
    assert not errors, "\n".join(str(f) for f in errors)


def test_the_full_deck_is_on_the_grid(tmp_path):
    """The defect that started this: the old builder worked in inches, so every
    coordinate landed off an 8-point grid and no design system could be met."""
    path = _full_deck(str(tmp_path / "all.pptx"))
    misaligned = [f for f in audit(path, _system()) if f.check == "grid-misalign"]
    assert not misaligned, "\n".join(str(f) for f in misaligned[:10])


def test_at_least_fifteen_layouts_are_advertised():
    assert len(doc_themes.Deck.LAYOUTS) >= 15
    for name in doc_themes.Deck.LAYOUTS:
        assert callable(getattr(doc_themes.Deck, name, None)), name


# ── palette injection ────────────────────────────────────────────────────────
def test_a_caller_supplied_palette_is_used(tmp_path):
    """The torture-test brief names Navy/Cyan explicitly. The old builder could
    only ever produce one of ten built-in themes, so that brief was unreachable."""
    palette = {"bg": "#0B1220", "text": "#F8FAFC", "muted": "#94A3B8",
               "primary": "#22D3EE", "accent": "#A3E635"}
    path = str(tmp_path / "branded.pptx")
    d = doc_themes.Deck(path, palette=palette, footer="Acme")
    d.hero("Branded", "With a caller's own colours")
    d.kpi("Numbers", [("Revenue", "$4.2M", "+12%")])
    d.save()

    derived = doc_themes.surfaces({k: v.lstrip("#").upper() for k, v in palette.items()})
    system = DesignSystem(
        palette=set(derived.values()),
        font_sizes=set(doc_themes.SCALE.values()),
        fonts={doc_themes.DISPLAY_FONT, doc_themes.BODY_FONT, doc_themes.MONO_FONT},
        grid_pt=doc_themes.GRID,
    )
    drift = [f for f in audit(path, system) if f.check == "color-drift"]
    assert not drift, "\n".join(str(f) for f in drift[:5])


def test_a_partial_palette_only_overrides_what_it_names(tmp_path):
    path = str(tmp_path / "partial.pptx")
    doc_themes.Deck(path, theme="midnight", palette={"accent": "FF0000"},
                    title="Partial").save()
    midnight = doc_themes.theme("midnight")
    system = DesignSystem(
        palette={midnight[k].upper() for k in ("bg", "text", "muted", "primary")}
                | {"FF0000"},
        font_sizes=set(doc_themes.SCALE.values()),
        fonts={doc_themes.DISPLAY_FONT, doc_themes.BODY_FONT, doc_themes.MONO_FONT},
    )
    assert not [f for f in audit(path, system) if f.check == "color-drift"]


# ── behaviour the spec demands ───────────────────────────────────────────────
def test_a_long_table_paginates_instead_of_shrinking(tmp_path):
    """The spec asks for 15+ row comparisons AND forbids shrinking below
    readability. Those can only both hold if the deck gains a slide."""
    rows = [["Region", "Q1", "Q2"]] + [[f"Row {i}", i, i * 2] for i in range(30)]
    path = str(tmp_path / "big.pptx")
    d = doc_themes.Deck(path, theme="midnight")
    d.table("Large", rows)
    d.save()

    assert slide_count(path) > 1, "30 rows were crammed onto one slide"
    below = [f for f in audit(path, _system()) if f.check == "below-readable"]
    assert not below


def test_text_boxes_do_not_silently_grow(tmp_path):
    """python-pptx defaults textboxes to SHAPE_TO_FIT_TEXT, which expands a box
    past the geometry it was placed on and breaks the grid at render time."""
    from pptx import Presentation

    path = _full_deck(str(tmp_path / "all.pptx"))
    grew = []
    for i, slide in enumerate(Presentation(path).slides, start=1):
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                mode = shape.text_frame.auto_size
                if mode is not None and getattr(mode, "value", None) == 1:
                    grew.append((i, shape.name))
    assert not grew, f"boxes set to grow: {grew[:5]}"


def test_headline_shrinks_to_fit_but_never_below_the_floor():
    long_title = "A deliberately enormous headline " * 6
    size = doc_themes._fit_size(long_title, doc_themes.CONTENT_W, 96, start="h1")
    assert size < doc_themes.SCALE["h1"], "did not shrink at all"
    assert size >= doc_themes.MIN_READABLE_PT


def test_snap_puts_everything_on_the_grid():
    for raw in (0, 3, 4, 5, 63, 64.8, 172.8, 828):
        assert doc_themes.snap(raw) % doc_themes.GRID == 0


# ── quality gates that are not about correctness ─────────────────────────────
@pytest.mark.parametrize("name", sorted(doc_themes.THEMES))
def test_every_theme_meets_contrast(name):
    """WCAG AA on every shipped palette.

    Found two failures when this was first run: `sand` at 3.79:1 and `mono` at
    4.51:1 for muted-on-background. Captions and footers use `muted`, so both
    themes had unreadable small text on a projector and nothing said so.
    """
    t = doc_themes.THEMES[name]
    for role, floor in (("text", 4.5), ("muted", 4.5),
                        ("primary", 3.0), ("accent", 3.0)):
        ratio = doc_themes.contrast_ratio(t[role], t["bg"])
        assert ratio >= floor, f"{name}.{role} is {ratio:.2f}:1, needs {floor}:1"


def test_derived_surfaces_stay_distinguishable_from_the_background():
    """A panel tint so faint it is invisible is worse than no panel: the card
    stops reading as an object and the layout looks like floating text."""
    for name, t in doc_themes.THEMES.items():
        full = doc_themes.surfaces(dict(t))
        assert full["surface"] != t["bg"], name
        assert doc_themes.contrast_ratio(full["text"], full["surface"]) >= 4.5, name


def test_leading_loosens_as_type_gets_smaller():
    """Large type needs proportionally less leading; small type needs more.
    A single flat multiplier makes headlines gappy or captions cramped."""
    sizes = [doc_themes.SCALE[k] for k in ("h1", "h3", "body", "caption")]
    ratios = [doc_themes.leading(s) for s in sizes]
    assert ratios == sorted(ratios), f"leading is not monotonic: {ratios}"
    assert ratios[0] < ratios[-1]


@pytest.mark.parametrize("n, most, expected_gaps", [
    (6, 4, 0),   # 3x2, not 4+2 with two holes
    (4, 4, 0),
    (8, 4, 0),
    (5, 4, 1),   # unavoidable
    (3, 3, 0),
])
def test_columns_are_chosen_to_leave_the_fewest_gaps(n, most, expected_gaps):
    cols = doc_themes._balanced_columns(n, most)
    rows = -(-n // cols)
    assert cols * rows - n == expected_gaps, f"{n} in {cols} cols"


def test_a_full_deck_has_no_monotonous_run(tmp_path):
    """The builder offers fifteen layouts; a deck using them should not read as
    a wall of one."""
    path = _full_deck(str(tmp_path / "all.pptx"))
    assert not [f for f in audit(path, _system()) if f.check == "monotonous-run"]


# ── the old API still works ──────────────────────────────────────────────────
def test_the_previous_api_still_builds(tmp_path):
    """`cover` and `slide` are what the themed-documents skill teaches and what
    already-generated scripts call."""
    path = str(tmp_path / "legacy.pptx")
    d = doc_themes.Deck(path, theme="midnight", title="Cover", subtitle="Sub")
    d.slide("A heading", ["first", "second"])
    d.save()
    assert slide_count(path) == 2
    assert not [f for f in audit(path, _system()) if f.severity == "error"]
