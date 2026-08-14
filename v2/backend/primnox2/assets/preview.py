"""Read-only previews of stored assets — one shape per family of format.

Why this is a server concern. The browser renders PDFs and images natively and
needs nothing from us; Word, Excel and PowerPoint it cannot read at all. The
libraries that *can* read them are already installed here — the sandbox writes
these formats with the very same ones — so the alternative is shipping a
second, JavaScript copy of each parser to the client. This is the cheaper half
of the job by a wide margin.

Nothing here can modify an asset. Every function reads bytes and returns a
description of them; there is no write path to reach even by accident.

Every parser is optional. A missing library downgrades one format to
`unsupported`, which the UI renders as "no preview, download it" — never as a
broken page.
"""
from __future__ import annotations

import csv
import io
import json
from pathlib import Path

# Generous enough for a real document, bounded enough that a pathological file
# cannot hand the client a hundred megabytes of JSON.
MAX_TEXT_CHARS = 200_000
MAX_ROWS_PER_SHEET = 500
MAX_COLS = 40
MAX_BLOCKS = 2_000
MAX_SLIDES = 200

TEXTUAL = {".txt", ".md", ".markdown", ".log", ".py", ".js", ".ts", ".tsx", ".jsx",
           ".json", ".html", ".htm", ".css", ".xml", ".yml", ".yaml", ".ini",
           ".toml", ".sql", ".sh", ".bat", ".ps1", ".rst", ".c", ".h", ".cpp",
           ".java", ".rs", ".go", ".rb", ".php"}
TABULAR = {".csv", ".tsv"}
IMAGES = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".svg", ".ico"}


def describe(asset: dict) -> dict:
    """What the client needs to display this asset, and nothing more."""
    name = asset.get("original_name") or ""
    suffix = Path(name).suffix.lower()
    base = {"asset_id": asset.get("id"), "name": name,
            "bytes": asset.get("bytes"), "mime": asset.get("mime")}

    # Handed straight to the browser, which already knows these. Sending the
    # bytes through here would only add a copy.
    if suffix == ".pdf" or (asset.get("mime") or "") == "application/pdf":
        return {**base, "kind": "pdf"}
    if suffix in IMAGES or (asset.get("mime") or "").startswith("image/"):
        return {**base, "kind": "image"}

    path = Path(asset.get("path") or "")
    if not path.is_file():
        return {**base, "kind": "missing"}

    try:
        if suffix in TABULAR:
            # The display name comes from the asset, never from the path: files
            # are stored content-addressed, so `path.stem` is a sha256 and the
            # viewer labelled the sheet with 64 hex characters.
            return {**base, **_delimited(path, "\t" if suffix == ".tsv" else ",",
                                         Path(name).stem)}
        if suffix == ".xlsx" or suffix == ".xlsm":
            return {**base, **_spreadsheet(path)}
        if suffix == ".docx":
            return {**base, **_word(path)}
        if suffix == ".pptx":
            return {**base, **_slides(path)}
        if suffix in TEXTUAL:
            return {**base, **_text(path, suffix)}
        if suffix in (".db", ".sqlite", ".sqlite3"):
            return {**base, **_sqlite(path)}
    except Exception as exc:
        # A file that cannot be parsed is still a file the user may download.
        # Reporting why beats rendering an empty viewer.
        return {**base, "kind": "unreadable", "error": f"{type(exc).__name__}: {exc}"}

    return {**base, "kind": "unsupported"}


# ── text ─────────────────────────────────────────────────────────────────────
def _text(path: Path, suffix: str) -> dict:
    raw = path.read_bytes()[: MAX_TEXT_CHARS * 2]
    text = raw.decode("utf-8", "replace")
    truncated = len(text) > MAX_TEXT_CHARS
    language = {".md": "markdown", ".markdown": "markdown"}.get(
        suffix, suffix.lstrip(".") or "text")
    return {"kind": "text", "language": language,
            "text": text[:MAX_TEXT_CHARS], "truncated": truncated}


# ── csv / tsv ────────────────────────────────────────────────────────────────
def _delimited(path: Path, delimiter: str, display_name: str) -> dict:
    text = path.read_bytes().decode("utf-8", "replace")
    reader = csv.reader(io.StringIO(text), delimiter=delimiter)
    rows: list[list[str]] = []
    total = 0
    for row in reader:
        total += 1
        if len(rows) < MAX_ROWS_PER_SHEET:
            rows.append([c for c in row[:MAX_COLS]])
    header, body = (rows[0], rows[1:]) if rows else ([], [])
    return {"kind": "sheets", "sheets": [{
        "name": display_name, "header": header, "rows": body,
        "total_rows": max(total - 1, 0),
        "truncated": total - 1 > len(body),
    }]}


# ── xlsx ─────────────────────────────────────────────────────────────────────
def _spreadsheet(path: Path) -> dict:
    import openpyxl

    # Handed the bytes, not the path. Assets are stored content-addressed —
    # the file on disk is named after its sha256 and has no extension — and
    # openpyxl decides the format from the filename, so it refuses a perfectly
    # good workbook it cannot see a `.xlsx` on. python-docx and python-pptx
    # sniff the zip instead, which is why only this one needed it.
    #
    # read_only keeps a large workbook from being materialised in full, and
    # data_only asks for the last computed value rather than the formula —
    # a preview should show what the sheet says, not how it says it.
    wb = openpyxl.load_workbook(io.BytesIO(path.read_bytes()),
                                read_only=True, data_only=True)
    sheets = []
    for ws in wb.worksheets:
        rows: list[list] = []
        total = 0
        for row in ws.iter_rows(values_only=True):
            if row is None or all(v is None for v in row):
                continue
            total += 1
            if len(rows) < MAX_ROWS_PER_SHEET + 1:
                rows.append(["" if v is None else _cell(v) for v in row[:MAX_COLS]])
        header, body = (rows[0], rows[1:]) if rows else ([], [])
        sheets.append({"name": ws.title, "header": header, "rows": body,
                       "total_rows": max(total - 1, 0),
                       "truncated": max(total - 1, 0) > len(body)})
    wb.close()
    return {"kind": "sheets", "sheets": sheets}


def _cell(value) -> str:
    if isinstance(value, float) and value.is_integer():
        return str(int(value))
    return str(value)


# ── docx ─────────────────────────────────────────────────────────────────────
def _word(path: Path) -> dict:
    import docx

    document = docx.Document(str(path))
    blocks: list[dict] = []
    for para in document.paragraphs:
        text = (para.text or "").strip()
        if not text:
            continue
        style = (para.style.name or "").lower() if para.style else ""
        if style.startswith("heading"):
            level = "".join(ch for ch in style if ch.isdigit())
            blocks.append({"type": "heading", "level": int(level or 1), "text": text})
        elif style.startswith("list") or style.startswith("bullet"):
            blocks.append({"type": "bullet", "text": text})
        else:
            blocks.append({"type": "paragraph", "text": text})
        if len(blocks) >= MAX_BLOCKS:
            break

    for table in document.tables:
        rows = [[(c.text or "").strip() for c in r.cells[:MAX_COLS]]
                for r in table.rows[:MAX_ROWS_PER_SHEET]]
        if rows:
            blocks.append({"type": "table", "rows": rows})

    return {"kind": "document", "blocks": blocks[:MAX_BLOCKS]}


# ── pptx ─────────────────────────────────────────────────────────────────────
def _slides(path: Path) -> dict:
    from pptx import Presentation

    deck = Presentation(str(path))
    slides = []
    for i, slide in enumerate(deck.slides, start=1):
        title = ""
        lines: list[str] = []
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text = (shape.text_frame.text or "").strip()
            if not text:
                continue
            # The title placeholder, when there is one; otherwise the first
            # line of text stands in, which is what the slide looks like
            # anyway.
            if not title and shape == getattr(slide.shapes, "title", None):
                title = text
            else:
                lines.extend(l.strip() for l in text.splitlines() if l.strip())
        if not title and lines:
            title = lines.pop(0)
        slides.append({"index": i, "title": title, "lines": lines})
        if len(slides) >= MAX_SLIDES:
            break
    return {"kind": "slides", "slides": slides}


# ── sqlite ───────────────────────────────────────────────────────────────────
def _sqlite(path: Path) -> dict:
    """Opened read-only through a URI, so previewing a database can never
    write to one — not even a journal file beside it."""
    import sqlite3

    conn = sqlite3.connect(f"file:{path}?mode=ro", uri=True)
    try:
        names = [r[0] for r in conn.execute(
            "SELECT name FROM sqlite_master WHERE type='table' "
            "AND name NOT LIKE 'sqlite_%' ORDER BY name")]
        sheets = []
        for table in names:
            cursor = conn.execute(f'SELECT * FROM "{table}" LIMIT {MAX_ROWS_PER_SHEET}')
            header = [d[0] for d in cursor.description][:MAX_COLS]
            rows = [["" if v is None else str(v) for v in row[:MAX_COLS]]
                    for row in cursor.fetchall()]
            total = conn.execute(f'SELECT COUNT(*) FROM "{table}"').fetchone()[0]
            sheets.append({"name": table, "header": header, "rows": rows,
                           "total_rows": total, "truncated": total > len(rows)})
    finally:
        conn.close()
    return {"kind": "sheets", "sheets": sheets}


def to_json(value) -> str:
    return json.dumps(value, default=str)
