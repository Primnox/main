"""Mechanical auditor for generated decks.

The PPT torture test asks for a 120-250 slide deck under a strict design system
and ends with a ten-item validation checklist. Every item on that checklist is
a property of the .pptx file, not a matter of taste, so it is checked here
rather than by looking at slides.

This is the scorer, not the generator. It takes a finished file and a design
system and returns findings. That split is deliberate: the same auditor grades
a deck from Opus, from a 7B, or from a hand-written fixture, and a run is only
comparable to another run if the thing measuring them never changes.

Approximation is confined to ONE check. Text overflow needs font metrics that
python-pptx cannot provide without a rendering engine, so `_overflows` estimates
and deliberately errs toward silence: a flagged overflow should be real, because
an auditor that cries wolf gets ignored and then the checklist is worthless.
Every other check is exact.
"""
from __future__ import annotations

import hashlib
import math
from collections import Counter, defaultdict
from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

EMU_PER_PT = 12700

# Severity is about whether a human must look, not about how ugly it is.
ERROR, WARN = "error", "warn"

# Consecutive slides sharing one layout before it reads as a wall. Three is a
# deliberate group; four is the audience realising nothing is going to change.
MAX_SAME_LAYOUT_RUN = 3


@dataclass
class Finding:
    slide: int          # 1-based; 0 means deck-wide
    check: str
    detail: str
    severity: str = ERROR

    def __str__(self) -> str:
        where = f"slide {self.slide}" if self.slide else "deck"
        return f"[{self.severity}] {where}: {self.check} - {self.detail}"


@dataclass
class DesignSystem:
    """The contract a deck is graded against."""
    palette: set[str] = field(default_factory=set)      # uppercase hex, no '#'
    font_sizes: set[float] = field(default_factory=set)  # points
    fonts: set[str] = field(default_factory=set)
    grid_pt: int = 8
    min_readable_pt: float = 16.0
    # Shapes covering this fraction of the slide are backgrounds: they are
    # expected to sit under everything, so they are exempt from overlap.
    background_area_ratio: float = 0.92

    @classmethod
    def torture_test(cls) -> "DesignSystem":
        """The system specified in the Primnox PPT Torture Test prompt."""
        return cls(
            palette={
                "0F172A",  # deep navy   - primary
                "334155",  # slate       - secondary
                "22D3EE",  # cyan        - accent
                "A3E635",  # lime        - highlight
                "10B981",  # emerald     - success
                "F59E0B",  # amber       - warning
                "F43F5E",  # rose        - error
                "0B1220",  # background
                "F8FAFC",  # text
                "94A3B8",  # muted
            },
            font_sizes={54, 40, 30, 22, 16},
            fonts={"Inter", "Arial", "JetBrains Mono"},
            grid_pt=8,
            min_readable_pt=16.0,
        )


# ── geometry helpers ─────────────────────────────────────────────────────────
def _pt(v) -> float | None:
    return None if v is None else Emu(v).pt


def _box(shape):
    for attr in ("left", "top", "width", "height"):
        if getattr(shape, attr, None) is None:
            return None
    l, t = _pt(shape.left), _pt(shape.top)
    return (l, t, l + _pt(shape.width), t + _pt(shape.height))


def _intersects(a, b, tol: float = 1.0) -> bool:
    """Boxes overlap by more than `tol` points on BOTH axes.

    The tolerance is not cosmetic: abutting shapes on an 8pt grid share an edge
    exactly, and a strict test would report every adjacent pair in the deck.
    """
    return (min(a[2], b[2]) - max(a[0], b[0]) > tol
            and min(a[3], b[3]) - max(a[1], b[1]) > tol)


def _runs(shape):
    if not getattr(shape, "has_text_frame", False):
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            yield para, run


def _colors(shape):
    """Every explicitly-set RGB on a shape: fill, line, and each run."""
    out: list[str] = []

    def _take(fmt):
        try:
            if fmt is not None and fmt.type is not None and fmt.rgb is not None:
                out.append(str(fmt.rgb).upper())
        except Exception:
            pass  # theme-inherited or scheme colors have no rgb; not drift

    try:
        _take(shape.fill.fore_color)
    except Exception:
        pass
    try:
        _take(shape.line.color)
    except Exception:
        pass
    for _, run in _runs(shape):
        try:
            _take(run.font.color)
        except Exception:
            pass
    return out


def _overflows(shape, system: DesignSystem) -> tuple[str, str] | None:
    """Estimate whether text exceeds its box. Errs toward silence.

    Returns (check, detail) or None. The autofit mode decides which of three
    different things an over-long string actually means, and they are not
    interchangeable:

      TEXT_TO_FIT_SHAPE  PowerPoint shrinks the type. Not an overflow; the risk
                         is dropping under the readability floor, which the
                         per-run size check already covers. Skipped.
      SHAPE_TO_FIT_TEXT  PowerPoint grows the BOX, past the geometry stored in
                         the file. Reported, because every other geometric check
                         here — grid alignment, overlap — reads that stored
                         geometry and is therefore unreliable for this shape.
      NONE               Nothing adjusts. The text is genuinely clipped.

    python-pptx defaults new textboxes to SHAPE_TO_FIT_TEXT, so treating "not
    NONE" as "already handled" silently disabled this check for every shape a
    generator produces — which is exactly what the planted-overflow test caught.
    """
    if not getattr(shape, "has_text_frame", False):
        return None
    tf = shape.text_frame
    text = (tf.text or "").strip()
    if not text:
        return None

    mode = None
    try:
        mode = tf.auto_size
    except Exception:
        pass
    if mode is not None and getattr(mode, "value", None) == 2:   # TEXT_TO_FIT_SHAPE
        return None
    grows = mode is not None and getattr(mode, "value", None) == 1  # SHAPE_TO_FIT_TEXT

    w, h = _pt(shape.width), _pt(shape.height)
    if not w or not h:
        return None
    inset_w = w - _pt(tf.margin_left or 0) - _pt(tf.margin_right or 0)
    inset_h = h - _pt(tf.margin_top or 0) - _pt(tf.margin_bottom or 0)
    if inset_w <= 0 or inset_h <= 0:
        return None

    needed = 0.0
    for para in tf.paragraphs:
        size = next((r.font.size.pt for r in para.runs
                     if r.font.size is not None), None) or 18.0
        line = (para.text or "")
        # 0.5em average advance is generous for Inter/Arial at sentence case;
        # a tighter figure would manufacture overflows on ordinary text.
        chars_per_line = max(1, int(inset_w / (size * 0.5)))
        lines = max(1, math.ceil(len(line) / chars_per_line)) if line else 1
        # The paragraph's OWN line spacing when it declares one. A flat 1.2
        # estimate under-counts any deck that sets looser leading, and this
        # missed a real defect: three lines of card content at 1.45/1.15/1.45
        # measured as fitting and rendered outside the panel.
        spacing = 1.2
        try:
            declared = para.line_spacing
            if isinstance(declared, (int, float)) and declared > 0:
                spacing = float(declared)
        except Exception:
            pass
        needed += lines * size * spacing

    # 5% headroom, not 15%. The wide tolerance was set when every line was
    # estimated at a flat 1.2 leading and the guess could be off by a fifth;
    # now that the paragraph's own line spacing is read, the estimate is close
    # enough that 15% only hides real defects — it passed a card whose last
    # line rendered outside the panel by two points.
    if needed <= inset_h * 1.05:
        return None
    if grows:
        return ("autogrow-geometry",
                f"box will expand to ~{needed:.0f}pt from a stored {inset_h:.0f}pt; "
                "grid and overlap results for this shape are not trustworthy")
    return ("text-overflow", f"~{needed:.0f}pt of text in a {inset_h:.0f}pt box")


# ── the checks ───────────────────────────────────────────────────────────────
def audit(path: str | Path, system: DesignSystem | None = None) -> list[Finding]:
    system = system or DesignSystem.torture_test()
    prs = Presentation(str(path))
    findings: list[Finding] = []

    slide_w, slide_h = _pt(prs.slide_width), _pt(prs.slide_height)
    slide_area = slide_w * slide_h

    seen_fonts: Counter = Counter()
    layout_fingerprints: dict[str, list[int]] = defaultdict(list)
    title_lefts: Counter = Counter()
    slide_shapes: list[tuple] = []

    for idx, slide in enumerate(prs.slides, start=1):
        shapes = [s for s in slide.shapes]
        boxes: list[tuple] = []

        # -- titles -------------------------------------------------------
        title = None
        try:
            title = slide.shapes.title
        except Exception:
            pass
        has_title_text = bool(title is not None and (title.text or "").strip())
        if not has_title_text:
            # Most generated decks never touch the title placeholder, so a
            # surrogate is required or every slide reads as untitled.
            #
            # The rule is "largest type in the upper two-thirds", not "anything
            # near the top". Position alone rejected hero and cover slides,
            # whose headline is deliberately centred - which is a real layout in
            # the spec, so flagging it was the auditor being wrong rather than
            # the deck. Size is what actually distinguishes a title.
            texts = [s for s in shapes
                     if getattr(s, "has_text_frame", False) and (s.text or "").strip()]
            sizes = [r.font.size.pt for s in texts for _, r in _runs(s)
                     if r.font.size is not None]
            biggest = max(sizes, default=None)
            surrogate = any(
                _pt(s.top) is not None and _pt(s.top) < slide_h * 0.66
                and (biggest is None
                     or any(r.font.size is not None and r.font.size.pt == biggest
                            for _, r in _runs(s)))
                for s in texts
            )
            if not surrogate:
                findings.append(Finding(idx, "missing-title", "no title text on the slide"))
        if title is not None and _pt(title.left) is not None:
            title_lefts[round(_pt(title.left), 1)] += 1

        for shape in shapes:
            box = _box(shape)

            # -- colour drift ---------------------------------------------
            if system.palette:
                for rgb in _colors(shape):
                    if rgb not in system.palette:
                        findings.append(Finding(
                            idx, "color-drift",
                            f"#{rgb} is not a design-system token "
                            f"(shape {shape.shape_type}, name {shape.name!r})"))

            # -- fonts and hierarchy --------------------------------------
            for para, run in _runs(shape):
                name = run.font.name
                if name:
                    seen_fonts[name] += 1
                    if system.fonts and name not in system.fonts:
                        findings.append(Finding(
                            idx, "font-drift", f"{name!r} is not in the type system"))
                size = run.font.size
                if size is not None:
                    pts = size.pt
                    if system.font_sizes and pts not in system.font_sizes:
                        findings.append(Finding(
                            idx, "type-scale",
                            f"{pts:g}pt is off the scale "
                            f"{sorted(system.font_sizes)}", WARN))
                    if pts < system.min_readable_pt:
                        findings.append(Finding(
                            idx, "below-readable",
                            f"{pts:g}pt is under the {system.min_readable_pt:g}pt floor"))

            # -- orphan bullets -------------------------------------------
            if getattr(shape, "has_text_frame", False):
                for para in shape.text_frame.paragraphs:
                    if para.level > 0 and not (para.text or "").strip():
                        findings.append(Finding(
                            idx, "orphan-bullet",
                            f"empty bullet at indent level {para.level}"))

            # -- grid alignment -------------------------------------------
            # A full-bleed element is exempt on the axis it bleeds along. The
            # slide is 540pt tall, which is not a multiple of 8, so a rule that
            # runs the full height CANNOT be on the grid — and a designer who
            # inset it by 4pt to satisfy the checker would have made the deck
            # worse to satisfy the tool.
            if box:
                bleeds_w = _pt(shape.width) is not None and _pt(shape.width) >= slide_w - 0.5
                bleeds_h = _pt(shape.height) is not None and _pt(shape.height) >= slide_h - 0.5
                exempt = {"width"} if bleeds_w else set()
                if bleeds_h:
                    exempt |= {"height", "top"}
                for label, value in (("left", _pt(shape.left)), ("top", _pt(shape.top)),
                                     ("width", _pt(shape.width)), ("height", _pt(shape.height))):
                    if value is None or label in exempt:
                        continue
                    off = abs(value - round(value / system.grid_pt) * system.grid_pt)
                    if off > 0.75:   # sub-point slack; PowerPoint stores EMU
                        findings.append(Finding(
                            idx, "grid-misalign",
                            f"{shape.name!r} {label}={value:.1f}pt is "
                            f"{off:.1f}pt off the {system.grid_pt}pt grid", WARN))

            # -- overflow --------------------------------------------------
            over = _overflows(shape, system)
            if over:
                check, detail = over
                findings.append(Finding(
                    idx, check, f"{shape.name!r}: {detail}",
                    ERROR if check == "text-overflow" else WARN))

            # Background-sized shapes are meant to sit under everything.
            if box and (box[2] - box[0]) * (box[3] - box[1]) < slide_area * system.background_area_ratio:
                if getattr(shape, "has_text_frame", False) and (shape.text or "").strip():
                    boxes.append((box, shape.name))
                elif shape.shape_type is not None and not getattr(shape, "has_text_frame", False):
                    boxes.append((box, shape.name))

        # -- overlap ------------------------------------------------------
        for i in range(len(boxes)):
            for j in range(i + 1, len(boxes)):
                if _intersects(boxes[i][0], boxes[j][0]):
                    findings.append(Finding(
                        idx, "overlap",
                        f"{boxes[i][1]!r} overlaps {boxes[j][1]!r}"))

        # -- accidental duplicate slides ----------------------------------
        body = "␟".join(
            (s.text or "").strip() for s in shapes if getattr(s, "has_text_frame", False))
        fingerprint = hashlib.sha256(
            f"{slide.slide_layout.name}␞{body}".encode("utf-8")).hexdigest()
        layout_fingerprints[fingerprint].append(idx)

        # Geometry only, no text: two bullet slides with different words are the
        # same layout, which is exactly what the rhythm check needs to notice.
        slide_shapes.append(tuple(sorted(
            (str(s.shape_type), _pt(s.left), _pt(s.top), _pt(s.width), _pt(s.height))
            for s in shapes if _box(s))))

    # ── deck-wide ────────────────────────────────────────────────────────
    for slides in layout_fingerprints.values():
        if len(slides) > 1:
            findings.append(Finding(
                slides[1], "duplicate-slide",
                f"identical layout and content to slide {slides[0]}"))

    # -- rhythm --------------------------------------------------------------
    # A deck-level property: nothing per-slide is wrong with six bullet slides
    # in a row, and it is still a bad deck. Shape signature stands in for layout
    # because generated decks all use the same blank master — what distinguishes
    # them is what got drawn on it.
    run_start, run_sig = 0, None
    for i, sig in enumerate(slide_shapes):
        if sig != run_sig:
            if run_sig is not None and i - run_start > MAX_SAME_LAYOUT_RUN:
                findings.append(Finding(
                    run_start + 1, "monotonous-run",
                    f"{i - run_start} consecutive slides share one layout", WARN))
            run_start, run_sig = i, sig
    if run_sig is not None and len(slide_shapes) - run_start > MAX_SAME_LAYOUT_RUN:
        findings.append(Finding(
            run_start + 1, "monotonous-run",
            f"{len(slide_shapes) - run_start} consecutive slides share one layout",
            WARN))

    families = {f for f in seen_fonts if f}
    if len(families) > len(system.fonts):
        findings.append(Finding(
            0, "font-count",
            f"{len(families)} font families in use: {sorted(families)}"))

    if len(title_lefts) > 1:
        common = title_lefts.most_common()
        findings.append(Finding(
            0, "inconsistent-margin",
            f"title left edge varies across slides: {common[:4]}", WARN))

    return findings


def summarise(findings: list[Finding], slide_count: int) -> dict:
    by_check: Counter = Counter(f.check for f in findings)
    errors = [f for f in findings if f.severity == ERROR]
    return {
        "slides": slide_count,
        "findings": len(findings),
        "errors": len(errors),
        "warnings": len(findings) - len(errors),
        "by_check": dict(by_check.most_common()),
        "clean": not errors,
    }


def slide_count(path: str | Path) -> int:
    return len(Presentation(str(path)).slides)
